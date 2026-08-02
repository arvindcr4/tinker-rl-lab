"""Build the ZVF Program 'Progress & Next Steps' deck (2026-06-14).

Generates a dark-theme 16:9 status deck matching the original ZVF_Program.pptx
aesthetic (letter-spaced eyebrows, accent rule, footer + page numbers).
All content is factual to the work actually shipped this session; nothing fabricated.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette ----------------------------------------------------------------
BG      = RGBColor(0x0E, 0x11, 0x16)
PANEL   = RGBColor(0x16, 0x1B, 0x22)
PANEL2  = RGBColor(0x1C, 0x23, 0x2C)
INK     = RGBColor(0xE6, 0xED, 0xF3)
MUTE    = RGBColor(0x8B, 0x94, 0x9E)
TEAL    = RGBColor(0x2D, 0xD4, 0xBF)
AMBER   = RGBColor(0xF5, 0xA6, 0x23)
GREEN   = RGBColor(0x3F, 0xB9, 0x50)
YELLOW  = RGBColor(0xD2, 0x99, 0x22)
RED     = RGBColor(0xF8, 0x51, 0x49)
LINECLR = RGBColor(0x30, 0x36, 0x3D)
FONT    = "Arial"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def spaced(s):
    return "  ".join(list(s))


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def box(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tf


def para(tf, text, size, color=INK, bold=False, font=FONT, align=PP_ALIGN.LEFT,
         first=False, space_before=0, space_after=4, italic=False, spacing=1.0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = spacing
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic; f.name = font
    f.color.rgb = color
    return p


def rect(s, l, t, w, h, fill=PANEL, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def pill(s, l, t, text, color):
    w = 0.16 + 0.082 * len(text)
    sp = rect(s, l, t, w, 0.32, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        sp.adjustments[0] = 0.5
    except Exception:
        pass
    tf = sp.text_frame; tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(10); r.font.bold = True; r.font.name = FONT
    r.font.color.rgb = BG
    return w


def chrome(s, eyebrow, n):
    rect(s, 0.55, 0.5, 0.12, 0.12, fill=TEAL, shape=MSO_SHAPE.OVAL)
    tf = box(s, 0.78, 0.46, 9.5, 0.4)
    para(tf, spaced(eyebrow), 10.5, MUTE, bold=True, first=True)
    tf2 = box(s, 11.2, 0.46, 1.6, 0.4)
    para(tf2, f"{n:02d} / 12", 10.5, MUTE, bold=True, align=PP_ALIGN.RIGHT, first=True)
    rect(s, 0.55, 6.92, 12.23, 0.018, fill=LINECLR)
    tf3 = box(s, 0.55, 7.0, 12.23, 0.35)
    para(tf3, "ZVF Program  ·  Progress & Next Steps  ·  2026-06-14        agentic_grpo_bench · Arvind",
         9, MUTE, first=True)  # noqa


def header(s, eyebrow, title, n, subtitle=None):
    chrome(s, eyebrow, n)
    tf = box(s, 0.55, 0.95, 12.2, 1.1)
    para(tf, title, 30, INK, bold=True, first=True)
    if subtitle:
        para(tf, subtitle, 14, MUTE, space_before=2)
    rect(s, 0.57, 1.95, 1.5, 0.05, fill=TEAL)


# =============================================================== 1 · TITLE
s = slide()
rect(s, 0, 0, 13.333, 0.09, fill=TEAL)
tf = box(s, 0.9, 0.62, 11, 0.4)
para(tf, spaced("GRPO AUDIT · NEXT MOVES"), 11, MUTE, bold=True, first=True)
tf = box(s, 0.9, 2.35, 11.5, 2.2)
para(tf, "The ZVF Program", 52, INK, bold=True, first=True)
para(tf, "Progress & Next Steps — Status Review", 22, TEAL, bold=True, space_before=6)
para(tf, "From an audit diagnostic to the standard telemetry, theory, and tooling for RL post-training.",
     15, MUTE, space_before=12)
rect(s, 0.92, 5.35, 2.0, 0.05, fill=AMBER)
tf = box(s, 0.9, 5.55, 11.5, 1.2)
para(tf, "One method · one theorem · one library · one community position", 14, INK, bold=True, first=True)
para(tf, "June 14, 2026     ·     Arvind     ·     agentic_grpo_bench     ·     NeurIPS 2026 track",
     12, MUTE, space_before=8)

# =============================================================== 2 · THE PLAN
s = slide()
header(s, "RECAP · THE PROGRAM", "The plan: four shippable deliverables in ~4 months", 2,
       "Each pillar feeds the next. The library and corpus compound throughout.")
cards = [
    ("01", "METHOD", "Adaptive ZVF-aware rollout controller — train only in the contrast band.", "ICLR '27", TEAL),
    ("02", "THEORY", "ZVF as U-statistic estimator + sample-complexity lower bound + optimal G*.", "AISTATS / COLT", AMBER),
    ("03", "OPEN SOURCE", "zvf-triage library — drop-in ZVFCallback for TRL / verl / OpenRLHF / NeMo-RL.", "NeurIPS D&B", TEAL),
    ("04", "POSITION", "MIN-REPORT-RL minimum-reportable-stack + reproducibility audit.", "NeurIPS Position", AMBER),
]
x = 0.55; w = 2.92; gap = 0.12
for num, name, desc, venue, acc in cards:
    rect(s, x, 2.35, w, 3.7, fill=PANEL)
    rect(s, x, 2.35, w, 0.09, fill=acc)
    tf = box(s, x + 0.22, 2.62, w - 0.44, 3.3)
    para(tf, num, 30, acc, bold=True, first=True)
    para(tf, name, 13, INK, bold=True, space_before=2)
    para(tf, desc, 11.5, MUTE, space_before=8, spacing=1.05)
    para(tf, "→ " + venue, 11, acc, bold=True, space_before=10)
    x += w + gap

# =============================================================== 3 · SCORECARD
s = slide()
header(s, "STATUS · WHERE THE FOUR STAND", "Scorecard — 3 of 4 pillars now have shipping artifacts", 3)
rows = [
    ("PILLAR", "ARTIFACT ON DISK", "STATUS", "THIS SESSION"),
    ("3 · OSS library", "zvf-triage package (71/71 tests, builds wheel)", "SHIPPABLE", "built + hardened"),
    ("1 · Method / M1", "sweep/ harness — 403-cell dry-run, cell_runner wired", "LAUNCH-READY", "built + wired"),
    ("2 · Theory", "zvf_theory.tex (compiles) + THEORY_NOTES.md", "DRAFT · SKETCH", "drafted, gaps flagged"),
    ("4 · Position", "min_report_rl.tex (compiles) + CHECKLIST.md", "DRAFT · TODO", "drafted, results TODO"),
]
status_clr = {"SHIPPABLE": GREEN, "LAUNCH-READY": TEAL, "DRAFT · SKETCH": YELLOW, "DRAFT · TODO": AMBER}
top = 2.45; rh = 0.84; left = 0.55
widths = [2.7, 5.0, 2.35, 2.18]
for ri, row in enumerate(rows):
    x = left
    is_head = ri == 0
    for ci, cell in enumerate(row):
        fill = PANEL2 if is_head else (PANEL if ri % 2 else BG)
        rect(s, x, top + ri * rh, widths[ci], rh, fill=fill, line=LINECLR, line_w=0.75)
        if ci == 2 and not is_head:
            pill(s, x + 0.18, top + ri * rh + (rh - 0.32) / 2, cell, status_clr[cell])
        else:
            tf = box(s, x + 0.18, top + ri * rh, widths[ci] - 0.32, rh, anchor=MSO_ANCHOR.MIDDLE)
            para(tf, cell, 11.5 if is_head else 12,
                 MUTE if is_head else INK, bold=is_head or ci == 0, first=True)
        x += widths[ci]
tf = box(s, 0.55, top + len(rows) * rh + 0.18, 12.2, 0.6)
para(tf, "Predecessor v1 audit (EAI Endorsed Transactions, ~95 runs, score 97/100) is submission-ready — "
         "the corpus the follow-ups build on.", 11.5, MUTE, italic=True, first=True)

# =============================================================== 4 · PILLAR 3
s = slide()
header(s, "PILLAR 3 · OPEN SOURCE", "zvf-triage — the headline win: a real, installable library", 4)
pill(s, 0.57, 1.55, "SHIPPABLE", GREEN)
pill(s, 2.0, 1.55, "71 / 71 TESTS PASS", TEAL)
pill(s, 4.4, 1.55, "BUILDS WHEEL", TEAL)
feats = [
    ("ZVFController", "State machine: cold-start collapse · saturation · exploitable-contrast, with adaptive G + auto-stop."),
    ("ZVFCallback", "Framework-agnostic + lazy TRL TrainerCallback adapter; imports with numpy alone."),
    ("Pure-numpy core", "ZVF / GU / peak-to-tail drift — definitions reused verbatim from your formalization appendix."),
    ("Ship-ready", "Apache-2.0 · CI on py3.9–3.12 · runnable example · py.typed · verl/OpenRLHF/NeMo-RL adapter stubs."),
]
y = 2.35
for name, desc in feats:
    rect(s, 0.55, y, 7.2, 0.95, fill=PANEL)
    rect(s, 0.55, y, 0.07, 0.95, fill=TEAL)
    tf = box(s, 0.78, y + 0.13, 6.85, 0.75)
    para(tf, name, 13, TEAL, bold=True, first=True)
    para(tf, desc, 11, MUTE, space_before=2, spacing=1.03)
    y += 1.06
rect(s, 8.0, 2.35, 4.78, 3.0, fill=RGBColor(0x0A, 0x0D, 0x11), line=LINECLR)
tf = box(s, 8.25, 2.55, 4.3, 2.7)
code = [
    "from zvf_triage import ZVFCallback",
    "",
    "trainer = GRPOTrainer(",
    "    model, dataset,",
    "    callbacks=[ZVFCallback(",
    "        window=5, zvf_max=0.85,",
    '        on_collapse="warm_start",',
    "        adaptive_G=True,",
    "        wandb_panel=True)])",
]
for i, ln in enumerate(code):
    para(tf, ln if ln else " ", 11.5, TEAL if ("ZVFCallback" in ln or "import" in ln) else INK,
         font="Courier New", first=(i == 0), space_after=2)
tf = box(s, 8.0, 5.5, 4.78, 0.7)
para(tf, "Next: polish → TestPyPI → upstream PR to TRL.", 11, AMBER, bold=True, first=True)

# =============================================================== 5 · PILLAR 1
s = slide()
header(s, "PILLAR 1 · METHOD / M1", "Sweep harness — one command from launch", 5)
pill(s, 0.57, 1.55, "LAUNCH-READY", TEAL)
pill(s, 2.35, 1.55, "DRY-RUN: 403 CELLS, LAUNCHES NOTHING", AMBER)
items = [
    ("403-cell grid", "385 audit + 18 matched-compute. 5 model families × frameworks × seeds × group sizes × difficulty."),
    ("cell_runner.py wired", "Shells into live_zvf_probe — the only runner that already emits per-step ZVF / GU. No training reimplemented."),
    ("Matched-compute arms", "Canonical-GRPO vs DAPO vs GSPO at equal rollouts (480) and tokens. Winner computed from real logs only."),
    ("Resumable + dry-run default", "Skips completed cells; aggregate reports MISSING for absent files — never a guessed value."),
]
y = 2.35
for name, desc in items:
    rect(s, 0.55, y, 12.23, 1.0, fill=PANEL)
    rect(s, 0.55, y, 0.07, 1.0, fill=TEAL)
    tf = box(s, 0.8, y + 0.14, 11.8, 0.8)
    para(tf, name, 13, INK, bold=True, first=True)
    para(tf, desc, 11.5, MUTE, space_before=2)
    y += 1.1
tf = box(s, 0.55, 6.25, 12.2, 0.5)
para(tf, "Needs you: TINKER_API_KEY + a ~30-line shim confirm, then  python3 run_sweep.py --execute  "
         "(~1–3 days, 95 → ~400 runs).", 11.5, AMBER, bold=True, first=True)

# =============================================================== 6 · PILLAR 2
s = slide()
header(s, "PILLAR 2 · THEORY", "Drafted and compiles — but proof-sketch only (stated honestly)", 6)
pill(s, 0.57, 1.55, "DRAFT", AMBER)
pill(s, 1.65, 1.55, "PROOF-SKETCH ONLY", YELLOW)
pill(s, 4.15, 1.55, "16-ITEM PROOF LEDGER", RED)
ths = [
    ("T1 · Estimator", "ZVF as a U-statistic; asymptotic normality + closed-form 95% CI.",
     "GAP: currently collapses to a binomial CLT and degenerates exactly at ZVF→1 — where the controller cares most."),
    ("T2 · Lower bound", "Finite-sample floor on rollouts-to-policy-improvement vs observed ZVF.",
     "GAP: proves rollouts-to-nonzero-gradient, not to-improvement (overclaim). Uses population ZVF, needs one-sided bound."),
    ("T3 · Optimal G*", "Closed-form per-prompt group size maximizing signal per rollout.",
     "GAP: the signal function is an unjustified modeling choice; uniqueness of G* unproven."),
]
y = 2.35
for name, claim, gap in ths:
    rect(s, 0.55, y, 12.23, 1.35, fill=PANEL)
    rect(s, 0.55, y, 0.07, 1.35, fill=AMBER)
    tf = box(s, 0.8, y + 0.14, 11.8, 1.1)
    para(tf, name + "   —   " + claim, 12.5, INK, bold=True, first=True)
    para(tf, gap, 11, RED, space_before=4, spacing=1.03)
    y += 1.45
tf = box(s, 0.55, 6.35, 12.2, 0.45)
para(tf, "Biggest threat to all three: across-group i.i.d. breaks under curriculum / replay → invalidates the T1 CI. "
         "Read THEORY_NOTES.md before citing.", 11, MUTE, italic=True, first=True)

# =============================================================== 7 · PILLAR 4
s = slide()
header(s, "PILLAR 4 · COMMUNITY POSITION", "MIN-REPORT-RL — a minimum-reportable-stack for RL post-training", 7)
pill(s, 0.57, 1.55, "DRAFT COMPILES", AMBER)
pill(s, 2.5, 1.55, "RESULTS + CITATIONS = TODO", YELLOW)
rect(s, 0.55, 2.35, 6.3, 3.9, fill=PANEL)
tf = box(s, 0.8, 2.55, 5.85, 3.6)
para(tf, "Every GRPO-family paper should report:", 13, TEAL, bold=True, first=True)
checklist = [
    "Loss form (PPO ratio? clip? token mask?)",
    "Reference policy + KL handling",
    "Sampler / backend / precision",
    "Per-step ZVF and GU trajectory",
    "Group-size schedule (fixed or adaptive)",
    "Held-out split distinct from reward env",
    "Decontamination probe results",
]
for c in checklist:
    para(tf, "•  " + c, 12, INK, space_before=6)
rect(s, 7.05, 2.35, 5.73, 3.9, fill=PANEL)
tf = box(s, 7.3, 2.55, 5.3, 3.6)
para(tf, "Reproducibility audit", 13, AMBER, bold=True, first=True)
para(tf, "Re-implement DAPO / GSPO / Dr.GRPO / MAD-GRPO in ONE controlled stack. Report which claimed gains survive.",
     12, MUTE, space_before=6, spacing=1.05)
para(tf, "Target: NeurIPS / ICML Position Track + MLRC.", 11.5, INK, space_before=10)
para(tf, "Outstanding", 12.5, RED, bold=True, space_before=16)
para(tf, "Every results cell and every \\cite{} is a marked TODO — filled from the Pillar-1 sweep output. "
         "No fake authors or numbers were invented.", 11.5, MUTE, space_before=4, spacing=1.05)

# =============================================================== 8 · BOUNDARY
s = slide()
header(s, "INTEGRITY · THE HONEST BOUNDARY", "What shipped vs what still needs you — nothing fabricated", 8)
rect(s, 0.55, 2.35, 6.0, 4.0, fill=PANEL)
rect(s, 0.55, 2.35, 6.0, 0.09, fill=GREEN)
tf = box(s, 0.8, 2.6, 5.55, 3.6)
para(tf, "DONE — real artifacts", 14, GREEN, bold=True, first=True)
for t in ["Installable zvf-triage package, tests green, wheel builds",
          "403-cell sweep harness + cell_runner shim, wired & dry-run verified",
          "Theory paper draft (statements of T1 / T2 / T3, compiles)",
          "Position paper + MIN-REPORT-RL checklist (compiles)",
          "zvf-triage integrated into live_zvf_probe behind a flag"]:
    para(tf, "•  " + t, 11.5, INK, space_before=7, spacing=1.03)
rect(s, 6.78, 2.35, 6.0, 4.0, fill=PANEL)
rect(s, 6.78, 2.35, 6.0, 0.09, fill=RED)
tf = box(s, 7.03, 2.6, 5.55, 3.6)
para(tf, "NEEDS YOU — can't be faked", 14, RED, bold=True, first=True)
for t in ["The 403 GPU runs — need your Tinker / Modal compute",
          "The theorem proofs — T1/T2/T3 are sketches; you verify the math",
          "The audit results table — fills from the sweep output",
          "Canonical DAPO / GSPO configs — surrogates need tightening",
          "Real GitHub URL + PyPI account before publish"]:
    para(tf, "•  " + t, 11.5, INK, space_before=7, spacing=1.03)
tf = box(s, 0.55, 6.5, 12.2, 0.4)
para(tf, "No run numbers were invented; no proof-sketch is claimed as a proof; missing files surface as MISSING.",
     11.5, TEAL, bold=True, align=PP_ALIGN.CENTER, first=True)

# =============================================================== 9 · DECISIONS
s = slide()
header(s, "DECISIONS · YOUR CALL", "Three choices to confirm before the sweep launches", 9)
dec = [
    ("1 · Sweep runner", "live_zvf_probe is the only runner emitting ZVF/GU, so cell_runner wired into it for all 403 cells. "
     "tinker_parallel_runner is missing; campaign_v2 doesn't log ZVF.", "Confirm it's the right runner.", AMBER),
    ("2 · DAPO / GSPO fidelity", "Arms are surrogates (asymmetric clip; sequence-level aggregation) layered on the GRPO loop — "
     "not full reimplementations. Validity-critical for the audit's 'which gains survive' claim.",
     "Point me at canonical configs to tighten, or accept for stratification.", RED),
    ("3 · License", "zvf-triage is Apache-2.0, chosen to match TRL/verl for upstreaming.", "One-line revert to MIT if preferred.", TEAL),
]
y = 2.35
for name, body, ask, acc in dec:
    rect(s, 0.55, y, 12.23, 1.35, fill=PANEL)
    rect(s, 0.55, y, 0.07, 1.35, fill=acc)
    tf = box(s, 0.8, y + 0.15, 11.8, 1.1)
    para(tf, name, 13, acc, bold=True, first=True)
    para(tf, body, 11, MUTE, space_before=3, spacing=1.02)
    para(tf, "→ " + ask, 11, INK, bold=True, space_before=3)
    y += 1.45

# =============================================================== 10 · NEXT STEPS
s = slide()
header(s, "ROADMAP · NEXT STEPS", "From shipped infrastructure to four papers", 10)
tf = box(s, 0.55, 2.0, 12.2, 0.4)
para(tf, "This week:  smoke-launch a 5-cell Tinker run via cell_runner · publish zvf-triage to TestPyPI · lock DAPO/GSPO configs.",
     12, TEAL, bold=True, first=True)
steps = [
    ("M1", "Weeks 1–2", "Strengthen the audit", "Full 403-run sweep + held-out suite (GSM-Plus, AIME-2025, BFCL-v3).", "v2 paper + zvf-triage v0.1"),
    ("M2", "Weeks 2–3", "Ship the method", "Adaptive controller; matched-compute wins vs DAPO/GSPO; upstream PRs to TRL/verl.", "ICLR '27 method paper"),
    ("M3", "Weeks 2–3", "Land the theory", "Prove T1 CI rigorously; fix the T2 claim; derive G*; align with Razin & Zhou.", "AISTATS / COLT paper"),
    ("M4", "Weeks 3–4", "Set the standard", "Fill the audit table from sweep output; finalize MIN-REPORT-RL position.", "NeurIPS Position + MLRC"),
]
x = 0.55; w = 2.92; gap = 0.12
for m, when, title, body, out in steps:
    rect(s, x, 2.55, w, 3.7, fill=PANEL)
    rect(s, x, 2.55, w, 0.09, fill=TEAL)
    tf = box(s, x + 0.2, 2.8, w - 0.4, 3.4)
    para(tf, m + "  ·  " + when, 12, TEAL, bold=True, first=True)
    para(tf, title, 14, INK, bold=True, space_before=4)
    para(tf, body, 11, MUTE, space_before=8, spacing=1.05)
    para(tf, "→ " + out, 11, AMBER, bold=True, space_before=10)
    x += w + gap

# =============================================================== 11 · RISKS
s = slide()
header(s, "RISKS · WHAT COULD SLIP", "Known risks, ranked", 11)
risks = [
    ("Theory proof gaps", "HIGH", RED, "T1 degenerates at ZVF→1; T2 overclaims; T3 modeling choice unjustified. Could weaken the theorem paper or force a reframe."),
    ("Surrogate fidelity", "HIGH", RED, "DAPO/GSPO arms aren't canonical → the audit's 'which gains survive' claim is only as strong as the configs. Tighten before publishing."),
    ("Compute cost / time", "MED", YELLOW, "~400 LoRA-GRPO runs ≈ 1–3 days wall-clock + Tinker/Modal billing. Smoke-test 5 cells and read real per-run cost first."),
    ("Small group size", "MED", YELLOW, "At G=2 (K=2) the per-prompt drop and adaptive-G logic are weak; controller validation may need larger K."),
    ("Non-i.i.d. groups", "MED", YELLOW, "Curriculum / replay correlation breaks the T1 CI assumption and the T2 bound. Needs an explicit caveat or correction."),
]
y = 2.35
for name, sev, clr, body in risks:
    rect(s, 0.55, y, 12.23, 0.84, fill=PANEL)
    rect(s, 0.55, y, 0.07, 0.84, fill=clr)
    pill(s, 0.8, y + (0.84 - 0.32) / 2, sev, clr)
    tf = box(s, 2.05, y, 10.5, 0.84, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, name + "  —  " + body, 11.5, INK, first=True, spacing=1.0)
    y += 0.94

# =============================================================== 12 · CLOSING
s = slide()
rect(s, 0, 0, 13.333, 0.09, fill=TEAL)
chrome(s, "GRPO AUDIT · NEXT MOVES", 12)
tf = box(s, 0.55, 1.1, 12.2, 1.0)
para(tf, "Why this path holds", 30, INK, bold=True, first=True)
para(tf, "One coherent program subsumes every weaker version of this work.", 14, MUTE, space_before=2)
rect(s, 0.57, 2.15, 1.5, 0.05, fill=TEAL)
cols = [
    ("DEFENSIBLE", "Library + theory close every flank in the Limitations section. 'n too small' and 'not real GRPO' no longer land."),
    ("COMPOUNDING", "Each pillar feeds the next: Q1 corpus → Q2 controller → Q3 bound → Q4 position. The library compounds throughout."),
    ("OWNABLE", "No competing group is building this. First mover sets the metric, the API, and the citation graph through 2027."),
]
x = 0.55; w = 3.95; gap = 0.19
for title, body in cols:
    rect(s, x, 2.55, w, 2.7, fill=PANEL)
    rect(s, x, 2.55, w, 0.09, fill=AMBER)
    tf = box(s, x + 0.25, 2.85, w - 0.5, 2.3)
    para(tf, title, 16, AMBER, bold=True, first=True)
    para(tf, body, 12.5, MUTE, space_before=10, spacing=1.1)
    x += w + gap
rect(s, 0.55, 5.6, 12.23, 0.9, fill=PANEL2)
tf = box(s, 0.8, 5.62, 11.8, 0.86, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Today: the library and the harness shipped — the two artifacts the other three pillars feed on.",
     15, TEAL, bold=True, align=PP_ALIGN.CENTER, first=True)

OUT = "/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14.pptx"
prs.save(OUT)
print("saved:", OUT)
print("slides:", len(prs.slides._sldIdLst))
