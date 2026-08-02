from pptx import Presentation

try:
    prs = Presentation('/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14.pptx')
    if len(prs.slides) >= 17:
        slide = prs.slides[16]
        for i, shape in enumerate(slide.shapes):
            print(f'Shape {i}: {shape.shape_type}')
            if hasattr(shape, 'text'):
                print(f'Text: {shape.text}')
    else:
        print(f'Only {len(prs.slides)} slides found.')
except Exception as e:
    print(f'Error: {e}')
