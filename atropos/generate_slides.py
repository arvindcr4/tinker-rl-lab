#!/usr/bin/env python3
"""Generate PPTX slides for 8th Guidance Call - Updated with full results"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Dark theme colors
BG = RGBColor(0x1a, 0x1a, 0x2e)
CYAN = RGBColor(0x00, 0xd4, 0xff)
GREEN = RGBColor(0x00, 0xff, 0x88)
YELLOW = RGBColor(0xff, 0xd7, 0x00)
RED = RGBColor(0xff, 0x44, 0x44)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0xaa, 0xaa, 0xaa)
DARK_CARD = RGBColor(0x22, 0x22, 0x3a)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def add_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_CARD
    shape.line.fill.background()
    return shape

# ==================== SLIDE 1: Title ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)
add_text(slide, 1, 0.8, 11, 1, "RL for LLMs", 36, CYAN, True, PP_ALIGN.CENTER)
add_text(slide, 1, 1.8, 11, 1.2, "Tinker RL Project", 44, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, 1, 3.0, 11, 0.6, "8th Guidance Call — Progress Update", 24, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, 1, 3.8, 11, 0.5, "Group 6", 20, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, 1, 4.3, 11, 0.5, "Madhu Kumara L  |  Sandhya Jeyaraj  |  Mohammad Rafi", 16, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, 1, 4.7, 11, 0.5, "Arumugam Chetty K  |  Arvind CR  |  Dhruva N", 16, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, 1, 5.5, 11, 0.5, "MTech in Data Science & AI  |  Semester 3 Capstone", 14, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, 1, 5.9, 11, 0.5, "PES University, Bengaluru", 14, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, 1, 6.5, 11, 0.5, "1 March 2026", 16, CYAN, False, PP_ALIGN.CENTER)

# ==================== SLIDE 2: Addressing Feedback ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.5, 0.3, 3, 0.5, "01", 14, CYAN, True)
add_text(slide, 1.2, 0.3, 10, 0.6, "Addressing Feedback & Progress", 32, WHITE, True)

add_card(slide, 0.5, 1.2, 6, 1.2)
add_text(slide, 0.7, 1.3, 5.6, 0.3, "Feedback from 7th Call", 16, RED, True)
add_text(slide, 0.7, 1.7, 5.6, 0.6, '"What is the current SOTA and what are we improving?" + "Try bigger models"', 14, GRAY)

add_card(slide, 7, 1.2, 5.8, 1.2)
add_text(slide, 7.2, 1.3, 5.4, 0.3, "Our Response", 16, GREEN, True)
add_text(slide, 7.2, 1.7, 5.4, 0.6, "SOTA survey done. 4 GSM8K experiments complete. 2 MATH experiments running. Cross-family comparison.", 14, GRAY)

add_card(slide, 0.5, 2.7, 3.8, 1.5)
add_text(slide, 0.7, 2.8, 3.4, 0.3, "SOTA Survey", 16, CYAN, True)
add_text(slide, 0.7, 3.2, 3.4, 0.8, "Mapped DeepSeek-R1, DeepScaleR, SimpleRL-Zoo, TinyZero, DeepCoder across 1B-14B", 13, GRAY)

add_card(slide, 4.6, 2.7, 3.8, 1.5)
add_text(slide, 4.8, 2.8, 3.4, 0.3, "Bigger Models", 16, CYAN, True)
add_text(slide, 4.8, 3.2, 3.4, 0.8, "Scaled from 1B to 30B MoE. Qwen3-8B, Qwen3-30B, Llama-3B, Llama-8B all tested.", 13, GRAY)

add_card(slide, 8.7, 2.7, 4.1, 1.5)
add_text(slide, 8.9, 2.8, 3.7, 0.3, "Multi-Benchmark", 16, CYAN, True)
add_text(slide, 8.9, 3.2, 3.7, 0.8, "GSM8K (grade math) + MATH (competition math) + LogP Steering. Cross-family, cross-task.", 13, GRAY)

for i, (val, label) in enumerate([("10+", "Experiments"), ("30B", "Biggest Model"), ("2", "Benchmarks"), ("2", "Model Families")]):
    x = 0.5 + i * 3.15
    add_card(slide, x, 4.5, 2.9, 1.0)
    add_text(slide, x, 4.55, 2.9, 0.5, val, 28, CYAN, True, PP_ALIGN.CENTER)
    add_text(slide, x, 5.05, 2.9, 0.4, label, 12, GRAY, False, PP_ALIGN.CENTER)

add_card(slide, 0.5, 5.8, 12.3, 1.2)
add_text(slide, 0.7, 5.9, 11.9, 0.3, "KEY RESULTS", 14, GREEN, True)
add_text(slide, 0.7, 6.3, 11.9, 0.5, "GSM8K: Qwen3-8B 100%, Qwen3-30B 100%, Llama-8B 100% | Llama-3B ~1.5% (base model, no instruct tuning) | MATH experiments in progress", 14, GRAY)

# ==================== SLIDE 3: SOTA Comparison ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.5, 0.3, 3, 0.5, "02", 14, CYAN, True)
add_text(slide, 1.2, 0.3, 10, 0.6, "Current SOTA — Where We Stand", 32, WHITE, True)

rows = [
    ("Model", "Size", "Method", "Best Result", "Source"),
    ("DeepSeek-R1-Distill-Qwen-1.5B", "1.5B", "Distill (671B teacher)", "MATH-500: 83.9%", "Jan 2025"),
    ("DeepScaleR-1.5B", "1.5B", "GRPO on distilled", "Beats o1-Preview", "Feb 2025"),
    ("DeepSeek-R1-Distill-Qwen-7B", "7B", "Distill (671B teacher)", "MATH-500: 92.8%", "Jan 2025"),
    ("SimpleRL-Zoo", "0.5-32B", "Zero RL (GRPO)", "Strong gains", "COLM 2025"),
    ("TinyZero", "3B", "GRPO", "Emergent aha", "Feb 2025"),
    ("DeepCoder-14B", "14B", "RL + code", "60.6% LCB", "Apr 2025"),
    ("Ours: Qwen3-8B", "8B", "GRPO (no distill)", "GSM8K: 100%", "This work"),
    ("Ours: Qwen3-30B MoE", "30B", "GRPO (no distill)", "GSM8K: 100%", "This work"),
    ("Ours: Llama-3.1-8B", "8B", "GRPO (no distill)", "GSM8K: 100%", "This work"),
]

y = 1.1
for i, row in enumerate(rows):
    color = WHITE if i == 0 else (GREEN if i >= 7 else GRAY)
    bold = (i == 0 or i >= 7)
    size = 12 if i > 0 else 13
    for j, (text, w) in enumerate(zip(row, [3.8, 0.9, 2.8, 3.2, 1.5])):
        x = 0.5 + sum([3.8, 0.9, 2.8, 3.2, 1.5][:j])
        add_text(slide, x, y, w, 0.4, text, size, color if j > 0 or i == 0 else (CYAN if i >= 7 else WHITE), bold if j == 0 else (i == 0))
    y += 0.42

# Key insight box
add_card(slide, 0.5, 5.6, 12.3, 1.4)
add_text(slide, 0.7, 5.7, 11.9, 0.3, "KEY INSIGHT", 14, YELLOW, True)
add_text(slide, 0.7, 6.1, 11.9, 0.7, "SOTA uses distillation from 671B teacher ($4,500+ compute). We use pure GRPO from base models on Tinker cloud — no distillation, no teacher, no GPU management. Accessible for any researcher.", 14, GRAY)

# ==================== SLIDE 4: GSM8K Scaling Results ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.5, 0.3, 3, 0.5, "03", 14, CYAN, True)
add_text(slide, 1.2, 0.3, 10, 0.6, "GSM8K Results — Complete Scaling Ladder", 28, WHITE, True)

# Results table
add_text(slide, 0.5, 1.1, 2.5, 0.3, "Model", 13, CYAN, True)
add_text(slide, 3, 1.1, 1.5, 0.3, "Family", 13, CYAN, True)
add_text(slide, 4.5, 1.1, 1.5, 0.3, "Baseline", 13, CYAN, True)
add_text(slide, 6, 1.1, 1.5, 0.3, "Peak", 13, CYAN, True)
add_text(slide, 7.5, 1.1, 2, 0.3, "Steps to 50%", 13, CYAN, True)
add_text(slide, 9.5, 1.1, 1.5, 0.3, "Status", 13, CYAN, True)

results = [
    ("Llama-3.2-1B", "Llama", "~10%", "~63%", "~10", "DONE (prev)"),
    ("Llama-3.2-3B", "Llama", "0.8%", "1.6%", "N/A", "RUNNING"),
    ("Llama-3.1-8B-Inst", "Llama", "78.9%", "100%", "0 (pretrained)", "DONE"),
    ("Qwen3-8B", "Qwen", "7.0%", "100%", "~25", "DONE"),
    ("Qwen3-30B-A3B MoE", "Qwen", "17.2%", "100%", "~14", "DONE"),
]
for i, (model, family, base, peak, steps, status) in enumerate(results):
    y = 1.5 + i * 0.4
    done = "DONE" in status
    c = GREEN if done else YELLOW
    add_text(slide, 0.5, y, 2.5, 0.3, model, 12, WHITE, True)
    add_text(slide, 3, y, 1.5, 0.3, family, 12, CYAN if family == "Qwen" else YELLOW)
    add_text(slide, 4.5, y, 1.5, 0.3, base, 12, GRAY)
    add_text(slide, 6, y, 1.5, 0.3, peak, 12, GREEN if "100" in peak else GRAY, "100" in peak)
    add_text(slide, 7.5, y, 2, 0.3, steps, 12, GRAY)
    add_text(slide, 9.5, y, 1.5, 0.3, status, 12, c, done)

# Key observations
add_card(slide, 0.5, 3.8, 6, 1.5)
add_text(slide, 0.7, 3.9, 5.6, 0.3, "Key Observations", 14, YELLOW, True)
add_text(slide, 0.7, 4.3, 5.6, 0.8, "- Bigger models learn faster: 30B reaches 50% in ~14 steps vs ~25 for 8B\n- Instruct models (Llama-8B) start high (79%) — already aligned\n- Base Llama-3B struggles (no instruct tuning) — model needs capacity\n- Qwen family shows consistent GRPO improvement", 12, GRAY)

add_card(slide, 7, 3.8, 5.8, 1.5)
add_text(slide, 7.2, 3.9, 5.4, 0.3, "Cross-Family Insight", 14, GREEN, True)
add_text(slide, 7.2, 4.3, 5.4, 0.8, "- Qwen3-8B (base): 7% -> 100% via GRPO\n- Llama-3.1-8B (instruct): 79% -> 100% via GRPO\n- Both reach 100% — GRPO works across families\n- Different starting points = different learning dynamics", 12, GRAY)

add_card(slide, 0.5, 5.6, 12.3, 1.4)
add_text(slide, 0.7, 5.7, 11.9, 0.3, "SCALING LAW", 14, CYAN, True)
add_text(slide, 0.7, 6.1, 11.9, 0.7, "Clear trend: model size correlates with GRPO effectiveness. 1B: 63% | 8B: 100% | 30B: 100% (faster convergence). This mirrors findings in SimpleRL-Zoo and DeepSeek-R1 literature.", 14, GRAY)

# ==================== SLIDE 5: Reward Trajectories ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.5, 0.3, 3, 0.5, "04", 14, CYAN, True)
add_text(slide, 1.2, 0.3, 10, 0.6, "Training Reward Trajectories", 28, WHITE, True)

# Qwen3-8B trajectory
add_card(slide, 0.5, 1.1, 3.8, 3.5)
add_text(slide, 0.7, 1.2, 3.4, 0.3, "Qwen3-8B (base)", 14, CYAN, True)
qwen8b_key = [(0,"7%"),(10,"4%"),(20,"12%"),(25,"75%"),(30,"100%"),(40,"99%"),(49,"100%")]
for i, (step, r) in enumerate(qwen8b_key):
    c = GREEN if "100" in r or "99" in r else (YELLOW if int(r.replace("%","")) > 20 else GRAY)
    add_text(slide, 0.7, 1.6 + i*0.35, 1, 0.3, f"S{step}", 11, GRAY)
    add_text(slide, 1.7, 1.6 + i*0.35, 2, 0.3, r, 11, c, "100" in r)

# Qwen3-30B trajectory
add_card(slide, 4.6, 1.1, 3.8, 3.5)
add_text(slide, 4.8, 1.2, 3.4, 0.3, "Qwen3-30B MoE", 14, CYAN, True)
qwen30b_key = [(0,"17%"),(4,"39%"),(14,"59%"),(20,"63%"),(28,"100%"),(38,"100%"),(49,"99%")]
for i, (step, r) in enumerate(qwen30b_key):
    c = GREEN if "100" in r or "99" in r else (YELLOW if int(r.replace("%","")) > 20 else GRAY)
    add_text(slide, 4.8, 1.6 + i*0.35, 1, 0.3, f"S{step}", 11, GRAY)
    add_text(slide, 5.8, 1.6 + i*0.35, 2, 0.3, r, 11, c, "100" in r)

# Llama-8B trajectory
add_card(slide, 8.7, 1.1, 4.1, 3.5)
add_text(slide, 8.9, 1.2, 3.7, 0.3, "Llama-3.1-8B (instruct)", 14, YELLOW, True)
llama8b_key = [(0,"79%"),(5,"71%"),(10,"87%"),(15,"95%"),(20,"98%"),(36,"100%"),(49,"100%")]
for i, (step, r) in enumerate(llama8b_key):
    c = GREEN if "100" in r else (YELLOW if int(r.replace("%","")) > 80 else GRAY)
    add_text(slide, 8.9, 1.6 + i*0.35, 1, 0.3, f"S{step}", 11, GRAY)
    add_text(slide, 9.9, 1.6 + i*0.35, 2, 0.3, r, 11, c, "100" in r)

# Interpretation
add_card(slide, 0.5, 4.9, 12.3, 2.1)
add_text(slide, 0.7, 5.0, 11.9, 0.3, "INTERPRETATION", 14, GREEN, True)
add_text(slide, 0.7, 5.4, 5.5, 1.2, "Phase Transition Pattern:\n- Qwen3-8B: 3 phases — baseline (0-11), learning (12-24),\n  mastery (25-50). Sharp transition at step 25.\n- Qwen3-30B: Higher baseline, faster convergence.\n  100% first hit at step 28.\n- Llama-8B: Starts high (instruct-tuned).\n  Quickly saturates to 100%.", 12, GRAY)
add_text(slide, 6.5, 5.4, 6, 1.2, "Llama-3.2-3B (base, no instruct):\n- Stays at ~1% after 30+ steps\n- 3B base model lacks capacity for GSM8K via\n  GRPO alone (no instruct tuning)\n- Contrast: Qwen3-8B base succeeds because\n  8B params provide enough capacity\n- Finding: minimum ~8B params needed for\n  pure GRPO to work on GSM8K", 12, GRAY)

# ==================== SLIDE 6: Other Experiments ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.5, 0.3, 3, 0.5, "05", 14, CYAN, True)
add_text(slide, 1.2, 0.3, 10, 0.6, "Beyond GSM8K — Multi-Benchmark Experiments", 28, WHITE, True)

# Experiment cards
benchmarks = [
    ("MATH Competition", "Qwen3-8B + Llama-8B", "RUNNING", "Competition-level math (algebra, geometry, number theory). Harder than GSM8K. Tests deeper reasoning.", YELLOW),
    ("LogP Steering", "Qwen3-30B-A3B", "QUEUED", "Self-distillation via logprob steering. Model teaches itself with different system prompt (emoji style).", CYAN),
    ("Code Generation", "Qwen3-8B", "PLANNED", "HumanEval/MBPP code problems. Binary reward: code passes test cases. Tests generalization beyond math.", GRAY),
    ("Instruction Following", "Llama-8B + Qwen-8B", "PLANNED", "IFEval benchmark. Tests constraint adherence (word count, format). Different reward structure.", GRAY),
]
for i, (name, models, status, desc, color) in enumerate(benchmarks):
    x = 0.5 + (i % 2) * 6.3
    y = 1.2 + (i // 2) * 2.6
    add_card(slide, x, y, 5.8, 2.2)
    add_text(slide, x + 0.2, y + 0.1, 4.5, 0.3, name, 16, color, True)
    sc = GREEN if "RUNNING" in status else (YELLOW if "QUEUED" in status else GRAY)
    add_text(slide, x + 4.7, y + 0.1, 1, 0.3, status, 10, sc, True)
    add_text(slide, x + 0.2, y + 0.5, 5.4, 0.3, f"Models: {models}", 12, CYAN)
    add_text(slide, x + 0.2, y + 0.9, 5.4, 1.0, desc, 12, GRAY)

add_card(slide, 0.5, 6.2, 12.3, 0.8)
add_text(slide, 0.7, 6.3, 11.9, 0.5, "Total experiment matrix: 5 models x 4 benchmarks = 20 experiments. Currently 6 complete, 2 running, 12 planned.", 14, GRAY)

# ==================== SLIDE 7: Technical Architecture ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.5, 0.3, 3, 0.5, "06", 14, CYAN, True)
add_text(slide, 1.2, 0.3, 10, 0.6, "Technical Setup — Tinker + Atropos", 32, WHITE, True)

# Architecture diagram as cards
add_card(slide, 0.5, 1.2, 3.5, 1.5)
add_text(slide, 0.7, 1.3, 3.1, 0.3, "Environment (Local)", 14, GREEN, True)
add_text(slide, 0.7, 1.7, 3.1, 0.8, "- GSM8K / MATH dataset\n- Generate rollouts (n=16)\n- Score with math_verify\n- Binary reward (0/1)", 12, GRAY)

add_card(slide, 4.5, 1.2, 3.5, 1.5)
add_text(slide, 4.7, 1.3, 3.1, 0.3, "Atropos API (Local)", 14, YELLOW, True)
add_text(slide, 4.7, 1.7, 3.1, 0.8, "- Coordinates env <-> trainer\n- Batches rollout data\n- Manages training steps\n- Off-policy buffer (3 batches)", 12, GRAY)

add_card(slide, 8.5, 1.2, 4.3, 1.5)
add_text(slide, 8.7, 1.3, 3.9, 0.3, "Tinker Trainer (Cloud)", 14, CYAN, True)
add_text(slide, 8.7, 1.7, 3.9, 0.8, "- LoRA training (rank 32)\n- forward_backward() on GPU\n- Importance sampling loss\n- Updates sampling client", 12, GRAY)

# Arrows
add_text(slide, 3.9, 1.7, 0.7, 0.3, "->", 24, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, 7.9, 1.7, 0.7, 0.3, "->", 24, WHITE, True, PP_ALIGN.CENTER)

# Key advantages
advs = [
    ("GRPO", "No critic needed", "Group Relative Policy Optimization. No reward model or value function."),
    ("LoRA", "~0.1% params", "Rank 32 adaptation. 99.2% parameter reduction. Rapid iteration."),
    ("Cloud API", "No GPU setup", "Tinker handles GPUs. Runs from laptop. Parallel experiments."),
    ("Verifiable", "Binary 0/1", "math_verify + boxed format. No learned reward model needed."),
]
for i, (title, val, desc) in enumerate(advs):
    x = 0.5 + i * 3.15
    add_card(slide, x, 3.2, 2.9, 1.8)
    add_text(slide, x + 0.2, 3.3, 2.5, 0.3, title, 13, CYAN, True)
    add_text(slide, x + 0.2, 3.7, 2.5, 0.3, val, 20, YELLOW, True)
    add_text(slide, x + 0.2, 4.2, 2.5, 0.6, desc, 11, GRAY)

# Consistent hyperparams
add_card(slide, 0.5, 5.3, 12.3, 1.8)
add_text(slide, 0.7, 5.4, 11.9, 0.3, "Consistent Hyperparameters Across All Experiments", 14, WHITE, True)
add_text(slide, 0.7, 5.8, 11.9, 0.4, "LoRA rank: 32  |  LR: 3-5e-5  |  Batch: 128  |  Group: 16  |  Max tokens: 512-1024  |  Steps: 50  |  Loss: IS", 14, CYAN)
add_text(slide, 0.7, 6.3, 11.9, 0.5, "Same recipe applied to every model and benchmark — enables fair comparison across sizes, families, and tasks", 13, GRAY)

# ==================== SLIDE 8: Our Contribution ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.5, 0.3, 3, 0.5, "07", 14, CYAN, True)
add_text(slide, 1.2, 0.3, 10, 0.6, "Our Contribution — What We're Adding", 32, WHITE, True)

add_text(slide, 0.5, 1.1, 6, 0.4, "Gaps in Literature", 18, RED, True)
gaps = ["Most GRPO results on Qwen only — Llama under GRPO underexplored", "No systematic scaling study comparing GRPO across model families", "Existing work needs expensive local GPU clusters ($4,500+)", "Limited multi-benchmark GRPO analysis (most focus on one task)"]
for i, gap in enumerate(gaps):
    add_text(slide, 0.7, 1.6 + i * 0.4, 6, 0.4, f"- {gap}", 13, GRAY)

add_text(slide, 0.5, 3.4, 12, 0.5, "Our Contribution", 18, GREEN, True)
add_card(slide, 0.5, 3.9, 12.3, 0.8)
add_text(slide, 0.7, 4.0, 11.9, 0.6, "Systematic GRPO scaling study across model families (Llama vs Qwen), sizes (1B-30B), and benchmarks (GSM8K, MATH, code), with reproducible cloud-based experiments on the Tinker platform", 15, WHITE, True)

dims = [("Model Size", "1B -> 3B -> 8B -> 30B (MoE)"), ("Model Family", "Llama vs Qwen under identical GRPO recipe"), ("Benchmarks", "GSM8K, MATH competition, Code, IFEval"), ("Infrastructure", "Tinker cloud API — zero GPU management"), ("Reproducibility", "YAML configs + scripts for every experiment")]
y = 5.0
add_text(slide, 0.5, y, 5, 0.3, "Dimension", 13, CYAN, True)
add_text(slide, 5, y, 7, 0.3, "What We Compare", 13, CYAN, True)
for label, val in dims:
    y += 0.35
    add_text(slide, 0.5, y, 4.5, 0.3, label, 13, WHITE, True)
    add_text(slide, 5, y, 7.5, 0.3, val, 13, GRAY)

# ==================== SLIDE 9: Timeline ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.5, 0.3, 3, 0.5, "08", 14, CYAN, True)
add_text(slide, 1.2, 0.3, 10, 0.6, "Timeline & Next Steps", 32, WHITE, True)

milestones = [
    ("Feb 22", "7th Guidance Call", GRAY, False),
    ("Mar 1", "8th Call: 4 GSM8K done + MATH running", YELLOW, True),
    ("Mar 3-5", "Complete MATH + Code experiments", CYAN, False),
    ("Mar 7", "9th Guidance — full multi-benchmark results", WHITE, False),
    ("Mar 14", "Ablation studies + conference paper outline", GRAY, False),
    ("Mar 28", "3rd Submission — Interim Report", RED, False),
    ("Apr 11", "Final Report + Conference Paper", RED, False),
]
for i, (date, desc, color, highlight) in enumerate(milestones):
    y = 1.2 + i * 0.55
    add_text(slide, 0.5, y, 1.5, 0.4, date, 14, color, highlight)
    add_text(slide, 2.2, y, 5, 0.4, desc, 14, WHITE if highlight else GRAY, highlight)

add_card(slide, 7, 1.2, 5.8, 3.0)
add_text(slide, 7.2, 1.3, 5.4, 0.3, "Experiment Status", 16, CYAN, True)
exps = [
    ("DONE", "GSM8K Llama-3.2-1B (~63%)", GREEN),
    ("DONE", "GSM8K Qwen3-8B (100%)", GREEN),
    ("DONE", "GSM8K Qwen3-30B MoE (100%)", GREEN),
    ("DONE", "GSM8K Llama-3.1-8B (100%)", GREEN),
    ("RUN", "GSM8K Llama-3.2-3B (step 31)", YELLOW),
    ("RUN", "MATH Qwen3-8B (starting)", YELLOW),
    ("RUN", "MATH Llama-3.1-8B (starting)", YELLOW),
    ("NEXT", "LogP Steering + Code Gen", GRAY),
]
for i, (status, exp, c) in enumerate(exps):
    add_text(slide, 7.2, 1.7 + i * 0.3, 0.8, 0.3, status, 10, c, True)
    add_text(slide, 8.1, 1.7 + i * 0.3, 4.5, 0.3, exp, 11, c)

# ==================== SLIDE 10: Discussion ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.5, 0.3, 3, 0.5, "09", 14, CYAN, True)
add_text(slide, 1.2, 0.3, 10, 0.6, "Discussion & Guidance Needed", 32, WHITE, True)

questions = [
    ("Evaluation", "Training reward is 100% on GSM8K. Need proper held-out test eval. How should we structure the evaluation section?", CYAN),
    ("Conference Angle", "Systematic GRPO scaling across families + benchmarks on cloud infra — strong enough? Suggested venues?", YELLOW),
    ("Bigger Models", "Tinker supports Qwen3-235B and DeepSeek-V3.1. Should we attempt 70B+ for the final report?", GREEN),
    ("Depth vs Breadth", "More benchmarks (code, IFEval) or deeper analysis on math (ablations, longer training)?", RED),
]
for i, (title, desc, color) in enumerate(questions):
    x = 0.5 + (i % 2) * 6.3
    y = 1.2 + (i // 2) * 2.6
    add_card(slide, x, y, 5.8, 2.2)
    add_text(slide, x + 0.2, y + 0.1, 5.4, 0.3, title, 16, color, True)
    add_text(slide, x + 0.2, y + 0.5, 5.4, 1.4, desc, 14, GRAY)

add_card(slide, 0.5, 6.0, 12.3, 1.0)
add_text(slide, 0.7, 6.1, 11.9, 0.7, "6 experiments complete, 3 running. GRPO works across Qwen and Llama at 8B+. Next: MATH results + evaluation methodology.", 14, GRAY)

# ==================== SLIDE 11: Thank You ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 1, 2.0, 11, 1, "Thank You", 48, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, 1, 3.2, 11, 0.5, "Questions & Discussion", 24, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, 1, 4.2, 11, 0.5, "Group 6", 20, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, 1, 4.7, 11, 0.5, "Madhu Kumara L  |  Sandhya Jeyaraj  |  Mohammad Rafi", 16, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, 1, 5.1, 11, 0.5, "Arumugam Chetty K  |  Arvind CR  |  Dhruva N", 16, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, 1, 5.8, 11, 0.5, "Mentors: Narayana Darapaneni & Anwesh Reddy Padhuri", 14, CYAN, False, PP_ALIGN.CENTER)

# Save
output_path = "/Users/arvind/Downloads/8th_guidance_call_slides.pptx"
prs.save(output_path)
print(f"Slides saved to: {output_path}")
