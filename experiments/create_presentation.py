#!/usr/bin/env python3
"""Create PowerPoint presentation for Tinker RL demo with TikZ-style diagrams."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Color scheme
DARK_BLUE = RGBColor(0x1a, 0x1a, 0x2e)
LIGHT_BLUE = RGBColor(0x41, 0x69, 0xe1)
GREEN = RGBColor(0x28, 0xa7, 0x45)
ORANGE = RGBColor(0xff, 0x8c, 0x00)
RED = RGBColor(0xdc, 0x35, 0x45)
PURPLE = RGBColor(0x6f, 0x42, 0xc1)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xf0, 0xf0, 0xf0)
WHITE = RGBColor(255, 255, 255)

def add_box(slide, left, top, width, height, text, fill_color, text_color=WHITE, font_size=14, bold=True):
    """Add a rounded rectangle box with text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_arrow_right(slide, x, y, length=0.8, color=GRAY):
    """Add a right-pointing arrow."""
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(length), Inches(0.25))
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()
    return arr

def add_arrow_down(slide, x, y, length=0.8, color=GRAY):
    """Add a down-pointing arrow."""
    arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(0.25), Inches(length))
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()
    return arr

def add_arrow_left(slide, x, y, length=0.8, color=GRAY):
    """Add a left-pointing arrow."""
    arr = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(x), Inches(y), Inches(length), Inches(0.25))
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()
    return arr

def add_arrow_up(slide, x, y, length=0.8, color=GRAY):
    """Add an up-pointing arrow."""
    arr = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(x), Inches(y), Inches(0.25), Inches(length))
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()
    return arr

def add_label(slide, x, y, text, font_size=11, color=GRAY, italic=True):
    """Add a text label."""
    label = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(1.5), Inches(0.4))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.italic = italic
    p.font.color.rgb = color
    return label

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1.2))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    bullet_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.5), Inches(5.5))
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.space_after = Pt(12)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return slide

def add_tinker_architecture_diagram(prs):
    """Diagram: How Tinker Works - Cloud Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "How Tinker Works: Cloud Architecture"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # User/Client side (left)
    add_box(slide, 0.3, 1.5, 2.2, 1.3, "Your Code\n(Python SDK)", LIGHT_GRAY, DARK_BLUE, 13)
    add_box(slide, 0.3, 3.2, 2.2, 0.8, "API Key", ORANGE, WHITE, 12)

    # Cloud boundary
    cloud = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.2), Inches(1.2), Inches(6.5), Inches(5))
    cloud.fill.solid()
    cloud.fill.fore_color.rgb = RGBColor(0xe8, 0xf4, 0xfc)
    cloud.line.color.rgb = LIGHT_BLUE
    cloud.line.width = Pt(2)

    cloud_label = slide.shapes.add_textbox(Inches(5.5), Inches(1.3), Inches(2), Inches(0.4))
    tf = cloud_label.text_frame
    p = tf.paragraphs[0]
    p.text = "☁️ Tinker Cloud"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE

    # Tinker components
    add_box(slide, 3.5, 2, 2, 1, "API Gateway\n& Auth", LIGHT_BLUE, WHITE, 12)
    add_box(slide, 6, 2, 2, 1, "Model\nRegistry", PURPLE, WHITE, 12)
    add_box(slide, 3.5, 3.5, 2, 1.2, "Training\nOrchestrator", GREEN, WHITE, 12)
    add_box(slide, 6, 3.5, 2, 1.2, "GPU\nCluster", ORANGE, WHITE, 12)
    add_box(slide, 4.75, 5.2, 2.5, 0.9, "Checkpoint\nStorage", GRAY, WHITE, 12)

    # Arrows
    add_arrow_right(slide, 2.5, 2, 0.9)
    add_arrow_right(slide, 5.5, 2.35, 0.4)
    add_arrow_down(slide, 4.35, 3.05, 0.4)
    add_arrow_right(slide, 5.5, 3.95, 0.4)
    add_arrow_down(slide, 5.85, 4.75, 0.4)

    # Labels
    add_label(slide, 2.5, 1.6, "requests", 10)
    add_label(slide, 5.3, 1.95, "load", 10)
    add_label(slide, 4.5, 3.1, "schedule", 10)
    add_label(slide, 5.5, 3.55, "train", 10)
    add_label(slide, 6.1, 4.8, "save", 10)

    # Caption
    caption = slide.shapes.add_textbox(Inches(0.3), Inches(6.4), Inches(9.4), Inches(0.5))
    tf = caption.text_frame
    p = tf.paragraphs[0]
    p.text = "Tinker handles GPU allocation, model loading, and checkpointing - you just write training code"
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_rl_loop_diagram(prs):
    """Diagram: RL Training Loop."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Reinforcement Learning Training Loop"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Central LLM
    add_box(slide, 4, 2.8, 2, 1.2, "LLM\n(Policy π)", LIGHT_BLUE, WHITE, 14)

    # Environment
    add_box(slide, 7.5, 2.8, 2, 1.2, "Environment\n(Task)", GREEN, WHITE, 14)

    # Reward
    add_box(slide, 7.5, 5, 2, 0.9, "Reward r(s,a)", ORANGE, WHITE, 13)

    # Optimizer
    add_box(slide, 0.5, 2.8, 2, 1.2, "Optimizer\n(LoRA SGD)", PURPLE, WHITE, 14)

    # Arrows forming the loop
    add_arrow_right(slide, 6, 3.25, 1.4)  # LLM -> Env
    add_arrow_down(slide, 8.35, 4.05, 0.85)  # Env -> Reward
    add_arrow_left(slide, 2.6, 5.2, 4.8)  # Reward -> back
    add_arrow_up(slide, 1.35, 4.05, 0.65)  # Up to optimizer
    add_arrow_right(slide, 2.5, 3.25, 1.4)  # Optimizer -> LLM

    # Labels on arrows
    add_label(slide, 6.3, 2.85, "action a", 11)
    add_label(slide, 8.6, 4.4, "feedback", 11)
    add_label(slide, 4.5, 4.85, "gradient ∇J", 11)
    add_label(slide, 2.8, 2.85, "θ update", 11)

    # State indicator
    state_box = slide.shapes.add_textbox(Inches(6.3), Inches(3.7), Inches(1.2), Inches(0.4))
    tf = state_box.text_frame
    p = tf.paragraphs[0]
    p.text = "state s'"
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = GRAY

    # Formula
    formula = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.5))
    tf = formula.text_frame
    p = tf.paragraphs[0]
    p.text = "Objective: J(θ) = E[Σ r(s,a)] → maximize expected reward"
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_distillation_diagram(prs):
    """Diagram: Knowledge Distillation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Knowledge Distillation: Teacher → Student"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Teacher (top)
    add_box(slide, 3.5, 1.3, 3, 1.2, "Teacher Model\n(Large, Frozen)", PURPLE, WHITE, 13)

    # Student (bottom)
    add_box(slide, 3.5, 4.8, 3, 1.2, "Student Model\n(Small, Training)", LIGHT_BLUE, WHITE, 13)

    # Input
    add_box(slide, 0.3, 3, 1.8, 0.9, "Input x", LIGHT_GRAY, DARK_BLUE, 12)

    # KL Divergence (center)
    kl = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.25), Inches(2.9), Inches(1.5), Inches(1.2))
    kl.fill.solid()
    kl.fill.fore_color.rgb = ORANGE
    kl.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    tf = kl.text_frame
    p = tf.paragraphs[0]
    p.text = "KL\nLoss"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Output distributions
    add_box(slide, 7.5, 1.6, 2, 0.8, "P_T(y|x)", PURPLE, WHITE, 12)
    add_box(slide, 7.5, 5.1, 2, 0.8, "P_S(y|x)", LIGHT_BLUE, WHITE, 12)

    # Arrows
    add_arrow_right(slide, 2.1, 1.8, 1.3)  # Input -> Teacher
    add_arrow_right(slide, 2.1, 5.25, 1.3)  # Input -> Student
    add_arrow_down(slide, 4.85, 2.55, 0.3)  # Teacher -> KL
    add_arrow_up(slide, 4.85, 4.15, 0.6)  # Student -> KL
    add_arrow_right(slide, 6.5, 1.85, 0.9)  # Teacher -> Output
    add_arrow_right(slide, 6.5, 5.35, 0.9)  # Student -> Output

    # Formula
    formula = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(0.5))
    tf = formula.text_frame
    p = tf.paragraphs[0]
    p.text = "Loss = KL(P_T || P_S) = Σ P_T(y) log(P_T(y)/P_S(y))"
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_preference_diagram(prs):
    """Diagram: Pairwise Preference Learning."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Pairwise Preference Learning"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Input
    add_box(slide, 0.2, 2.3, 1.8, 1, "Prompt", LIGHT_GRAY, DARK_BLUE, 12)

    # LLM samples
    add_box(slide, 2.5, 2.3, 1.6, 1, "LLM\n×4", LIGHT_BLUE, WHITE, 12)

    # Responses (bracket style)
    responses = [
        (4.6, 1.3, "A: 120 tok", RED),
        (4.6, 2.3, "B: 80 tok", ORANGE),
        (4.6, 3.3, "C: 50 tok", RGBColor(0x90, 0xee, 0x90)),
        (4.6, 4.3, "D: 35 tok", GREEN),
    ]
    for x, y, text, color in responses:
        text_color = WHITE if color not in [RGBColor(0x90, 0xee, 0x90)] else DARK_BLUE
        add_box(slide, x, y, 1.8, 0.8, text, color, text_color, 11)

    # Comparator
    comp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7), Inches(2.5), Inches(1.4), Inches(1.2))
    comp.fill.solid()
    comp.fill.fore_color.rgb = PURPLE
    tf = comp.text_frame
    p = tf.paragraphs[0]
    p.text = "Compare\nPairs"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Result
    add_box(slide, 8.7, 2.5, 1.1, 1.2, "Winner\n→ D", GREEN, WHITE, 12)

    # Arrows
    add_arrow_right(slide, 2, 2.65, 0.4)
    add_arrow_right(slide, 4.1, 2.65, 0.4)
    add_arrow_right(slide, 6.4, 2.95, 0.5)
    add_arrow_right(slide, 8.4, 2.95, 0.25)

    # Win/Loss labels
    win_label = slide.shapes.add_textbox(Inches(6.5), Inches(4.3), Inches(0.8), Inches(0.3))
    tf = win_label.text_frame
    p = tf.paragraphs[0]
    p.text = "✓ Win"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = GREEN

    lose_label = slide.shapes.add_textbox(Inches(6.5), Inches(1.3), Inches(0.8), Inches(0.3))
    tf = lose_label.text_frame
    p = tf.paragraphs[0]
    p.text = "✗ Lose"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RED

    # Reward formula
    formula = slide.shapes.add_textbox(Inches(0.3), Inches(5.5), Inches(9.4), Inches(0.9))
    tf = formula.text_frame
    p = tf.paragraphs[0]
    p.text = "reward = (wins - losses) / comparisons + format_bonus"
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Shorter responses beat longer ones → model learns conciseness"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_lora_diagram(prs):
    """Diagram: LoRA Fine-tuning."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "LoRA: Low-Rank Adaptation"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Original weight matrix
    add_box(slide, 0.5, 2, 2.5, 2.5, "W₀\n(Frozen)\n\nd × d", GRAY, WHITE, 14)

    # Plus sign
    plus = slide.shapes.add_textbox(Inches(3.1), Inches(3), Inches(0.5), Inches(0.5))
    tf = plus.text_frame
    p = tf.paragraphs[0]
    p.text = "+"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # LoRA matrices
    add_box(slide, 3.8, 1.8, 0.8, 2.5, "B\nd×r", LIGHT_BLUE, WHITE, 12)
    add_box(slide, 4.8, 2.5, 2, 0.8, "A (r×d)", LIGHT_BLUE, WHITE, 12)

    # Multiplication
    mult = slide.shapes.add_textbox(Inches(4.5), Inches(2), Inches(0.5), Inches(0.4))
    tf = mult.text_frame
    p = tf.paragraphs[0]
    p.text = "×"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Equals
    eq = slide.shapes.add_textbox(Inches(7), Inches(3), Inches(0.5), Inches(0.5))
    tf = eq.text_frame
    p = tf.paragraphs[0]
    p.text = "="
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Result
    add_box(slide, 7.5, 2, 2, 2.5, "W₀ + BA\n(Effective)\n\nd × d", GREEN, WHITE, 14)

    # Annotations
    annotations = [
        (0.5, 4.7, "Original model\nweights (billions)", GRAY),
        (3.8, 4.7, "LoRA adapters\n(millions)", LIGHT_BLUE),
        (7.5, 4.7, "Combined for\ninference", GREEN),
    ]
    for x, y, text, color in annotations:
        ann = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.2), Inches(0.8))
        tf = ann.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

    # Key insight
    insight = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.8))
    tf = insight.text_frame
    p = tf.paragraphs[0]
    p.text = "rank r << d: Train only ~0.1% of parameters while preserving model capabilities"
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Example: d=4096, r=32 → 99.2% fewer trainable parameters"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_results_chart(prs):
    """Visual chart of Math RL results."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Math RL Results: Learning Curve"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Y-axis
    y_axis = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.3), Inches(0.03), Inches(4.2))
    y_axis.fill.solid()
    y_axis.fill.fore_color.rgb = DARK_BLUE
    y_axis.line.fill.background()

    # X-axis
    x_axis = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(5.5), Inches(5.5), Inches(0.03))
    x_axis.fill.solid()
    x_axis.fill.fore_color.rgb = DARK_BLUE
    x_axis.line.fill.background()

    # Y labels
    y_labels = [("1.0", 1.3), ("0.8", 2.15), ("0.6", 3.0), ("0.4", 3.85), ("0.2", 4.7)]
    for label, y in y_labels:
        lbl = slide.shapes.add_textbox(Inches(1.0), Inches(y - 0.1), Inches(0.5), Inches(0.3))
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY

    # X labels
    x_labels = [("0", 1.5), ("5", 2.6), ("10", 3.7), ("15", 4.8), ("20", 5.9)]
    for label, x in x_labels:
        lbl = slide.shapes.add_textbox(Inches(x - 0.1), Inches(5.6), Inches(0.5), Inches(0.3))
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY

    # Axis labels
    y_title = slide.shapes.add_textbox(Inches(0.3), Inches(3.0), Inches(0.8), Inches(0.5))
    tf = y_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Reward"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    x_title = slide.shapes.add_textbox(Inches(3.8), Inches(5.9), Inches(1.5), Inches(0.4))
    tf = x_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Step"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Data points
    points = [(1.7, 4.0, "0.67"), (2.8, 2.6, "0.85"), (3.9, 1.9, "0.95"), (5.0, 1.55, "0.99"), (6.1, 1.35, "1.00")]

    for x, y, val in points:
        pt = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.1), Inches(y - 0.1), Inches(0.2), Inches(0.2))
        pt.fill.solid()
        pt.fill.fore_color.rgb = GREEN
        pt.line.color.rgb = DARK_BLUE

        lbl = slide.shapes.add_textbox(Inches(x - 0.2), Inches(y - 0.45), Inches(0.5), Inches(0.3))
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = GREEN

    # Annotation box
    ann_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(1.5), Inches(2.4), Inches(2.5))
    ann_box.fill.solid()
    ann_box.fill.fore_color.rgb = LIGHT_GRAY
    ann_box.line.color.rgb = GRAY

    ann = slide.shapes.add_textbox(Inches(7.4), Inches(1.6), Inches(2.2), Inches(2.3))
    tf = ann.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Key Results:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    results = ["• 69.5% → 100%", "• Converged in 20 steps", "• ~2 min training", "• Perfect accuracy"]
    for r in results:
        p = tf.add_paragraph()
        p.text = r
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
    return slide

def add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    cols = len(headers)
    table = slide.shapes.add_table(len(rows) + 1, cols, Inches(0.3), Inches(1.2), Inches(9.4), Inches(0.5 * (len(rows) + 1))).table

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BLUE

    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
    return slide

def add_code_slide(prs, title, code):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    code_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(0.9), Inches(9.4), Inches(5.8))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = RGBColor(0x1e, 0x1e, 0x1e)
    code_bg.line.fill.background()

    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(5.4))
    tf = code_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(12)
    p.font.name = "Courier New"
    p.font.color.rgb = RGBColor(0x00, 0xff, 0x00)
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Title
    add_title_slide(prs,
        "Reinforcement Learning for LLMs",
        "Training Language Models with Tinker\n\nPES LLM Research Project")

    # 2. What is Tinker
    add_content_slide(prs, "What is Tinker?", [
        "Cloud-based RL training platform for LLMs",
        "Supports LoRA fine-tuning with reinforcement learning",
        "Pre-built recipes: Math RL, Preference, Distillation",
        "No GPU required locally - training runs in the cloud",
        "Python SDK with simple API"
    ])

    # 3. Tinker Architecture Diagram
    add_tinker_architecture_diagram(prs)

    # 4. LoRA Diagram
    add_lora_diagram(prs)

    # 5. RL Loop Diagram
    add_rl_loop_diagram(prs)

    # 6. Distillation Diagram
    add_distillation_diagram(prs)

    # 7. Preference Diagram
    add_preference_diagram(prs)

    # 8. Experiments Table
    add_table_slide(prs, "Experiments Conducted",
        ["Recipe", "Task", "Model", "Status", "Result"],
        [
            ["Math RL", "Arithmetic", "Llama-3.2-1B", "✓ Complete", "100% accuracy"],
            ["Chat SL", "Conversation", "Llama-3.2-1B", "✓ Running", "NLL: 2.5→1.8"],
            ["Preference", "Shorter", "Qwen-0.6B", "✓ Running", "37% done"],
            ["Distill (Off)", "Reasoning", "Llama-3.2-1B", "✓ Running", "SFT active"],
            ["Distill (On)", "KL Match", "Llama-3.2-1B", "✓ Running", "KL: 2.88"],
            ["GSM8K", "Word Problems", "Llama-3.2-1B", "✓ Running", "Multi-step"],
        ])

    # 9. Results Chart
    add_results_chart(prs)

    # 10. Code Example
    add_code_slide(prs, "Example: Math RL Training",
"""# Install Tinker
pip install tinker tinker-cookbook

# Set API key
export TINKER_API_KEY="your-key-here"

# Run RL training
python -m tinker_cookbook.recipes.math_rl.train \\
    model_name="meta-llama/Llama-3.2-1B" \\
    env=arithmetic \\
    group_size=4 \\
    groups_per_batch=100 \\
    learning_rate=1e-4

# Output:
# Step 0:  reward=0.676, correct=69.5%
# Step 10: reward=0.998, correct=99.8%
# Step 20: reward=1.000, correct=100%""")

    # 11. Notebooks
    add_content_slide(prs, "Jupyter Notebooks Created", [
        "01_math_rl_arithmetic.ipynb - Basic RL for addition",
        "02_chat_sl_sft.ipynb - Supervised fine-tuning",
        "03_preference_shorter.ipynb - Pairwise preference learning",
        "04_distillation_off_policy.ipynb - SFT distillation",
        "05_distillation_on_policy.ipynb - KL distillation",
        "06_math_rl_gsm8k.ipynb - Word problem solving"
    ])

    # 12. Key Findings
    add_content_slide(prs, "Key Findings", [
        "Simple tasks (arithmetic): RL achieves 100% in ~20 steps",
        "Complex tasks (GSM8K): Requires more training, larger models",
        "Preference learning: Effective for stylistic control",
        "Distillation: Efficient knowledge transfer",
        "Cloud training: Accessible without local GPUs"
    ])

    # 13. Future Work
    add_content_slide(prs, "Future Directions", [
        "Code generation with sandbox (DeepCoder)",
        "Multi-agent RL training",
        "RLHF with human preference data",
        "Scaling to larger models (8B, 32B)",
        "Custom reward functions"
    ])

    # 14. Thank You
    add_title_slide(prs, "Thank You!", "Questions?\n\npes-llm-research/tinker-experiments")

    output_path = '/home/ubuntu/pes-llm-research/tinker-experiments/Tinker_RL_Demo.pptx'
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
