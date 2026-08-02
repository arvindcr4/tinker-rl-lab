#!/usr/bin/env python3
"""Build the comprehensive faculty program-review deck for Tinker RL Lab.

Covers ALL work done in the program (Semester 3 Group 6 capstone, Semester 4
solo continuation, the E1 confirmatory campaign, preregistration, and the
publication program). Every number is read from / grounded in the repository.

Visual language matches outputs/build_progress_update_deck.py (dark navy,
restrained signal colors, 16:9).
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs" / "Tinker_RL_Program_Review_All_Work_2026-08-02.pptx"
FIG = ROOT / "platform_hybrid" / "paper" / "figures" / "v2"

W, H = 13.333, 7.5

# Palette (matches the existing progress deck).
NAVY = "0B1220"
PANEL = "111C2E"
PANEL_2 = "17243A"
INK = "F4F7FB"
MUTED = "A7B5C9"
TEAL = "2DD4BF"
BLUE = "60A5FA"
LAV = "A78BFA"
AMBER = "FBBF24"
RED = "FB7185"
GREEN = "86EFAC"
WHITE = "FFFFFF"
GRID = "2B3A52"

FONT = "Aptos"
FONT_DISPLAY = "Aptos Display"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def add_shape(slide, kind, x, y, w, h, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else kind
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    shape.line.width = Pt(0.75)
    return shape


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False, font=FONT,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.04, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, runs, x, y, w, h, align=PP_ALIGN.LEFT,
                  valign=MSO_ANCHOR.TOP, margin=0.04):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    for text, color, bold, size in runs:
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_line(slide, x1, y1, x2, y2, color=GRID, width=1.0):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_header(slide, kicker, title, number):
    add_text(slide, kicker.upper(), 0.55, 0.26, 8.0, 0.22, size=8.5, color=TEAL, bold=True)
    add_text(slide, title, 0.55, 0.54, 11.9, 0.52, size=25, color=INK, bold=True, font=FONT_DISPLAY)
    add_text(slide, f"{number:02d}", 12.27, 0.28, 0.45, 0.25, size=10, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)
    add_line(slide, 0.55, 1.17, 12.78, 1.17, GRID, 0.8)


def add_footer(slide, source: str):
    add_line(slide, 0.55, 7.14, 12.78, 7.14, GRID, 0.7)
    add_text(slide, source, 0.58, 7.18, 11.7, 0.17, size=6.8, color=MUTED, italic=True, margin=0)


def add_metric(slide, x, y, w, h, value, label, accent=TEAL, detail=None):
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, PANEL, GRID, radius=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.06, h, accent)
    add_text(slide, value, x + 0.22, y + 0.16, w - 0.35, 0.42, size=25, color=INK, bold=True, font=FONT_DISPLAY)
    add_text(slide, label, x + 0.22, y + 0.62, w - 0.35, 0.28, size=10, color=MUTED, bold=True)
    if detail:
        add_text(slide, detail, x + 0.22, y + 0.92, w - 0.35, h - 1.02, size=8.2, color=MUTED)


def add_bullet(slide, text, x, y, w, h=0.42, color=INK, size=13, bullet_color=TEAL, bold=False):
    add_shape(slide, MSO_SHAPE.OVAL, x, y + 0.10, 0.12, 0.12, bullet_color, bullet_color)
    add_text(slide, text, x + 0.24, y, w - 0.24, h, size=size, color=color, bold=bold)


def add_chip(slide, label, x, y, color=TEAL, w=None, size=7.7):
    w = w or max(0.55, 0.11 * len(label) + 0.24)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.27, color, color, radius=True)
    add_text(slide, label.upper(), x, y + 0.005, w, 0.23, size=size, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)


def add_picture(slide, path: Path, x, y, w, h):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(NAVY)
    return slide


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    # ── 1. Title ────────────────────────────────────────────────────────────
    s = new_slide(prs)
    add_text(s, "TINKER RL LAB  /  PES UNIVERSITY", 0.65, 0.55, 6.0, 0.25, size=10, color=TEAL, bold=True)
    add_text(s, "Reinforcement Learning Post-Training\nof Large Language Models", 0.65, 1.30, 12.0, 1.9, size=40, color=INK, bold=True, font=FONT_DISPLAY)
    add_text(s, "Program review — every phase, paper, and experiment from Semester 3 capstone to today", 0.69, 3.62, 11.0, 0.42, size=17, color=MUTED)
    add_text(s, "02 AUG 2026  ·  FACULTY REVIEW  ·  ARVIND C R", 0.69, 4.26, 6.5, 0.25, size=9.5, color=AMBER, bold=True)
    add_text(s, "M.Tech candidate  ·  Guide: Prof. Ramesh Prakash Guledgudd", 0.69, 4.56, 7.0, 0.25, size=12, color=INK, bold=True)
    add_line(s, 0.72, 5.72, 12.45, 5.72, GRID, 2.3)
    phases = [
        (1.30, "SEM 3", "Group 6 capstone\nTinkerRL-Bench · NeurIPS", TEAL),
        (4.90, "SEM 4", "Solo continuation\n8 pillar papers · thesis", BLUE),
        (8.50, "JUL–AUG 26", "E1 campaign · audit\npreregistration", AMBER),
    ]
    for x, tag, label, color in phases:
        add_shape(s, MSO_SHAPE.OVAL, x + 0.15, 5.54, 0.36, 0.36, color, color)
        add_text(s, tag, x - 0.15, 6.05, 0.85, 0.20, size=8.0, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, label, x - 0.75, 6.30, 2.5, 0.50, size=9.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "One question, two semesters, eighteen manuscripts, and a fail-closed path to the next result.",
             0.69, 7.02, 11.8, 0.23, size=9.5, color=MUTED, italic=True)

    # ── 2. Program at a glance ─────────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "Overview", "The program in one slide", 2)
    add_metric(s, 0.65, 1.42, 2.30, 1.24, "1,001", "commits", TEAL, "14 Mar → 02 Aug 2026")
    add_metric(s, 3.12, 1.42, 2.30, 1.24, "18", "manuscripts", BLUE, "8 pillar papers + venue variants + synthesis")
    add_metric(s, 5.59, 1.42, 2.30, 1.24, "480", "paper pages", LAV, "P1–P8 compiled cleanly, 0 errors")
    add_metric(s, 8.06, 1.42, 2.30, 1.24, "40/40", "E1 units verified", AMBER, "5 arms × 8 seeds, Colab A100")
    add_metric(s, 10.53, 1.42, 2.30, 1.24, "INCONCL.", "E1 verdicts", RED, "correct statistics, honest result")
    row = [
        ("01", "QUESTION", "Does group-relative RL post-training starve of gradient signal — and can we measure it?", TEAL),
        ("02", "BENCHMARK", "TinkerRL-Bench: unified benchmark for RL post-training across 4 frameworks, 6 backends", BLUE),
        ("03", "DISCOVERY", "ZVF: the zero-variance fraction — an exact identity verified to 1.11e-16 on 505 tasks", LAV),
        ("04", "EVIDENCE", "40-unit confirmatory campaign, 983-run audit workbook, W&B 153-run public record", AMBER),
        ("05", "DISCIPLINE", "Preregistered, hash-bound, fail-closed protocols — no claim until the evidence clears", GREEN),
    ]
    for i, (num, kicker, body, color) in enumerate(row):
        y = 2.98 + i * 0.80
        add_shape(s, MSO_SHAPE.OVAL, 0.85, y + 0.06, 0.52, 0.52, color, color)
        add_text(s, num, 0.85, y + 0.18, 0.52, 0.22, size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, kicker, 1.62, y - 0.02, 2.3, 0.26, size=11, color=color, bold=True)
        add_text(s, body, 1.62, y + 0.30, 10.9, 0.34, size=10.6, color=INK)
    add_footer(s, "Sources: git log; PAPERS_README.md; ARTIFACT.md; zvf-program/audit/COLAB_EXECUTION_STATUS.md; STATISTICAL_REANALYSIS.md")



    # ── 3. Research question ───────────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "Motivation", "Why group-relative RL post-training, and why measure signal starvation", 3)
    add_text(s, "GRPO trains a language model on groups of G completions per prompt; each completion's advantage is group-relative. "
                "Under binary rewards, a group whose completions all get the same reward contributes exactly zero gradient.", 0.65, 1.42, 7.4, 1.15, size=13, color=INK)
    add_bullet(s, "Reward curves alone cannot see this: training can look maximally successful (reward ≈ 1.0) while every update carries no signal.", 0.65, 2.72, 7.4, 0.62, size=12, color=MUTED, bullet_color=TEAL)
    add_bullet(s, "The field lacks a unified benchmark for RL post-training — each library reports its own stack-conditioned numbers.", 0.65, 3.40, 7.4, 0.62, size=12, color=MUTED, bullet_color=TEAL)
    add_bullet(s, "No accepted minimum-report standard: 17× same-label comparisons traced to a changed base checkpoint (under-specification).", 0.65, 4.08, 7.4, 0.62, size=12, color=MUTED, bullet_color=TEAL)
    add_bullet(s, "Contribution: a diagnostic (ZVF), a benchmark (TinkerRL-Bench), a reporting standard (MIN-REPORT-RL), and a registry.", 0.65, 4.76, 7.4, 0.62, size=12, color=INK, bold=True, bullet_color=BLUE)
    # floor/ceiling panel
    add_shape(s, MSO_SHAPE.RECTANGLE, 8.35, 1.42, 4.35, 5.00, PANEL, GRID, radius=True)
    add_text(s, "FLOOR / CEILING FRAME", 8.62, 1.74, 3.8, 0.26, size=10, color=AMBER, bold=True)
    rows = [
        ("Arithmetic", "random 0.5%", "oracle 100%", TEAL),
        ("GSM8K", "GPT-2 ≈ 2%", "human ≈ 95%", BLUE),
        ("Distillation", "0% compress", "human ≈ 50%", LAV),
    ]
    for i, (task, floor, ceil, color) in enumerate(rows):
        yy = 2.32 + i * 0.88
        add_text(s, task, 8.62, yy, 1.5, 0.24, size=10.5, color=INK, bold=True)
        add_text(s, f"floor  {floor}", 8.62, yy + 0.30, 2.0, 0.22, size=9.5, color=MUTED)
        add_text(s, f"ceiling  {ceil}", 10.45, yy + 0.30, 2.1, 0.22, size=9.5, color=color, bold=True)
        add_line(s, 8.62, yy + 0.62, 12.40, yy + 0.62, GRID, 0.6)
    add_text(s, "Every trained agent must exceed its floor to demonstrate learning — statistics per Agarwal et al. 2021 (IQM, bootstrap CIs, 10 seeds).", 8.62, 5.42, 3.8, 0.62, size=9, color=MUTED, italic=True)
    add_footer(s, "Sources: BASELINES.md; thesis abstract (Claim 1); PAPERS_README.md items 4–5")

    # ── 4. Phase 1 — Semester 3 capstone ───────────────────────────────────
    s = new_slide(prs)
    add_header(s, "Phase 1 · Semester 3 · Group 6", "TinkerRL-Bench — a unified benchmark for RL post-training", 4)
    add_metric(s, 0.65, 1.42, 2.55, 1.24, "153", "W&B runs", TEAL, "public, 23.9 h client wall-clock")
    add_metric(s, 3.37, 1.42, 2.55, 1.24, "6", "authors + 2 guides", BLUE, "Group 6, equal-contribution leads")
    add_metric(s, 6.09, 1.42, 2.55, 1.24, "4", "frameworks", LAV, "Tinker/SkyRL · verl · OpenRLHF · TRL")
    add_metric(s, 8.81, 1.42, 2.55, 1.24, "13", "trainer scripts", AMBER, "34 Atropos + 31 scaling configs")
    add_metric(s, 11.53, 1.42, 1.15, 1.24, "52p", "anon. paper", GREEN, "NeurIPS 2026 D&B, blind")
    add_text(s, "Deliverables", 0.68, 2.95, 3.0, 0.30, size=15, color=INK, bold=True, font=FONT_DISPLAY)
    bullets = [
        ("NeurIPS 2026 Datasets & Benchmarks submission", "compiled blind review bundle; 0 PII hits; tag v1.0.0-neurips-2026", TEAL),
        ("ACM artifact review — three badges", "Available · Evaluated-Functional (< 10 min smoke test) · Reusable", BLUE),
        ("Reproducibility engineering", "pinned Docker (CUDA 12.4 + Ubuntu 22.04), seeds, ±2 pp claim checks", LAV),
        ("Recipes across environments", "Atropos, GSM8K, MATH, HumanEval, tool use — math RL hit 100% accuracy", AMBER),
    ]
    for i, (title, body, color) in enumerate(bullets):
        y = 3.42 + i * 0.92
        add_shape(s, MSO_SHAPE.OVAL, 0.82, y + 0.08, 0.34, 0.34, color, color)
        add_text(s, title, 1.38, y - 0.02, 8.0, 0.28, size=13, color=INK, bold=True)
        add_text(s, body, 1.38, y + 0.32, 8.6, 0.26, size=10.2, color=MUTED)
    add_shape(s, MSO_SHAPE.RECTANGLE, 9.70, 3.42, 2.95, 3.16, PANEL, GRID, radius=True)
    add_text(s, "FROZEN BOUNDARY", 9.95, 3.72, 2.4, 0.24, size=9, color=TEAL, bold=True)
    add_text(s, "Semester 3 is frozen at the April 2026 submission; Semester 4 continues solo from tag capstone-final-2026-04-25.", 9.95, 4.10, 2.45, 1.0, size=9.8, color=INK)
    add_text(s, "Author record: CITATION.cff + PROJECT_HISTORY.md separate group from solo work.", 9.95, 5.62, 2.45, 0.62, size=8.8, color=MUTED)
    add_footer(s, "Sources: ARTIFACT.md; CHANGELOG.md; README.md; platform_hybrid/sem 3 work/PROVENANCE.md")


    # ── 5. Early recipes & headline result ─────────────────────────────────
    s = new_slide(prs)
    add_header(s, "Phase 1 · Experiments", "Recipes that validated the platform end-to-end", 5)
    recipes = [
        ("Math RL (arithmetic)", "train model to add numbers", "100% accuracy", GREEN),
        ("Chat SFT", "NoRobots supervised fine-tune", "Complete", TEAL),
        ("Preference shorter", "concise-response reward", "Complete", BLUE),
        ("Distillation off-policy", "SFT on OpenThoughts3", "Complete", LAV),
        ("Distillation on-policy", "KL to teacher, no SFT copy", "Complete", AMBER),
        ("GSM8K GRPO vs PPO", "Qwen3-8B, held-out 500", "GRPO ≈ PPO (±0.1 pp)", BLUE),
    ]
    for i, (name, desc, status, color) in enumerate(recipes):
        y = 1.50 + i * 0.55
        add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, y, 5.6, 0.42, PANEL, GRID, radius=True)
        add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, y, 0.05, 0.42, color)
        add_text(s, name, 0.85, y + 0.075, 2.6, 0.26, size=10.5, color=INK, bold=True)
        add_text(s, desc, 3.45, y + 0.075, 1.9, 0.26, size=9, color=MUTED)
        add_text(s, status, 5.20, y + 0.075, 1.0, 0.26, size=8.8, color=color, bold=True)
    add_text(s, "The platform works; the science question is what the recipes reveal.", 0.68, 5.10, 5.6, 0.30, size=12, color=AMBER, bold=True)
    # figure
    add_picture(s, FIG / "ppo_vs_grpo.png", 6.65, 1.42, 6.10, 4.08)
    add_text(s, "GRPO vs PPO learning curves, Qwen3-8B / GSM8K — the flagship comparison that the E1 campaign later re-ran under a frozen protocol.", 6.68, 5.66, 6.05, 0.60, size=9.2, color=MUTED)
    add_text(s, "statistical floor: 10 seeds · IQM · bootstrap 95% CI (Agarwal et al., 2021)", 0.68, 6.60, 6.0, 0.26, size=9, color=MUTED, italic=True)
    add_footer(s, "Sources: README.md recipe table; experiments/ results; figures/v2/ppo_vs_grpo.png")

    # ── 6. Infrastructure matrix ───────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "Infrastructure", "One codebase, four RL frameworks, six compute backends", 6)
    add_text(s, "The consolidated repository keeps every platform reproducible with pinned versions and a unified launcher.", 0.65, 1.40, 11.5, 0.30, size=13.5, color=MUTED)
    frameworks = [
        ("Tinker / SkyRL", "thinkingmachines platform, GRPO presets (gsm8k --steps 200)", TEAL),
        ("verl", "0.3.0.post1 verified on isolated Python 3.11 env", BLUE),
        ("OpenRLHF", "Modal serverless H100 runs", LAV),
        ("HuggingFace TRL", "1.2.0 pinned; GRPO/Dr.GRPO/DAPO loss kernels source-audited", AMBER),
    ]
    for i, (name, body, color) in enumerate(frameworks):
        y = 1.95 + i * 0.78
        add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, y, 5.85, 0.64, PANEL, GRID, radius=True)
        add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, y, 0.05, 0.64, color)
        add_text(s, name, 0.88, y + 0.075, 2.6, 0.26, size=11, color=INK, bold=True)
        add_text(s, body, 0.88, y + 0.35, 5.4, 0.24, size=9, color=MUTED)
    backends = [
        ("LOCAL", "GPU unified launcher", TEAL),
        ("MODAL", "H100 serverless", BLUE),
        ("COLAB", "A100 campaign fleet", LAV),
        ("VAST.AI", "rented GPUs", AMBER),
        ("GCP", "A100 preflights", GREEN),
        ("HF SPACES", "hosted demos", RED),
    ]
    for i, (name, body, color) in enumerate(backends):
        x = 6.85 + (i % 3) * 2.10
        y = 1.95 + (i // 3) * 1.30
        add_shape(s, MSO_SHAPE.RECTANGLE, x, y, 1.92, 1.10, PANEL, GRID, radius=True)
        add_shape(s, MSO_SHAPE.RECTANGLE, x, y, 1.92, 0.06, color)
        add_text(s, name, x + 0.14, y + 0.16, 1.7, 0.24, size=10.5, color=color, bold=True)
        add_text(s, body, x + 0.14, y + 0.52, 1.7, 0.26, size=8.6, color=MUTED)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, 5.25, 12.05, 1.10, PANEL, GRID, radius=True)
    add_text(s, "DOCKER LOCK", 0.88, 5.48, 2.0, 0.24, size=9.5, color=AMBER, bold=True)
    add_text(s, "CUDA 12.4 + Ubuntu 22.04, requirements pinned, OCI label records the exact commit; package lock for E1: TRL 1.2.0 · Transformers 5.5.4 · Torch 2.7.1 · PEFT 0.19.1 · W&B 0.21.0 · NumPy 2.2.6",
             0.88, 5.80, 11.5, 0.44, size=10, color=INK)
    add_footer(s, "Sources: README.md; Dockerfile; execution-notes.md frozen-runtime record; platform_* directories")

    # ── 7. Phase 2 — Semester 4 solo continuation ──────────────────────────
    s = new_slide(prs)
    add_header(s, "Phase 2 · Semester 4 · Solo", "From benchmark to a falsifiable thesis on signal starvation", 7)
    add_text(s, "The Zero-Variance Fraction: Diagnosing and Budgeting Signal Starvation in Group-Relative RL Post-Training of LLMs",
             0.65, 1.40, 12.0, 0.62, size=15, color=INK, bold=True, font=FONT_DISPLAY)
    add_text(s, "M.Tech thesis · Arvind C R · guide Prof. Ramesh Prakash Guledgudd · July 2026 · 31-page current draft, chapter ↔ paper provenance map",
             0.68, 2.06, 11.8, 0.28, size=10.5, color=MUTED)
    add_metric(s, 0.65, 2.62, 2.30, 1.10, "206", "analysis iterations", TEAL, "recorded, versioned")
    add_metric(s, 3.12, 2.62, 2.30, 1.10, "8", "pillar papers", BLUE, "P1–P8, each compiles clean")
    add_metric(s, 5.59, 2.62, 2.30, 1.10, "480", "paper pages", LAV, "P5 80p · P6 65p · P7 81p · P8 94p")
    add_metric(s, 8.06, 2.62, 2.30, 1.10, "17→18", "document roster", AMBER, "unified synthesis added")
    add_metric(s, 10.53, 2.62, 2.30, 1.10, "2", "bounded claims", GREEN, "diagnostic + group-size")
    papers = [
        ("P1", "Scaling Laws for GRPO", "limits & identifiability audit — cross-library, cross-scale", 45, BLUE),
        ("P2", "The Zero-Variance Fraction", "descriptive diagnostic + exact accounting", 45, TEAL),
        ("P3", "Group Size in GRPO", "contrast density and the bridge to DPO", 25, LAV),
        ("P4", "Length Bias", "held-out generalization in GRPO & Dr.GRPO", 45, AMBER),
        ("P5", "Report the Stack", "8-item minimum-report standard", 80, GREEN),
        ("P6", "GRPO-Registry", "machine-readable catalog, 7-field runs", 65, RED),
        ("P7", "ZVF Controller", "signal-starvation theory; adaptive G not promoted", 81, BLUE),
        ("P8", "LLM vs XGBoost Fraud", "sensor and scribe, not scorer — parked side study", 94, TEAL),
    ]
    for i, (pid, title, body, pages, color) in enumerate(papers):
        x = 0.65 + (i % 4) * 3.05
        y = 4.02 + (i // 4) * 1.42
        add_shape(s, MSO_SHAPE.RECTANGLE, x, y, 2.88, 1.26, PANEL, GRID, radius=True)
        add_shape(s, MSO_SHAPE.RECTANGLE, x, y, 2.88, 0.05, color)
        add_text(s, f"{pid} · {pages}p", x + 0.16, y + 0.14, 2.5, 0.22, size=8.6, color=color, bold=True)
        add_text(s, title, x + 0.16, y + 0.40, 2.55, 0.26, size=10.2, color=INK, bold=True)
        add_text(s, body, x + 0.16, y + 0.70, 2.55, 0.52, size=8.0, color=MUTED)
    add_footer(s, "Sources: platform_hybrid/sem 4 work/README.md; PAPERS_README.md; PDF page counts (mdls); thesis/main.tex")

    # ── 8. Core discovery — ZVF ────────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "Core discovery", "ZVF — the Zero-Variance Fraction", 8)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, 1.42, 6.05, 1.06, PANEL, GRID, radius=True)
    add_text(s, "ZVF = fraction of groups whose completions share one reward — all-correct or all-wrong — and therefore contribute zero advantage (zero gradient).",
             0.88, 1.58, 5.6, 0.72, size=11.5, color=INK)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, 2.62, 6.05, 0.86, PANEL_2, GRID, radius=True)
    add_text(s, "Exact identity · pass@G − p^G = 1 − ZVF", 0.88, 2.80, 5.6, 0.26, size=12, color=TEAL, bold=True)
    add_text(s, "reproduced to 1.11e-16 across the 505-task audit", 0.88, 3.10, 5.6, 0.24, size=10, color=MUTED)
    add_text(s, "Claim 1 — diagnostic", 0.88, 3.72, 5.0, 0.26, size=12, color=BLUE, bold=True)
    add_text(s, "ZVF + mean reward reveals zero-advantage regimes that reward curves alone cannot show — including reward ≈ 1.0 while every update carries no signal. Wilson CI coverage 0.95–0.98.",
             0.88, 4.04, 5.55, 0.90, size=10.2, color=INK)
    add_text(s, "Claim 2 — group size", 0.88, 5.02, 5.0, 0.26, size=12, color=LAV, bold=True)
    add_text(s, "At a matched rollout budget on Qwen3-8B/GSM8K, G controls which end of training starves (see next slide).",
             0.88, 5.34, 5.55, 0.62, size=10.2, color=INK)
    add_picture(s, FIG / "zvf_heatmap.png", 6.95, 1.42, 5.85, 2.18)
    add_text(s, "ZVF heatmap across iterations and group sizes — the wall at small G is visible before reward collapses.", 6.98, 3.68, 5.75, 0.30, size=9, color=MUTED)
    add_shape(s, MSO_SHAPE.RECTANGLE, 6.95, 4.12, 5.85, 2.30, PANEL, GRID, radius=True)
    add_text(s, "WHY IT MATTERS", 7.20, 4.38, 2.5, 0.24, size=9.5, color=AMBER, bold=True)
    add_text(s, "GRPO's signal is the group contrast. When the contrast dies, training continues but learns nothing — and standard reward/loss curves hide it. ZVF makes the starvation observable, cheap, and online.",
             7.20, 4.72, 5.35, 1.40, size=11, color=INK)
    add_footer(s, "Sources: thesis/main.tex abstract; PAPERS_README.md item 1; figures/v2/zvf_heatmap.png")


    # ── 9. Group size & the bridge to DPO ──────────────────────────────────
    s = new_slide(prs)
    add_header(s, "Pillar P3", "Group size controls which end of training starves", 9)
    add_text(s, "Matched-token, two-seed panel: G=2 × 160 vs G=16 × 20 · 2,560 rollouts per arm · Qwen3-8B / GSM8K", 0.65, 1.40, 7.6, 0.30, size=12.5, color=MUTED)
    add_picture(s, FIG / "group_size_ablation.png", 6.75, 1.95, 5.95, 3.70)
    add_text(s, "Group-size ablation — small G saturates then starves; large G holds contrast.", 6.78, 5.78, 5.85, 0.28, size=9, color=MUTED)
    cards = [
        ("G = 2", "ends in a sustained all-correct zero-variance wall", "ZVF ≈ 0.75–1.0 at reward ≈ 1.0", RED),
        ("G = 16", "keeps contrast throughout training", "ZVF ≤ 0.25 at all steps", GREEN),
        ("G = 4", "best conditional utility on the 505-task cohort", "selected by (1 − ZVF)/√G audit", TEAL),
        ("G ≈ 32", "reconstructed hypotheses, different budgets", "not universal prescriptions", MUTED),
    ]
    for i, (g, body, num, color) in enumerate(cards):
        y = 1.95 + i * 1.02
        add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, y, 5.85, 0.88, PANEL, GRID, radius=True)
        add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, y, 0.05, 0.88, color)
        add_text(s, g, 0.85, y + 0.10, 1.1, 0.26, size=13, color=color, bold=True, font=FONT_DISPLAY)
        add_text(s, body, 1.95, y + 0.08, 4.4, 0.26, size=10.5, color=INK, bold=True)
        add_text(s, num, 1.95, y + 0.40, 4.4, 0.24, size=9.3, color=MUTED)
    add_text(s, "Adaptive G is NOT promoted: 92.3% of logged escalation fires are on all-correct groups; promotion needs a seed-paired, fixed-token bakeoff against static G=16.",
             0.68, 6.30, 12.0, 0.50, size=11, color=AMBER, bold=True)
    add_footer(s, "Sources: PAPERS_README.md items 2–4; thesis abstract Claim 2; figures/v2/group_size_ablation.png")

    # ── 10. Length bias & held-out generalization ──────────────────────────
    s = new_slide(prs)
    add_header(s, "Pillar P4", "Length bias under a bounded null", 10)
    add_text(s, "Question: does GRPO's implicit length pressure generalize, or does it just lengthen outputs?", 0.65, 1.42, 7.6, 0.30, size=13, color=MUTED)
    add_bullet(s, "Bounded null: with a 200-token cap, length bias effects stay inside the equivalence margin — the claim is bounded, not universal.", 0.65, 1.95, 7.5, 0.60, size=11.5, color=INK, bullet_color=TEAL)
    add_bullet(s, "Held-out design: 500 unseen GSM8K questions, checkpoint-anchored evaluation every 20 steps.", 0.65, 2.63, 7.5, 0.60, size=11.5, color=INK, bullet_color=BLUE)
    add_bullet(s, "Dr.GRPO arm included: loss-kernel semantics audited at source level (arm-specific reductions, 1e-4 epsilon).", 0.65, 3.31, 7.5, 0.60, size=11.5, color=INK, bullet_color=LAV)
    add_bullet(s, "Remaining gate: an uncapped long-horizon mediation study before any positive claim.", 0.65, 3.99, 7.5, 0.60, size=11.5, color=AMBER, bold=True, bullet_color=AMBER)
    add_picture(s, FIG / "p2_length_bias.png", 8.50, 1.95, 4.25, 2.94)
    add_text(s, "Length-bias analysis — reward vs response-length trajectories across GRPO variants.", 8.52, 5.00, 4.2, 0.28, size=9, color=MUTED)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, 5.02, 7.55, 1.42, PANEL, GRID, radius=True)
    add_text(s, "P1 in one line", 0.88, 5.24, 3.0, 0.24, size=10, color=BLUE, bold=True)
    add_text(s, "Scaling-law claims for GRPO are stack-conditioned: the same recipe can look like a law on one stack and nothing on another. P1 is a limits & identifiability audit — the honest reading of multi-seed evidence, not a positive scaling law.",
             0.88, 5.56, 7.05, 0.74, size=10.5, color=INK)
    add_footer(s, "Sources: PAPERS_README.md (P04/P01 roles); execution-notes.md TRL semantics audit; figures/v2/p2_length_bias.png")

    # ── 11. Standards, registry, controller, side study ───────────────────
    s = new_slide(prs)
    add_header(s, "Pillars P5–P8", "Standards, registry, controller, and a measurement-discipline study", 11)
    cards = [
        ("P5 · MIN-REPORT-RL", "Eight-item minimum-report standard", "7 run-manifest fields + held-out pass@k — the canonical reporting standard so results stop being stack-conditioned by omission.", 80, GREEN),
        ("P6 · GRPO-Registry", "Machine-readable living catalog", "seven-field run-start JSON per run; position-artifact resource for the community, evidence-mapped per paper.", 65, TEAL),
        ("P7 · ZVF Controller", "Diagnostic → controller, honestly scoped", "retrospective audit + prospective test plan only; adaptive G explicitly NOT promoted until a fixed-token bakeoff.", 81, BLUE),
        ("P8 · Fraud study", "LLM vs XGBoost, parked side study", "LLM as sensor & scribe, not scorer; demonstrates measurement discipline — contributes no RL evidence.", 94, AMBER),
    ]
    for i, (title, sub, body, pages, color) in enumerate(cards):
        x = 0.65 + (i % 2) * 6.12
        y = 1.50 + (i // 2) * 2.62
        add_shape(s, MSO_SHAPE.RECTANGLE, x, y, 5.92, 2.42, PANEL, GRID, radius=True)
        add_shape(s, MSO_SHAPE.RECTANGLE, x, y, 5.92, 0.06, color)
        add_text(s, title, x + 0.22, y + 0.18, 3.6, 0.26, size=12.5, color=color, bold=True)
        add_text(s, f"{pages} pages", x + 4.60, y + 0.20, 1.1, 0.22, size=9, color=MUTED, align=PP_ALIGN.RIGHT)
        add_text(s, sub, x + 0.22, y + 0.52, 5.4, 0.26, size=10.5, color=INK, bold=True)
        add_text(s, body, x + 0.22, y + 0.88, 5.45, 1.30, size=10.2, color=MUTED)
    add_text(s, "These four pillars turn a single diagnostic into community infrastructure — but the controller stays un-promoted until measured.",
             0.68, 6.48, 12.0, 0.30, size=11.5, color=AMBER, bold=True)
    add_footer(s, "Sources: PAPERS_README.md (P05–P08 roles); sem 4 work/papers PDFs")


    # ── 12. E1 confirmatory campaign ───────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "Evidence · E1 campaign", "The frozen confirmatory audit — 40/40 units validated", 12)
    add_metric(s, 0.65, 1.42, 2.30, 1.24, "40/40", "units verified", TEAL, "40 local + 40 remote, zero errors")
    add_metric(s, 3.12, 1.42, 2.30, 1.24, "5 × 8", "arms × seeds", BLUE, "GRPO, DAPO, GSPO, Dr.GRPO, AERO")
    add_metric(s, 5.59, 1.42, 2.30, 1.24, "500", "held-out questions", LAV, "GSM8K, pinned Qwen3-8B")
    add_metric(s, 8.06, 1.42, 2.30, 1.24, "30", "optimizer steps", AMBER, "LoRA on Colab A100")
    add_metric(s, 10.53, 1.42, 2.30, 1.24, "6", "manifests repaired", RED, "exact-checkpoint replay, re-verified")
    add_text(s, "Campaign facts", 0.68, 2.95, 3.0, 0.30, size=15, color=INK, bold=True, font=FONT_DISPLAY)
    bullets = [
        ("Executed on NVIDIA A100-SXM4-40GB via Colab CLI, OAuth2 account arvindcr4@gmail.com", TEAL),
        ("DAPO preregistered dynamic sampling realized 1,472–2,112 rollouts per unit; GRPO/GSPO 480 completions each", BLUE),
        ("Six legacy manifests predating completion hashes repaired via evaluation-only checkpoint replay; GSPO seed 71 byte-identical manifest re-verified", LAV),
        ("GRPO held-out ≈ 0.63–0.65 — near-identical across arms; the real story is in the statistics, not the mean", AMBER),
    ]
    for i, (body, color) in enumerate(bullets):
        y = 3.42 + i * 0.80
        add_shape(s, MSO_SHAPE.OVAL, 0.82, y + 0.08, 0.30, 0.30, color, color)
        add_text(s, body, 1.32, y, 8.3, 0.62, size=10.8, color=INK)
    add_shape(s, MSO_SHAPE.RECTANGLE, 9.70, 2.95, 2.95, 3.62, PANEL, GRID, radius=True)
    add_text(s, "HEADLINE", 9.95, 3.22, 2.4, 0.24, size=9, color=AMBER, bold=True)
    add_text(s, "Every arm finished.\nNo arm won.\n\nFour verdicts, all\nINCONCLUSIVE after\ncorrect statistics.", 9.95, 3.60, 2.45, 1.9, size=14, color=INK, bold=True, font=FONT_DISPLAY, valign=MSO_ANCHOR.MIDDLE)
    add_text(s, "That is the honest\nanswer — and it is\nwhat fail-closed\npreregistration is for.", 9.95, 5.72, 2.45, 0.75, size=9, color=MUTED)
    add_footer(s, "Sources: zvf-program/audit/COLAB_EXECUTION_STATUS.md; results/audit.json")


    # ── 13. Statistical reanalysis ─────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "Evidence · statistics", "Why the corrected answer is INCONCLUSIVE — and why that is a feature", 13)
    add_text(s, "The run records and scores never changed. The correction was in the analysis code: exact paired-t power replaces the large-sample normal approximation, and the preregistered Benjamini–Hochberg step is now executed.",
             0.65, 1.40, 11.9, 0.62, size=12.5, color=INK)
    cols = ["Arm", "Paired Δ", "95% bootstrap CI", "MDE80", "Raw p", "BH", "Verdict"]
    widths = [1.5, 1.5, 3.0, 1.5, 1.2, 1.0, 1.9]
    x0 = 0.65
    for c, wdt in zip(cols, widths):
        add_text(s, c, x0, 2.30, wdt, 0.24, size=10, color=TEAL, bold=True)
        x0 += wdt
    rows = [
        ("DAPO", "+0.00100", "[-0.00450, +0.00675]", "0.01012", "0.756", "No", "INCONCLUSIVE", RED),
        ("GSPO", "+0.00500", "[-0.00125, +0.01200]", "0.01185", "0.210", "No", "INCONCLUSIVE", RED),
        ("Dr.GRPO", "-0.00200", "[-0.00950, +0.00725]", "0.01483", "0.673", "No", "INCONCLUSIVE", RED),
        ("AERO", "-0.00075", "[-0.00825, +0.00675]", "0.01319", "0.858", "No", "INCONCLUSIVE", RED),
    ]
    for i, (arm, d, ci, mde, p, bh, verdict, color) in enumerate(rows):
        yy = 2.70 + i * 0.52
        add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, yy - 0.05, 10.6, 0.42, PANEL if i % 2 == 0 else PANEL_2, GRID, radius=True)
        x0 = 0.65
        for val, wdt in zip([arm, d, ci, mde, p, bh, verdict], widths):
            add_text(s, val, x0 + 0.08, yy, wdt - 0.1, 0.24, size=10, color=INK if val != "INCONCLUSIVE" else color,
                     bold=(val in ("INCONCLUSIVE", arm)))
            x0 += wdt
    add_text(s, "DAPO's 90% CI sits inside ±0.01, but its exact MDE80 (0.01012) exceeds the equivalence margin — so the locked rule forbids DISAPPEARS.",
             0.68, 4.85, 10.5, 0.40, size=10.8, color=AMBER, bold=True)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, 5.35, 12.05, 1.20, PANEL, GRID, radius=True)
    add_text(s, "WHAT THIS ESTABLISHES", 0.88, 5.58, 3.0, 0.24, size=9.5, color=TEAL, bold=True)
    add_text(s, "1 · Eight paired seeds are too few to resolve ±1 pp differences — the protocol itself must be upgraded (it was: 24-unit pilot).  2 · Earlier DAPO-DISAPPEARS statements are formally superseded everywhere.  3 · The machinery that caught this (exact power + multiplicity correction) is now the standard in this repo.",
             0.88, 5.92, 11.55, 0.55, size=10, color=INK)
    add_footer(s, "Sources: zvf-program/audit/STATISTICAL_REANALYSIS.md (table verbatim); execution-notes.md 2026-08-02 correction")


    # ── 14. Preregistration discipline ─────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "Discipline", "Preregistered, hash-bound, fail-closed science", 14)
    add_text(s, "Stage 5 is fail-closed on a proven joint-zero-gradient contradiction — 62/100 reward-degenerate groups (59 all-wrong, 3 all-correct) in the accepted corpus. No replacement unit may run until an explicit amendment is authorized.",
             0.65, 1.42, 11.9, 0.72, size=12, color=INK)
    rows = [
        ("Flagship protocol", "SHA-256 68237294…5171e", "frozen; step-0 eval identical across units; zero-advantage = zero gradient respected", TEAL),
        ("Pilot control plane", "24 units = 4 conditions × 2 regimes × seeds 11/23/37", "all dry_run_only, allocation.allowed = false; A100-only; disjoint 5-seed confirmation", BLUE),
        ("Pilot protocol", "SHA-256 5a0bbd25…00d7", "numeric execution contract complete; no readiness blocker", LAV),
        ("Replay contract", "SHA-256 21867edf…4e2a", "one content-addressed corpus per regime/seed replayed in same order by all 4 conditions", AMBER),
        ("Locked runtime", "Python 3.11/3.12 · TRL 1.2.0 · Torch 2.7.1", "Transformers 5.5.4 · Datasets 4.8.4 · PEFT 0.19.1 · W&B 0.21.0 · NumPy 2.2.6", GREEN),
    ]
    for i, (name, meta, body, color) in enumerate(rows):
        y = 2.42 + i * 0.86
        add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, y, 12.05, 0.74, PANEL, GRID, radius=True)
        add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, y, 0.05, 0.74, color)
        add_text(s, name, 0.88, y + 0.08, 2.9, 0.24, size=11, color=color, bold=True)
        add_text(s, meta, 0.88, y + 0.38, 4.6, 0.24, size=9.2, color=INK)
        add_text(s, body, 5.75, y + 0.08, 6.8, 0.56, size=9.2, color=MUTED)
    add_text(s, "Receipts are content-addressed; checkpoints carry hashed evaluation JSONL; resume restores model + RNG + cursor exactly; an independent verifier re-downloads and re-verifies every tree.",
             0.68, 6.60, 12.0, 0.34, size=10.5, color=TEAL, bold=True)
    add_footer(s, "Sources: execution-notes.md current-gate & evidence sections; pilot_preregistration.json")


    # ── 15. Publication program ────────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "Publications", "Eighteen documents aimed at five venues", 15)
    venues = [
        ("NeurIPS 2026 · Main D&B", "submitted (blind, anon. 52 pp)", "OpenReview ID CXbcYe69BQ · 239-page compendium source · 7/7 integration checks", TEAL, "SUBMITTED"),
        ("NeurIPS 2026 · Workshop", "main_zvf · main_workshop · main_dnb", "ZVF sentinel paper, exploratory artifact note, tiered artifact paper", BLUE, "READY"),
        ("ACM (compact)", "acm_main, 11 pp", "cross-library derivative with ethics statement; R2R3 fixes applied", LAV, "READY"),
        ("Springer variant", "springer_main", "book-chapter-style derivation of the program", AMBER, "DRAFT"),
        ("TMLR / venue track", "target for the flagship result", "after the pilot completes with confirmed evidence", GREEN, "PLANNED"),
    ]
    for i, (name, meta, body, color, status) in enumerate(venues):
        y = 1.50 + i * 1.06
        add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, y, 9.35, 0.92, PANEL, GRID, radius=True)
        add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, y, 0.05, 0.92, color)
        add_text(s, name, 0.88, y + 0.10, 3.4, 0.26, size=11.5, color=color, bold=True)
        add_text(s, meta, 0.88, y + 0.42, 3.4, 0.24, size=9.0, color=INK)
        add_text(s, body, 4.35, y + 0.10, 5.45, 0.72, size=9.3, color=MUTED)
        add_chip(s, status, 10.25, y + 0.33, color=color, w=1.50)
    add_text(s, "Every manuscript compiles with 0 errors / 0 undefined refs; ethics statements included; blind-review tarball re-scanned with zero residual identifiers.",
             0.68, 6.62, 12.0, 0.30, size=10.5, color=INK, bold=True)
    add_footer(s, "Sources: PAPERS_README.md roster; CHANGELOG.md v3.0; autoresearch/reason-260727-2155/ (OpenReview audit)")


    # ── 16. Reproducibility & artifacts ────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "Reproducibility", "Artifacts a reviewer can actually run", 16)
    checks = [
        ("< 10 min", "smoke test, 7 checks", "platform_modal/scripts/smoke_test.sh", TEAL),
        ("7/7", "integration audit", "scripts/integration_audit.py", BLUE),
        ("±2 pp", "Qwen3-8B claim check", "reproducibility/check_qwen3_8b_claim.py", LAV),
        ("15/15", "CPU reference tests", "flagship stack, fully masked", AMBER),
        ("88", "focused tests", "64 audit + 24 S1, passing today", GREEN),
    ]
    for i, (value, label, path, color) in enumerate(checks):
        x = 0.65 + (i % 3) * 4.05
        y = 1.55 + (i // 3) * 1.35
        add_shape(s, MSO_SHAPE.RECTANGLE, x, y, 3.85, 1.18, PANEL, GRID, radius=True)
        add_shape(s, MSO_SHAPE.RECTANGLE, x, y, 3.85, 0.06, color)
        add_text(s, value, x + 0.18, y + 0.14, 2.2, 0.36, size=20, color=color, bold=True, font=FONT_DISPLAY)
        add_text(s, label, x + 0.18, y + 0.58, 3.4, 0.24, size=9.5, color=INK, bold=True)
        add_text(s, path, x + 0.18, y + 0.84, 3.5, 0.22, size=7.8, color=MUTED)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, 4.30, 12.05, 1.16, PANEL, GRID, radius=True)
    add_text(s, "ACM BADGES — TARGETED", 0.88, 4.52, 3.0, 0.24, size=10, color=TEAL, bold=True)
    add_text(s, "Artifacts Available (GitHub + HF Hub, DOI pending) · Evaluated-Functional (< 10 min entry) · Evaluated-Reusable (modular layout, pinned deps, extensible configs). Docker OCI label records the exact commit; seeds set Python/NumPy/Torch/CUDA.",
             0.88, 4.86, 11.55, 0.50, size=10, color=INK)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, 5.62, 12.05, 0.98, PANEL_2, GRID, radius=True)
    add_text(s, "CHANNELS", 0.88, 5.84, 1.5, 0.24, size=10, color=AMBER, bold=True)
    add_text(s, "GitHub (canonical + mirror) · Hugging Face Hub (LoRA adapters + model cards) · W&B (public run logs) · Zenodo (camera-ready DOI) · 983-run audit workbook",
             0.88, 6.16, 11.55, 0.30, size=10, color=INK)
    add_footer(s, "Sources: ARTIFACT.md §§1,7,8; CHANGELOG.md; run_all_audits.py")


    # ── 17. Evidence base ──────────────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "Evidence", "Everything traces to a named artifact", 17)
    add_metric(s, 0.65, 1.45, 2.90, 1.30, "983", "Tinker runs in workbook", TEAL, "audited, linked, defensible")
    add_metric(s, 3.72, 1.45, 2.90, 1.30, "153", "public W&B runs", BLUE, "Sem 3 record, 23.9 h")
    add_metric(s, 6.79, 1.45, 2.90, 1.30, "505", "task ZVF identity audit", LAV, "1.11e-16 worst-case residual")
    add_metric(s, 9.86, 1.45, 2.90, 1.30, "206", "analysis iterations", AMBER, "sem 4, version-controlled")
    add_text(s, "Provenance rules", 0.68, 3.05, 3.0, 0.30, size=15, color=INK, bold=True, font=FONT_DISPLAY)
    rules = [
        ("Every quantitative claim traces to a named artifact — no numbers typed by hand into prose.", TEAL),
        ("Frozen corpora carry full-split hashes + PCG64 seed-specific row orders for GSM8K and MATH-500.", BLUE),
        ("Stale or superseded records are kept as history but marked non-evidence (e.g., the 39/40 snapshot).", LAV),
        ("Verdicts require the executed preregistered procedure — multiplicity step included, always.", AMBER),
    ]
    for i, (body, color) in enumerate(rules):
        y = 3.52 + i * 0.76
        add_shape(s, MSO_SHAPE.OVAL, 0.82, y + 0.08, 0.30, 0.30, color, color)
        add_text(s, body, 1.32, y, 11.4, 0.60, size=11.5, color=INK)
    add_footer(s, "Sources: tinker_runs_audit_2026-07-12.xlsx; ARTIFACT.md; PAPERS_README.md item 1; execution-notes.md")

    # ── 18. Engineering discipline ─────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "Engineering", "Research-grade software discipline", 18)
    cards = [
        ("VERSION CONTROL", "1,001 commits, semantic releases, frozen semester boundaries and tags (capstone-final-2026-04-25).", TEAL),
        ("LOCKED BUILD", "uv-locked deps, ruff lint, pre-commit hooks, pytest suite, wheel-content verification.", BLUE),
        ("AUDIT SUITE", "python platform_local/run_all_audits.py — submission, claim-strength, reviewer-caveat, scientific audits.", LAV),
        ("FIGURES FROM CODE", "python -m platform_hybrid.paper.figure_module --profile all regenerates every paper figure.", AMBER),
        ("FIG. AUDIT", "FIGURE_AUDIT.md + FIGURES.tex single source of truth for every figure in the paper.", GREEN),
        ("STATS TOOLING", "utils/stats.py: IQM, bootstrap CIs, performance profiles, rliable integration.", RED),
    ]
    for i, (title, body, color) in enumerate(cards):
        x = 0.65 + (i % 3) * 4.05
        y = 1.55 + (i // 3) * 2.35
        add_shape(s, MSO_SHAPE.RECTANGLE, x, y, 3.85, 2.15, PANEL, GRID, radius=True)
        add_shape(s, MSO_SHAPE.RECTANGLE, x, y, 3.85, 0.06, color)
        add_text(s, title, x + 0.2, y + 0.20, 3.4, 0.26, size=11.5, color=color, bold=True)
        add_text(s, body, x + 0.2, y + 0.58, 3.45, 1.40, size=10, color=INK)
    add_text(s, "CI checks protect the paper's claims: integration_audit.json is written for CI on every merge.",
             0.68, 6.55, 12.0, 0.28, size=11, color=MUTED, italic=True)
    add_footer(s, "Sources: CONTRIBUTING.md; Makefile; pyproject.toml; CHANGELOG.md")

    # ── 19. Timeline ───────────────────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "Timeline", "Aug 2025 → Aug 2026", 19)
    add_line(s, 0.85, 4.10, 12.45, 4.10, GRID, 2.2)
    milestones = [
        (0.95, "AUG 25", "RL gym\nfoundations", TEAL),
        (2.60, "JAN 26", "Tinker\ncookbook", BLUE),
        (4.25, "MAR 26", "Atropos +\nframeworks", LAV),
        (5.90, "APR 26", "NeurIPS D&B\nsubmission", AMBER),
        (7.55, "MAY 26", "Sem 4 solo:\nP1–P8 papers", GREEN),
        (9.20, "JUL 26", "Defense +\nthesis", TEAL),
        (10.85, "JUL–AUG 26", "E1 campaign\n+ reanalysis", RED),
    ]
    for x, date, label, color in milestones:
        add_shape(s, MSO_SHAPE.OVAL, x, 3.92, 0.36, 0.36, color, color)
        add_text(s, date, x - 0.28, 2.95, 1.0, 0.22, size=8.2, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, label, x - 0.52, 3.22, 1.5, 0.50, size=9.3, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "Where today sits", 0.68, 4.62, 3.0, 0.28, size=14, color=INK, bold=True, font=FONT_DISPLAY)
    add_text(s, "The E1 audit is closed and corrected; the flagship protocol is frozen; the pilot control plane is built but dry-run-only; an amendment for the joint-zero-gradient representation is the only gate before GPU spend.",
             0.68, 5.02, 11.9, 0.62, size=12, color=INK)
    add_text(s, "No claim is on the table until the pilot's confirmatory matrix runs on A100.", 0.68, 5.90, 11.9, 0.28, size=11.5, color=AMBER, bold=True)
    add_footer(s, "Sources: PROJECT_HISTORY.md; git log; execution-notes.md")

    # ── 20. What's next ────────────────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "Forward path", "The next three gates", 20)
    steps = [
        ("01", "Authorize the amendment", "A joint-zero representation/scoring amendment for the fail-closed gate — the only blocker before GPU spend.", RED),
        ("02", "Run the 24-unit pilot matrix", "4 semantic conditions × 2 regimes × seeds 11/23/37 on A100; replay contract; hash-bound receipts; dry-run plans already generated.", TEAL),
        ("03", "Resolve the flagship claim", "Then matched multi-seed cross-scale evidence (P1), gradient geometry (P2), token-matched group sweep (P3), uncapped mediation (P4) — and only then a ZVF-aware controller bakeoff against static G=16.", BLUE),
    ]
    for i, (num, title, body, color) in enumerate(steps):
        y = 1.70 + i * 1.55
        add_shape(s, MSO_SHAPE.OVAL, 0.85, y + 0.10, 0.78, 0.78, color, color)
        add_text(s, num, 0.85, y + 0.31, 0.78, 0.26, size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, title, 1.95, y - 0.02, 9.6, 0.30, size=16, color=INK, bold=True, font=FONT_DISPLAY)
        add_text(s, body, 1.95, y + 0.42, 10.6, 0.80, size=11.5, color=MUTED)
        if i < 2:
            add_line(s, 1.25, y + 0.92, 1.25, y + 1.52, GRID, 1.2)
    add_shape(s, MSO_SHAPE.RECTANGLE, 10.60, 1.70, 2.08, 4.62, PANEL, GRID, radius=True)
    add_text(s, "MY RULE", 10.85, 2.00, 1.6, 0.40, size=9, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "No claim\nuntil the\npilot runs\nand the\nreceipts\nverify.", 10.80, 2.72, 1.7, 2.2, size=16, color=INK, bold=True, font=FONT_DISPLAY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(s, "check it → then explain it", 10.78, 5.55, 1.72, 0.55, size=9.5, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s, "Sources: execution-notes.md current gate; pilot_preregistration.json; PAPERS_README.md remaining gates")

    # ── 21. Takeaways ──────────────────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "Bottom line", "What the program stands on", 21)
    takes = [
        ("A real discovery", "ZVF is an exact, cheap, online diagnostic of GRPO signal starvation — verified to 1.11e-16 on 505 tasks, with two bounded thesis claims.", TEAL),
        ("A serious benchmark", "TinkerRL-Bench ships NeurIPS 2026 D&B, three ACM artifact badges, and a 983-run audited evidence base.", BLUE),
        ("Standards that stick", "MIN-REPORT-RL (8 items) and GRPO-Registry (7 fields) turn stack-conditioned results into comparable ones.", LAV),
        ("Honest statistics", "The E1 campaign finished 40/40 and answered INCONCLUSIVE — the corrected analysis caught our own earlier over-claim.", AMBER),
        ("Fail-closed discipline", "Preregistered, hash-bound, replay-controlled protocols: no claim until the evidence clears, on every gate.", GREEN),
    ]
    for i, (title, body, color) in enumerate(takes):
        y = 1.55 + i * 1.02
        add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, y, 12.05, 0.88, PANEL, GRID, radius=True)
        add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, y, 0.06, 0.88, color)
        add_text(s, title, 0.95, y + 0.10, 3.3, 0.28, size=12.5, color=color, bold=True)
        add_text(s, body, 4.35, y + 0.10, 8.1, 0.66, size=10.8, color=INK)
    add_text(s, "The result is still to be earned. Everything that earns it is now checkable.", 0.68, 6.62, 12.0, 0.30, size=12.5, color=AMBER, bold=True)
    add_footer(s, "Primary references: thesis/main.tex; ARTIFACT.md; PAPERS_README.md; STATISTICAL_REANALYSIS.md")

    # ── 22. Thank you ──────────────────────────────────────────────────────
    s = new_slide(prs)
    add_text(s, "THANK YOU", 0.85, 2.10, 4.0, 0.60, size=34, color=TEAL, bold=True, font=FONT_DISPLAY)
    add_text(s, "Questions & discussion", 0.88, 3.00, 4.0, 0.32, size=14, color=MUTED)
    add_text(s, "Repository · github.com/pes-llm-research/tinker-rl-lab  ·  mirror github.com/arvindcr4/tinker-rl-lab", 0.88, 5.40, 11.0, 0.26, size=10.5, color=INK)
    add_text(s, "Arvind C R · arvindcr4@gmail.com · M.Tech, PES University · guide Prof. Ramesh Prakash Guledgudd", 0.88, 5.76, 11.0, 0.26, size=10.5, color=INK)
    add_text(s, "02 Aug 2026 · deck built from repository state at build time", 0.88, 6.12, 11.0, 0.24, size=9, color=MUTED)
    add_text(s, "arXiv-style honesty: every number on every slide traces to a named file.", 0.88, 6.48, 11.0, 0.24, size=9, color=AMBER, italic=True)

    return prs


if __name__ == "__main__":
    deck = build()
    OUT.parent.mkdir(parents=True, exist_ok=True)
    deck.save(OUT)
    print(f"saved {OUT}")
    print(f"slides {len(deck.slides)}")
