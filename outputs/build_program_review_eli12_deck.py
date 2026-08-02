#!/usr/bin/env python3
"""Build the ELI12 (plain-words) faculty deck for Tinker RL Lab.

Same visual language and evidence as the full program-review deck, but every
slide is one big idea in plain words with everyday analogies — matching the
established ELI12 defense-transcript style (outputs/defense_transcript_eli12_2026-07-12.md).

Reuses the helpers from build_program_review_deck.py.
"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from build_program_review_deck import (
    FIG,
    NAVY,
    PANEL,
    PANEL_2,
    INK,
    MUTED,
    TEAL,
    BLUE,
    LAV,
    AMBER,
    RED,
    GREEN,
    GRID,
    FONT,
    FONT_DISPLAY,
    add_bullet,
    add_footer,
    add_header,
    add_line,
    add_metric,
    add_picture,
    add_shape,
    add_text,
    new_slide,
)

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs" / "Tinker_RL_Program_Review_ELI12_2026-08-02.pptx"

W, H = 13.333, 7.5


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    # ── 1. Title ────────────────────────────────────────────────────────────
    s = new_slide(prs)
    add_text(s, "TINKER RL LAB  /  PES UNIVERSITY  /  FACULTY REVIEW", 0.65, 0.55, 7.0, 0.25, size=10, color=TEAL, bold=True)
    add_text(s, "Teaching models to think —\nand catching when the teaching secretly stops", 0.65, 1.20, 12.2, 1.9, size=38, color=INK, bold=True, font=FONT_DISPLAY)
    add_text(s, "The whole story in plain words: what I built, what I found, what I got wrong, and how I fixed it.", 0.69, 3.55, 11.5, 0.40, size=16, color=MUTED)
    add_text(s, "ARVIND C R  ·  M.TECH  ·  GUIDE: PROF. RAMESH PRAKASH GULEDGUDD  ·  02 AUG 2026", 0.69, 4.15, 9.0, 0.25, size=10, color=AMBER, bold=True)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, 4.80, 12.05, 1.30, PANEL, GRID, radius=True)
    add_text(s, "THE ONE-PICTURE IDEA", 0.88, 5.02, 3.0, 0.24, size=9.5, color=TEAL, bold=True)
    add_text(s, "Imagine a basketball coach who only compares your shots inside one drill. If all eight shots go in — or all eight miss — there is nothing to compare, and the drill taught you zero. That is what happens inside GRPO training, and the scoreboard can't show it.",
             0.88, 5.34, 11.55, 0.62, size=11.5, color=INK)
    add_text(s, "I built a number that counts those empty drills — the Zero-Variance Fraction (ZVF) — and two years of evidence around it.", 0.88, 6.10, 11.55, 0.28, size=10.5, color=MUTED, italic=True)
    add_line(s, 0.72, 6.55, 12.45, 6.55, GRID, 1.6)
    phases = [
        (1.35, "SEM 3", "built the test ground", TEAL),
        (4.70, "SEM 4", "found the discovery", BLUE),
        (8.05, "NOW", "made the claims honest", AMBER),
    ]
    for x, tag, label, color in phases:
        add_shape(s, MSO_SHAPE.OVAL, x + 0.20, 6.37, 0.30, 0.30, color, color)
        add_text(s, tag, x - 0.10, 6.82, 0.90, 0.20, size=8.0, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, label, x - 0.85, 7.06, 2.4, 0.24, size=9.0, color=INK, bold=True, align=PP_ALIGN.CENTER)

    # ── 2. The one problem ──────────────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "The problem", "The scoreboard says everything is fine. The coach is asleep.", 2)
    add_text(s, "GRPO teaches a language model by making it write several answers to the same question, then comparing the answers with each other. When all the answers agree — all right or all wrong — there is nothing to compare, so the model learns nothing.", 0.65, 1.42, 7.3, 1.30, size=13.5, color=INK)
    add_shape(s, MSO_SHAPE.RECTANGLE, 8.30, 1.42, 4.40, 1.30, PANEL, GRID, radius=True)
    add_text(s, "REWARD = 1.0  (perfect!)", 8.55, 1.66, 3.9, 0.34, size=15, color=GREEN, bold=True, font=FONT_DISPLAY)
    add_text(s, "…while every update carries zero learning signal", 8.55, 2.16, 3.9, 0.34, size=10.5, color=RED, bold=True)
    add_text(s, "The reward curve — the one number everyone watches — literally cannot show this.", 8.55, 2.62, 3.9, 0.30, size=9.5, color=MUTED, italic=True)
    add_bullet(s, "As the model gets better, more and more groups are all-correct — so the empty drills get MORE common exactly when things look best.", 0.68, 3.05, 7.2, 0.75, size=12.5, color=INK, bullet_color=TEAL)
    add_bullet(s, "Same recipe, different software → wildly different results (I measured a 17× span from one undisclosed change).", 0.68, 3.90, 7.2, 0.75, size=12.5, color=INK, bullet_color=BLUE)
    add_bullet(s, "So we need a second number beside the reward — one that counts empty drills — and we need to report exactly what ran.", 0.68, 4.75, 7.2, 0.75, size=12.5, color=INK, bold=True, bullet_color=AMBER)
    add_shape(s, MSO_SHAPE.RECTANGLE, 8.30, 3.05, 4.40, 2.85, PANEL_2, GRID, radius=True)
    add_text(s, "WHY IT MATTERS", 8.55, 3.32, 2.5, 0.24, size=9.5, color=AMBER, bold=True)
    add_text(s, "Every lab that trains these models watches the reward. If the reward lies at the exact moment training dies, entire GPU budgets can be spent polishing nothing. Measuring the blindness is the first step to fixing it.", 8.55, 3.70, 3.9, 1.30, size=10.5, color=INK)
    add_text(s, "My answer: ZVF.", 8.55, 5.30, 2.5, 0.30, size=14, color=TEAL, bold=True, font=FONT_DISPLAY)
    add_footer(s, "Sources: thesis abstract; BASELINES.md; PAPERS_README.md")


    # ── 3. GRPO in plain words ──────────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "The method", "GRPO: eight answers, one question, compare — no referee", 3)
    steps = [
        ("01", "One question", "The model writes eight answers to the same math question (a \"group\").", TEAL),
        ("02", "Grade them", "Each answer is marked right or wrong by a simple check — no expensive second AI.", BLUE),
        ("03", "Compare", "Better than the group average → push up. Worse → push down. That push is the learning.", LAV),
        ("04", "Repeat", "Next question, next group, again and again.", AMBER),
    ]
    for i, (num, title, body, color) in enumerate(steps):
        y = 1.55 + i * 1.18
        add_shape(s, MSO_SHAPE.OVAL, 0.85, y + 0.08, 0.70, 0.70, color, color)
        add_text(s, num, 0.85, y + 0.27, 0.70, 0.24, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, title, 1.85, y - 0.02, 4.0, 0.30, size=15, color=color, bold=True, font=FONT_DISPLAY)
        add_text(s, body, 1.85, y + 0.38, 8.2, 0.55, size=11.5, color=INK)
        if i < 3:
            add_line(s, 1.20, y + 0.82, 1.20, y + 1.14, GRID, 1.1)
    add_shape(s, MSO_SHAPE.RECTANGLE, 8.95, 1.55, 3.75, 4.40, PANEL, GRID, radius=True)
    add_text(s, "THE CATCH", 9.20, 1.82, 3.2, 0.26, size=10, color=RED, bold=True)
    add_text(s, "All eight right — or all eight wrong — means nothing to compare. The group contributes exactly zero gradient. Training continues, billing continues, learning is zero.", 9.20, 2.20, 3.25, 1.30, size=11, color=INK)
    add_text(s, "That's what I call signal starvation — and it's invisible on the reward curve.", 9.20, 3.90, 3.25, 0.62, size=10.5, color=AMBER, bold=True)
    add_text(s, "Dr.GRPO (2025) criticizes GRPO from another angle — I tested that critique too.", 9.20, 4.85, 3.25, 0.62, size=9.5, color=MUTED, italic=True)
    add_footer(s, "Sources: thesis abstract; platform_tinker/reports/final; execution-notes.md")


    # ── 4. My measuring stick: ZVF ──────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "The measurement", "ZVF: one number that counts the empty drills", 4)
    add_text(s, "ZVF = the fraction of groups where every answer got the same grade. High ZVF → most drills teach nothing. Low ZVF → plenty of contrast, plenty of learning.", 0.65, 1.42, 7.4, 0.90, size=13, color=INK)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, 2.50, 7.40, 0.90, PANEL_2, GRID, radius=True)
    add_text(s, "The math is exact, not approximate:", 0.88, 2.68, 6.9, 0.26, size=12, color=TEAL, bold=True)
    add_text(s, "pass@G − p^G = 1 − ZVF — verified to 1.11e-16 (machine precision) across a 505-question audit", 0.88, 3.00, 6.9, 0.26, size=10.5, color=INK)
    add_bullet(s, "Cheap: computed from the same batch the trainer already has — no extra runs.", 0.68, 3.75, 7.3, 0.55, size=12, color=INK, bullet_color=TEAL)
    add_bullet(s, "Trustworthy: Wilson confidence interval covers 0.95–0.98 of cases in every setting I tested.", 0.68, 4.45, 7.3, 0.55, size=12, color=INK, bullet_color=BLUE)
    add_bullet(s, "Online: you can watch it during training, not after.", 0.68, 5.15, 7.3, 0.55, size=12, color=INK, bullet_color=LAV)
    add_bullet(s, "It separates two things the reward smashes together: \"the model is good\" vs \"training is still moving\".", 0.68, 5.85, 7.3, 0.55, size=12, color=INK, bold=True, bullet_color=AMBER)
    add_picture(s, FIG / "zvf_heatmap.png", 8.40, 1.42, 4.30, 1.60)
    add_text(s, "The wall at small group sizes is visible in the heatmap before the reward ever collapses.", 8.45, 3.10, 4.2, 0.30, size=8.8, color=MUTED)
    add_shape(s, MSO_SHAPE.RECTANGLE, 8.40, 3.55, 4.30, 2.90, PANEL, GRID, radius=True)
    add_text(s, "WHAT IT'S FOR", 8.65, 3.82, 2.5, 0.24, size=9.5, color=AMBER, bold=True)
    add_text(s, "Not a dashboard for show — an early-warning number a trainer can act on: resample the prompt, change group size, warm-start, or stop the run. Acting on it safely is the next experiment (slide 12).", 8.65, 4.18, 3.8, 1.30, size=10.5, color=INK)
    add_text(s, "No controller claim yet — we test before we sell.", 8.65, 5.80, 3.8, 0.28, size=10, color=RED, bold=True)
    add_footer(s, "Sources: thesis abstract; PAPERS_README.md item 1; figures/v2/zvf_heatmap.png")


    # ── 5. Two semesters of work ────────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "The program", "Two semesters, one question, twelve live manuscripts", 5)
    add_metric(s, 0.65, 1.42, 2.30, 1.20, "1,012", "commits", TEAL, "every claim versioned")
    add_metric(s, 3.12, 1.42, 2.30, 1.20, "12", "active manuscripts", BLUE, "P1–P12; 6 absorbed")
    add_metric(s, 5.59, 1.42, 2.30, 1.20, "486", "paper pages", LAV, "today's portfolio audit")
    add_metric(s, 8.06, 1.42, 2.30, 1.20, "40/40", "audit runs done", AMBER, "the honest big test")
    add_metric(s, 10.53, 1.42, 2.30, 1.20, "1", "thesis", GREEN, "ZVF as M.Tech thesis")
    phases = [
        ("SEMESTER 3 · GROUP OF SIX", "Built the test ground", TEAL),
        ("SEMESTER 4 · SOLO", "Found the discovery", BLUE),
    ]
    for i, (kicker, title, color) in enumerate(phases):
        x = 0.65 + i * 6.12
        add_shape(s, MSO_SHAPE.RECTANGLE, x, 2.90, 5.92, 1.86, PANEL, GRID, radius=True)
        add_shape(s, MSO_SHAPE.RECTANGLE, x, 2.90, 5.92, 0.06, color)
        add_text(s, kicker, x + 0.25, 3.14, 5.4, 0.24, size=10, color=color, bold=True)
        add_text(s, title, x + 0.25, 3.50, 5.4, 0.34, size=15, color=INK, bold=True, font=FONT_DISPLAY)
        add_text(s, ("TinkerRL-Bench: a unified test ground for RL training of language models — 4 frameworks, 6 backends, math tasks from addition to grade-school word problems. Submitted to NeurIPS 2026 with a full reproducible artifact." if i == 0
                     else "The Zero-Variance Fraction: diagnosing signal starvation, group-size dial, reporting standard (MIN-REPORT-RL), registry — plus the E1 audit that made the claims honest."),
                 x + 0.25, 3.96, 5.45, 0.72, size=10.2, color=MUTED)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, 5.00, 12.05, 1.10, PANEL_2, GRID, radius=True)
    add_text(s, "THE THREAD", 0.88, 5.24, 2.0, 0.24, size=9.5, color=AMBER, bold=True)
    add_text(s, "Phase 1 proved we could run RL training anywhere and measure it. Phase 2 found something worth measuring — and then spent the last month making sure every claim about it would survive a skeptical reviewer.", 0.88, 5.58, 11.5, 0.44, size=10.5, color=INK)
    add_footer(s, "Sources: PROJECT_HISTORY.md; ARTIFACT.md; sem 4 work/README.md; autoresearch/deli-neurips-tmlr-260802/audits/18_PAPER_PORTFOLIO_REVIEW.md")


    # ── 6. The big findings ─────────────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "The findings", "Four things I measured — in plain words", 6)
    cards = [
        ("STARVATION IS REAL", "Small groups (G=2) end in a wall where every drill is empty — ZVF ≈ 0.75–1.0 while reward reads 1.0. Big groups (G=16) keep contrast the whole run (ZVF ≤ 0.25).", RED),
        ("GROUP SIZE IS THE DIAL", "G=4 wins a 505-task utility audit. But there is NO single best G — I proved the per-rollout math always favors small groups, and I report that negative result, not hide it.", AMBER),
        ("SAME RECIPE, DIFFERENT KITCHEN", "One undisclosed change (backend + checkpoint) moved final reward across a 17× span — 86% down to 5% under the same label. That's why I demand you report what actually ran.", BLUE),
        ("DR.GRPO'S CRITIQUE: REAL MATH, NO FOOTPRINT HERE", "The 1/L length penalty should cause rambling. I ran 6 versions × 3 seeds: lengths shrank 3.8–12.2% everywhere — no verbosity trap at my scale.", TEAL),
    ]
    for i, (title, body, color) in enumerate(cards):
        x = 0.65 + (i % 2) * 6.12
        y = 1.50 + (i // 2) * 2.62
        add_shape(s, MSO_SHAPE.RECTANGLE, x, y, 5.92, 2.42, PANEL, GRID, radius=True)
        add_shape(s, MSO_SHAPE.RECTANGLE, x, y, 5.92, 0.06, color)
        add_text(s, title, x + 0.22, y + 0.20, 5.4, 0.46, size=13, color=color, bold=True)
        add_text(s, body, x + 0.22, y + 0.76, 5.45, 1.45, size=10.8, color=INK)
    add_text(s, "Every one of these numbers traces to a named file — nothing typed by hand into prose.", 0.68, 6.60, 12.0, 0.28, size=11, color=MUTED, italic=True)
    add_footer(s, "Sources: thesis abstract Claim 2; PAPERS_README.md items 2–5; defense transcript ELI12 slides 5–11")

    # ── 7. The honest big test (E1) ─────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "The honest big test", "40 runs on Google's free A100s — every method finished, no method won", 7)
    add_metric(s, 0.65, 1.42, 2.27, 1.30, "40/40", "runs finished", TEAL, "8 seeds × 5 methods")
    add_metric(s, 3.09, 1.42, 2.27, 1.30, "5", "methods raced", BLUE, "GRPO, DAPO, GSPO, Dr.GRPO, AERO")
    add_metric(s, 5.53, 1.42, 2.27, 1.30, "500", "fresh questions each", LAV, "same questions for every method")
    add_metric(s, 7.97, 1.42, 2.27, 1.30, "≈63–65%", "held-out scores", AMBER, "all methods nearly tied")
    add_metric(s, 10.41, 1.42, 2.27, 1.30, "4", "verdicts", RED, "all INCONCLUSIVE")
    add_text(s, "What the test was", 0.68, 3.00, 3.0, 0.30, size=15, color=INK, bold=True, font=FONT_DISPLAY)
    add_bullet(s, "Same model (Qwen3-8B), same questions (GSM8K), same 30 steps, same LoRA, on NVIDIA A100s — the only thing that changed was the training method.", 0.68, 3.50, 7.6, 0.62, size=11.5, color=INK, bullet_color=TEAL)
    add_bullet(s, "Six older run records were missing checksums — I repaired all six by replaying the exact checkpoints before counting them.", 0.68, 4.22, 7.6, 0.62, size=11.5, color=INK, bullet_color=BLUE)
    add_bullet(s, "Everything finished. Scores are near-identical. The correct statistics say: we cannot tell the methods apart at 8 seeds.", 0.68, 4.94, 7.6, 0.62, size=11.5, color=INK, bold=True, bullet_color=AMBER)
    add_shape(s, MSO_SHAPE.RECTANGLE, 8.55, 3.00, 4.15, 2.62, PANEL, GRID, radius=True)
    add_text(s, "THE HONEST RESULT", 8.80, 3.26, 3.6, 0.24, size=9.5, color=AMBER, bold=True)
    add_text(s, "No method won.\nWe say so out loud.\n\nThat is what fail-closed preregistration is for: you write the stop rule BEFORE you spend the GPU hours — and then you obey it.", 8.80, 3.62, 3.7, 1.85, size=11, color=INK)
    add_text(s, "This audit is now the adoption gate: don't believe an algorithm claim until it passes a test like this.", 0.68, 5.95, 12.0, 0.34, size=11.5, color=TEAL, bold=True)
    add_footer(s, "Sources: zvf-program/audit/COLAB_EXECUTION_STATUS.md; results/audit.json")


    # ── 8. We caught our own mistake ────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "The correction", "We caught our own over-claim — and fixed it before anyone could", 8)
    add_text(s, "Our first analysis said DAPO \"DISAPPEARS\" — a real claim. Then we re-checked the statistics the way we promised in the preregistration:", 0.65, 1.42, 11.9, 0.55, size=13, color=INK)
    steps = [
        ("01", "The old math used a shortcut", "A large-sample approximation that over-estimates power at only 8 paired seeds.", RED),
        ("02", "The promised check was never run", "The Benjamini–Hochberg multiplicity step existed but was never called. Now it is.", AMBER),
        ("03", "Exact math says: not enough evidence", "Exact paired-t MDE80 = 0.01012 — just above the 0.01 margin the rules require.", BLUE),
        ("04", "Verdict flipped, out loud", "DAPO, GSPO, Dr.GRPO, AERO → all INCONCLUSIVE. Superseded everywhere, including our own talks.", TEAL),
    ]
    for i, (num, title, body, color) in enumerate(steps):
        y = 2.18 + i * 1.10
        add_shape(s, MSO_SHAPE.OVAL, 0.85, y + 0.06, 0.62, 0.62, color, color)
        add_text(s, num, 0.85, y + 0.22, 0.62, 0.22, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, title, 1.75, y - 0.02, 6.5, 0.28, size=13.5, color=color, bold=True)
        add_text(s, body, 1.75, y + 0.34, 7.4, 0.44, size=10.8, color=INK)
        if i < 3:
            add_line(s, 1.16, y + 0.72, 1.16, y + 1.06, GRID, 1.1)
    add_shape(s, MSO_SHAPE.RECTANGLE, 9.45, 2.18, 3.25, 4.26, PANEL, GRID, radius=True)
    add_text(s, "WHY I TELL YOU THIS", 9.70, 2.46, 2.8, 0.40, size=9.5, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "A result you can't defend isn't a result — it's a liability.\n\nThe machinery that caught our mistake (exact power + multiplicity correction) is now the standard for every verdict in this lab.", 9.70, 3.20, 2.75, 1.60, size=10.5, color=INK)
    add_text(s, "Same runs. Same scores.\nOnly the statistics\nchanged — and with them,\nthe honest answer.", 9.70, 5.10, 2.75, 1.10, size=10.5, color=TEAL, bold=True)
    add_footer(s, "Sources: zvf-program/audit/STATISTICAL_REANALYSIS.md; execution-notes.md 2026-08-02 correction")

    # ── 9. The NeurIPS review story ─────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "The review", "NeurIPS said: good idea, over-claimed paper. We agreed and fixed it.", 9)
    add_text(s, "Reviewers found the idea valuable (separate online reward, proxies, held-out capability, and labels) — and listed 17 weaknesses, 5 critical.", 0.65, 1.42, 7.3, 0.62, size=12.5, color=INK)
    add_bullet(s, "\"ZVF might just be reward in disguise\" — answered with a formalization and partial-correlation evidence.", 0.68, 2.22, 7.2, 0.52, size=11.5, color=INK, bullet_color=RED)
    add_bullet(s, "\"G=32 contradicts your own G=8 table\" — reconciled on token-budget-normalized sweeps.", 0.68, 2.86, 7.2, 0.52, size=11.5, color=INK, bullet_color=RED)
    add_bullet(s, "\"Single-seed runs power your headline\" — acknowledged; single-seed rows became exploratory only.", 0.68, 3.50, 7.2, 0.52, size=11.5, color=INK, bullet_color=AMBER)
    add_bullet(s, "\"You never compared with AERO/CPPO/NGRPO/Scaf-GRPO\" — added to related work and the E1 campaign.", 0.68, 4.14, 7.2, 0.52, size=11.5, color=INK, bullet_color=AMBER)
    add_bullet(s, "Reviewer 9kjk: \"missing cells are missing evidence\" — we answered directly, asked for no score change.", 0.68, 4.78, 7.2, 0.52, size=11.5, color=INK, bold=True, bullet_color=TEAL)
    add_shape(s, MSO_SHAPE.RECTANGLE, 8.25, 1.42, 4.45, 4.06, PANEL, GRID, radius=True)
    add_text(s, "THE HARD PART — WE AUDITED OURSELVES", 8.50, 1.70, 4.0, 0.28, size=10, color=AMBER, bold=True)
    self_pts = [
        "Withdrew the 92.6/92.1 five-seed claim: 3 of 5 W&B seeds were zero-runtime backfills with no upstream IDs.",
        "Quarantined the PPO row: 0.225 vs 0.350 were two different runs misattributed to one row.",
        "Removed an AUROC built on synthetic anchors; re-described p=.256 as a one-sample test.",
        "Replaced all three responses + the AC comment. Response sizes: 4,854 / 5,661 / 4,733 chars.",
    ]
    for i, body in enumerate(self_pts):
        y = 2.16 + i * 0.82
        add_shape(s, MSO_SHAPE.OVAL, 8.48, y + 0.08, 0.16, 0.16, AMBER, AMBER)
        add_text(s, body, 8.72, y, 3.85, 0.74, size=9.6, color=INK)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, 5.60, 12.05, 0.92, PANEL_2, GRID, radius=True)
    add_text(s, "THE LESSON", 0.88, 5.82, 2.0, 0.24, size=9.5, color=TEAL, bold=True)
    add_text(s, "The review killed an over-claimed omnibus paper — not the diagnostic. Concede scope · correct the record · preserve the methodology. The narrow methods paper is publishable now; the big empirical claims wait for new prospective evidence.", 0.88, 6.14, 11.55, 0.34, size=10.5, color=INK)
    add_footer(s, "Sources: reviewer_points.yaml; NEURIPS_2026_OPENREVIEW_REBUTTAL_FINAL.md; NEURIPS_2026_REVIEWER_9KJK_FOLLOWUP.md; autoresearch/reason-260727-2155/")


    # ── 10. The rules I follow now ──────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "The rules", "So every claim has to survive five checkpoints", 10)
    rules = [
        ("REPORT WHAT RAN", "MIN-REPORT-RL: eight items every RL training report must state — each one earned its place by flipping a result in my own data.", TEAL),
        ("PREREGISTER BEFORE RUNNING", "The protocol (hash-locked) is written before the GPU hours are spent: what we'll test, what counts, what stops us.", BLUE),
        ("FAIL CLOSED", "If a gate can't pass, we say blocked — not \"almost there\". Right now the flagship gate is fail-closed on a real mathematical contradiction we proved.", RED),
        ("HASH EVERY RECEIPT", "Checkpoints, manifests, corpora carry SHA-256 hashes. An independent verifier re-downloads and re-checks everything.", LAV),
        ("CHECK IT, THEN EXPLAIN IT", "No claim goes into a talk, thesis, or paper until the evidence file exists. That rule caught our own DAPO mistake.", AMBER),
    ]
    for i, (title, body, color) in enumerate(rules):
        y = 1.55 + i * 1.00
        add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, y, 12.05, 0.86, PANEL, GRID, radius=True)
        add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, y, 0.06, 0.86, color)
        add_text(s, title, 0.95, y + 0.10, 3.6, 0.28, size=12.5, color=color, bold=True)
        add_text(s, body, 4.65, y + 0.10, 7.9, 0.64, size=10.6, color=INK)
    add_text(s, "These aren't bureaucracy — each rule exists because a real mistake taught it to me.", 0.68, 6.60, 12.0, 0.28, size=11.5, color=MUTED, italic=True)
    add_footer(s, "Sources: PAPERS_README.md (P05/P06); execution-notes.md current gate; NEURIPS_2026_REVIEWER_9KJK_FOLLOWUP.md")


    # ── 11. Check my work yourself ──────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "Verification", "You can check my work in under ten minutes", 11)
    add_metric(s, 0.65, 1.45, 2.90, 1.30, "< 10 min", "smoke test", TEAL, "7 checks, one command")
    add_metric(s, 3.72, 1.45, 2.90, 1.30, "88", "focused tests", BLUE, "64 audit + 24 S1, passing")
    add_metric(s, 6.79, 1.45, 2.90, 1.30, "7/7", "integration checks", LAV, "paper ↔ results ↔ checklists")
    add_metric(s, 9.86, 1.45, 2.90, 1.30, "±2 pp", "claim re-check", AMBER, "headline number reproduced")
    add_text(s, "Everything is public", 0.68, 3.05, 3.0, 0.30, size=15, color=INK, bold=True, font=FONT_DISPLAY)
    add_bullet(s, "GitHub (pes-llm-research/tinker-rl-lab + mirror) — 1,001 commits, every claim versioned.", 0.68, 3.55, 11.8, 0.42, size=12, color=INK, bullet_color=TEAL)
    add_bullet(s, "W&B — 153 public runs; the 983-run audit workbook with links.", 0.68, 4.10, 11.8, 0.42, size=12, color=INK, bullet_color=BLUE)
    add_bullet(s, "Hugging Face — frozen checkpoints and model cards; Docker OCI label records the exact commit.", 0.68, 4.65, 11.8, 0.42, size=12, color=INK, bullet_color=LAV)
    add_bullet(s, "ACM artifact review targeted: Available · Functional (< 10 min) · Reusable.", 0.68, 5.20, 11.8, 0.42, size=12, color=INK, bullet_color=AMBER)
    add_text(s, "A reviewer who can't run it isn't reviewing — they're guessing. So I made running it the easy part.", 0.68, 6.10, 11.9, 0.30, size=12, color=TEAL, bold=True)
    add_footer(s, "Sources: ARTIFACT.md; CHANGELOG.md; platform_local/run_all_audits.py")

    # ── 12. Where I am — next three gates ───────────────────────────────────
    s = new_slide(prs)
    add_header(s, "What's next", "Three gates between here and the flagship result", 12)
    steps = [
        ("01", "Authorize one small amendment", "The fail-closed gate needs an explicit fix for the joint-zero-gradient case (59 all-wrong + 3 all-correct groups of 100). It's a paperwork + math fix — the only blocker before any GPU spend.", RED),
        ("02", "Run the 24-unit pilot", "4 conditions × 2 regimes × 3 seeds on A100, replay-locked corpora, hash receipts on every step. Dry-run plans already exist — execution is the remaining step.", TEAL),
        ("03", "Then earn the flagship claim", "Matched multi-seed evidence (P1), gradient geometry (P2), token-matched group sweep (P3), uncapped mediation (P4) — then, and only then, a ZVF-aware controller bakeoff against static G=16.", BLUE),
    ]
    for i, (num, title, body, color) in enumerate(steps):
        y = 1.70 + i * 1.55
        add_shape(s, MSO_SHAPE.OVAL, 0.85, y + 0.10, 0.78, 0.78, color, color)
        add_text(s, num, 0.85, y + 0.31, 0.78, 0.26, size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, title, 1.95, y - 0.02, 8.35, 0.30, size=16, color=INK, bold=True, font=FONT_DISPLAY)
        add_text(s, body, 1.95, y + 0.42, 8.35, 0.95, size=11.5, color=MUTED)
        if i < 2:
            add_line(s, 1.25, y + 0.92, 1.25, y + 1.52, GRID, 1.2)
    add_shape(s, MSO_SHAPE.RECTANGLE, 10.55, 1.70, 2.15, 4.62, PANEL, GRID, radius=True)
    add_text(s, "MY RULE", 10.80, 2.00, 1.65, 0.40, size=9, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "No claim until the pilot runs and the receipts verify.", 10.68, 2.72, 1.90, 2.2, size=16, color=INK, bold=True, font=FONT_DISPLAY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(s, "check it → then explain it", 10.63, 5.55, 2.0, 0.55, size=8.5, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s, "Sources: execution-notes.md current gate; pilot_preregistration.json")

    # ── 13. Remember ────────────────────────────────────────────────────────
    s = new_slide(prs)
    add_header(s, "Bottom line", "Five things to remember", 13)
    takes = [
        ("A real discovery", "ZVF counts the empty drills — exact math, verified to machine precision on 505 questions, two bounded claims in the thesis.", TEAL),
        ("A serious benchmark", "TinkerRL-Bench: NeurIPS 2026 submitted, ACM artifact badges targeted, 983 runs audited.", BLUE),
        ("Standards that stick", "MIN-REPORT-RL (8 items) + GRPO-Registry (7 fields) — built from mistakes I actually made.", LAV),
        ("Honest statistics", "We corrected our own over-claim (DAPO) with exact math — before anyone else could.", AMBER),
        ("Fail-closed discipline", "Preregistered, hash-bound, replay-controlled: no claim until the evidence clears — on every gate.", GREEN),
    ]
    for i, (title, body, color) in enumerate(takes):
        y = 1.55 + i * 1.02
        add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, y, 12.05, 0.88, PANEL, GRID, radius=True)
        add_shape(s, MSO_SHAPE.RECTANGLE, 0.65, y, 0.06, 0.88, color)
        add_text(s, title, 0.95, y + 0.10, 3.3, 0.28, size=12.5, color=color, bold=True)
        add_text(s, body, 4.35, y + 0.10, 8.1, 0.66, size=10.8, color=INK)
    add_text(s, "The result is still to be earned. Everything that earns it is now checkable.", 0.68, 6.62, 12.0, 0.30, size=12.5, color=AMBER, bold=True)
    add_footer(s, "Primary references: thesis/main.tex; ARTIFACT.md; PAPERS_README.md; STATISTICAL_REANALYSIS.md")

    # ── 14. Thank you ───────────────────────────────────────────────────────
    s = new_slide(prs)
    add_text(s, "THANK YOU", 0.85, 1.90, 4.0, 0.60, size=34, color=TEAL, bold=True, font=FONT_DISPLAY)
    add_text(s, "Questions are the best part — I'll answer with the evidence, not the hope.", 0.88, 2.80, 7.0, 0.32, size=14, color=MUTED)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0.85, 3.60, 6.5, 1.75, PANEL, GRID, radius=True)
    add_text(s, "IN ONE SENTENCE", 1.10, 3.88, 2.5, 0.24, size=9.5, color=TEAL, bold=True)
    add_text(s, "I found a hidden failure in how we train thinking models, built a number that sees it, proved the number is trustworthy, and set up rules so no claim outruns the evidence.", 1.10, 4.24, 6.0, 0.95, size=12, color=INK)
    add_text(s, "Repository · github.com/pes-llm-research/tinker-rl-lab · mirror github.com/arvindcr4/tinker-rl-lab", 0.88, 5.75, 11.5, 0.26, size=10.5, color=INK)
    add_text(s, "Arvind C R · arvindcr4@gmail.com · M.Tech, PES University · guide Prof. Ramesh Prakash Guledgudd", 0.88, 6.11, 11.5, 0.26, size=10.5, color=INK)
    add_text(s, "02 Aug 2026 · every number on every slide traces to a named file", 0.88, 6.47, 11.5, 0.24, size=9, color=MUTED)

    return prs


if __name__ == "__main__":
    deck = build()
    OUT.parent.mkdir(parents=True, exist_ok=True)
    deck.save(OUT)
    print(f"saved {OUT}")
    print(f"slides {len(deck.slides)}")
