#!/usr/bin/env python3
"""Build a concise, repository-grounded three-week progress update deck.

The deck intentionally distinguishes:
  * completed evidence and verification work;
  * prospective design / contract work; and
  * preflight artifacts that are not scientific evidence.

All visible numbers are read from the checkout at build time where practical.
"""

from __future__ import annotations

import json
import re
import statistics
import subprocess
from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs" / "Tinker_RL_Progress_Update_2026-08-02.pptx"

W, H = 13.333, 7.5

# Palette: dark research-notebook navy with restrained signal colors.
NAVY = "0B1220"
PANEL = "111C2E"
PANEL_2 = "17243A"
INK = "F4F7FB"
MUTED = "A7B5C9"
TEAL = "2DD4BF"
BLUE = "60A5FA"
LAV = "A78BFA"
AMBER = "FBBF24"
RED = "FB7185"
GREEN = "86EFAC"
WHITE = "FFFFFF"
GRID = "2B3A52"

FONT = "Aptos"
FONT_DISPLAY = "Aptos Display"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def read_json(path: str) -> dict:
    return json.loads((ROOT / path).read_text(encoding="utf-8"))


def git_count() -> int:
    before = subprocess.check_output(
        ["git", "rev-list", "-1", "--before=2026-07-12 00:00", "HEAD"],
        cwd=ROOT,
        text=True,
    ).strip()
    return int(
        subprocess.check_output(
            ["git", "rev-list", "--count", f"{before}..HEAD"], cwd=ROOT, text=True
        ).strip()
    )


def e1_rows() -> dict[str, list[float]]:
    text = (ROOT / "zvf-program/audit/COLAB_EXECUTION_STATUS.md").read_text(encoding="utf-8")
    rows: dict[str, list[float]] = {}
    for line in text.splitlines():
        match = re.match(
            r"\|\s*(GRPO|DAPO|GSPO|Dr\.GRPO|AERO)\s+\|\s*\d+\s+\|\s*\d+/500\s*=\s*([0-9.]+)",
            line,
        )
        if match:
            rows.setdefault(match.group(1), []).append(float(match.group(2)))
    return rows


def e1_summary() -> dict:
    rows = e1_rows()
    return {
        "units": 40,
        "heldout": 500,
        "seeds_per_arm": 8,
        "means": {key: statistics.mean(value) for key, value in rows.items()},
        "verdicts": {"DAPO": "INCONCLUSIVE", "GSPO": "INCONCLUSIVE", "Dr.GRPO": "INCONCLUSIVE", "AERO": "INCONCLUSIVE"},
    }


def add_shape(slide, kind, x, y, w, h, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else kind
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    shape.line.width = Pt(0.75)
    return shape


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False, font=FONT,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.04, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, runs: Sequence[tuple[str, str, bool, float]], x, y, w, h,
                  align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.04):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    for text, color, bold, size in runs:
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_line(slide, x1, y1, x2, y2, color=GRID, width=1.0, dash=None):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = dash
    return line


def add_tag(slide, label, x, y, color=TEAL, w=None):
    w = w or max(0.55, 0.11 * len(label) + 0.22)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.27, color)
    add_text(slide, label.upper(), x, y + 0.015, w, 0.23, size=8.5, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)


def add_header(slide, kicker, title, number):
    add_text(slide, kicker.upper(), 0.55, 0.26, 4.5, 0.22, size=8.5, color=TEAL, bold=True)
    add_text(slide, title, 0.55, 0.54, 11.9, 0.52, size=26, color=INK, bold=True, font=FONT_DISPLAY)
    add_text(slide, f"{number:02d}", 12.27, 0.28, 0.45, 0.25, size=10, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)
    add_line(slide, 0.55, 1.17, 12.78, 1.17, GRID, 0.8)


def add_footer(slide, source: str):
    add_line(slide, 0.55, 7.14, 12.78, 7.14, GRID, 0.7)
    add_text(slide, source, 0.58, 7.18, 11.7, 0.17, size=6.8, color=MUTED, italic=True, margin=0)


def add_metric(slide, x, y, w, h, value, label, accent=TEAL, detail=None):
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, PANEL, GRID, radius=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.06, h, accent)
    add_text(slide, value, x + 0.22, y + 0.18, w - 0.35, 0.42, size=26, color=INK, bold=True, font=FONT_DISPLAY)
    add_text(slide, label, x + 0.22, y + 0.64, w - 0.35, 0.28, size=10.5, color=MUTED, bold=True)
    if detail:
        add_text(slide, detail, x + 0.22, y + 0.94, w - 0.35, h - 1.02, size=8.5, color=MUTED)


def add_bullet(slide, text, x, y, w, color=INK, size=13, bullet_color=TEAL, h=0.42, bold_prefix=None):
    add_shape(slide, MSO_SHAPE.OVAL, x, y + 0.10, 0.12, 0.12, bullet_color, bullet_color)
    if bold_prefix and text.startswith(bold_prefix):
        add_rich_text(slide, [(bold_prefix, color, True, size), (text[len(bold_prefix):], color, False, size)],
                      x + 0.24, y, w - 0.24, h)
    else:
        add_text(slide, text, x + 0.24, y, w - 0.24, h, size=size, color=color)


def add_status_chip(slide, label, x, y, status, w=1.55):
    colors = {"DONE": GREEN, "IN PROGRESS": AMBER, "BLOCKED": RED, "PROPOSED": LAV}
    c = colors[status]
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.27, c, c, radius=True)
    add_text(slide, label.upper(), x, y + 0.005, w, 0.23, size=7.7, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)


def new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(NAVY)
    return slide


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    commits = git_count()
    e1 = e1_summary()
    gate = read_json("zvf-program/next-submission/results/preflight/preflight_gate.json")
    followup = read_json("zvf-program/experiments-next/rlhfbook_followup_preregistration.json")
    offline = read_json("zvf-program/experiments-next/offline_falsification_packet.json")

    # 1 — title
    s = new_slide(prs)
    add_text(s, "TINKER RL LAB", 0.65, 0.53, 3.0, 0.25, size=10, color=TEAL, bold=True)
    add_text(s, "My last three weeks", 0.65, 1.32, 7.0, 1.20, size=41, color=INK, bold=True, font=FONT_DISPLAY)
    add_text(s, "What I finished, what I learned, and what I still need to test", 0.69, 3.14, 7.5, 0.42, size=17, color=MUTED)
    add_text(s, "12 JUL — 02 AUG 2026  /  PRESENTATION TODAY", 0.69, 3.76, 5.4, 0.25, size=9.5, color=AMBER, bold=True)
    # Timeline ribbon
    add_line(s, 0.72, 5.60, 12.45, 5.60, GRID, 2.3)
    milestones = [
        (1.12, "12 JUL", "Defense\npackaged", TEAL),
        (4.08, "20–21 JUL", "Old-run audit\n40/40", BLUE),
        (7.05, "29 JUL", "New\nchecklist", LAV),
        (10.02, "30 JUL–02 AUG", "Remote setup\n& truth audit", AMBER),
    ]
    for x, date, label, color in milestones:
        add_shape(s, MSO_SHAPE.OVAL, x, 5.42, 0.36, 0.36, color, color)
        add_text(s, date, x - 0.34, 5.97, 1.05, 0.22, size=8.2, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, label, x - 0.55, 6.23, 1.5, 0.45, size=9.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "The short version: I finished the checks and plans. The next experiment is not ready to call a win.",
             0.69, 6.90, 11.8, 0.23, size=9, color=MUTED, italic=True)

    # 2 — executive summary
    s = new_slide(prs)
    add_header(s, "The short version", "What I did in three weeks", 2)
    add_text(s, "I moved from finishing the presentation to making sure the next experiment can be checked properly.",
             0.65, 1.45, 11.7, 0.35, size=15, color=INK)
    cards = [
        ("FINISHED", "Defense pack", "Slides, demo, workbook and evidence links were ready on 12 Jul.", GREEN, "DONE"),
        ("FINISHED", "Old-run audit", "I checked all 40 old test runs and kept the unclear results marked unclear.", BLUE, "DONE"),
        ("PLANNED", "Next experiment", "I wrote down the assumptions, checks and stop rules before more GPU time.", LAV, "DONE"),
        ("WAITING", "New test run", "The setup works, but three of four cells still lack a real update.", RED, "BLOCKED"),
    ]
    xs = [0.65, 3.78, 6.91, 10.04]
    for (kicker, title, body, accent, status), x in zip(cards, xs):
        add_shape(s, MSO_SHAPE.RECTANGLE, x, 2.08, 2.65, 3.65, PANEL, GRID, radius=True)
        add_shape(s, MSO_SHAPE.RECTANGLE, x, 2.08, 2.65, 0.08, accent)
        add_text(s, kicker, x + 0.2, 2.34, 2.2, 0.24, size=8.7, color=accent, bold=True)
        add_text(s, title, x + 0.2, 2.73, 2.2, 0.52, size=16, color=INK, bold=True, font=FONT_DISPLAY)
        add_text(s, body, x + 0.2, 3.50, 2.2, 1.15, size=10.2, color=MUTED)
        add_status_chip(s, status, x + 0.2, 5.16, status, w=1.02 if status != "BLOCKED" else 1.15)
    add_text(s, "This is a progress update, not a claim that the new method works.", 0.68, 6.27, 8.5, 0.35, size=14, color=AMBER, bold=True)
    add_footer(s, "Sources: git history 2026-07-12..2026-08-02; COLAB_EXECUTION_STATUS.md; rlhfbook_followup_preregistration.json; preflight_gate.json")

    # 3 — timeline
    s = new_slide(prs)
    add_header(s, "Timeline", "What happened when", 3)
    add_text(s, "Each step made the next experiment a little easier to check.",
             0.65, 1.42, 11.4, 0.30, size=14.5, color=MUTED)
    # Vertical axis
    add_line(s, 1.36, 2.14, 1.36, 6.42, GRID, 1.2)
    phases = [
        ("12 JUL", "Finish the defense", "I rebuilt the 17-slide defense, added the demo, and linked each important result to its source.", TEAL, "delivery"),
        ("20–21 JUL", "Check the old runs", "I finished the 40/40 audit and kept the unclear results marked unclear.", BLUE, "evidence"),
        ("29 JUL", "Write the new checklist", "I turned the RLHF Book and CS2824 notes into simple checks and stop rules.", LAV, "design"),
        ("30 JUL–02 AUG", "Lock the publication story", "I checked the reviewer claims, corrected the small-sample statistics, and kept the new experiment blocked until real updates exist.", AMBER, "publication"),
    ]
    ys = [2.12, 3.18, 4.24, 5.30]
    for (date, title, body, accent, _), y in zip(phases, ys):
        add_shape(s, MSO_SHAPE.OVAL, 1.17, y + 0.10, 0.38, 0.38, accent, accent)
        add_text(s, date, 0.65, y, 0.75, 0.26, size=8.6, color=accent, bold=True, align=PP_ALIGN.RIGHT)
        add_text(s, title, 1.78, y - 0.02, 5.4, 0.28, size=14.2, color=INK, bold=True)
        add_text(s, body, 1.78, y + 0.33, 8.8, 0.48, size=10.5, color=MUTED)
    add_shape(s, MSO_SHAPE.RECTANGLE, 10.95, 2.04, 1.75, 4.35, PANEL, GRID, radius=True)
    add_text(s, "MY RULE", 11.15, 2.32, 1.35, 0.45, size=8, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "I will not\ncall it\nlearning\nuntil I\nsee the\nupdate.", 11.13, 3.20, 1.4, 1.55, size=18, color=INK, bold=True, font=FONT_DISPLAY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(s, "check it →\nthen explain it", 11.16, 5.18, 1.38, 0.58, size=10, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s, "Sources: commit log; outputs/build_session1_deck.py; zvf-program/audit/COLAB_EXECUTION_STATUS.md; experiments-next audit files")

    # 4 — shipped evidence
    s = new_slide(prs)
    add_header(s, "What I finished", "The important pieces are in place", 4)
    add_metric(s, 0.65, 1.54, 2.35, 1.27, str(commits), "commits in window", TEAL, "12 Jul → 02 Aug in the live checkout")
    add_metric(s, 3.22, 1.54, 2.35, 1.27, "40/40", "old runs checked", BLUE, "8 seeds × 5 methods")
    add_metric(s, 5.79, 1.54, 2.35, 1.27, "500", "test questions / run", LAV, "same questions for each method")
    add_metric(s, 8.36, 1.54, 2.35, 1.27, "88", "focused tests", AMBER, "64 audit + 24 S1 passed today")
    add_metric(s, 10.93, 1.54, 1.78, 1.27, "4", "verdicts", RED, "all 4 are INCONCLUSIVE")
    add_text(s, "What I shipped", 0.68, 3.25, 3.0, 0.30, size=16, color=INK, bold=True, font=FONT_DISPLAY)
    add_bullet(s, "A defense deck with the demo, the method, and links back to the results.", 0.72, 3.78, 5.65, size=12, h=0.76, bullet_color=TEAL)
    add_bullet(s, "A workbook covering all 983 Tinker runs, with the useful links attached.", 0.72, 4.68, 5.65, size=12, h=0.76, bullet_color=TEAL)
    add_bullet(s, "A cleaned-up audit that shows what is finished, missing, or still unclear.", 0.72, 5.58, 5.65, size=12, h=0.76, bullet_color=TEAL)
    add_shape(s, MSO_SHAPE.RECTANGLE, 7.00, 3.25, 5.70, 3.16, PANEL, GRID, radius=True)
    add_text(s, "The honest takeaway", 7.30, 3.60, 5.1, 0.30, size=15, color=AMBER, bold=True)
    add_text(s, "The old campaign is fully documented.\n\nThat does not mean any method won: after the small-sample correction, DAPO, GSPO, Dr.GRPO and AERO are all still unclear.",
             7.30, 4.18, 4.95, 1.40, size=14, color=INK, bold=False)
    add_text(s, "Finishing the audit is not the same as proving a gain.", 7.30, 5.90, 4.9, 0.27, size=10.5, color=RED, bold=True)
    add_footer(s, "Sources: zvf-program/audit/COLAB_EXECUTION_STATUS.md; autoresearch/orchestrator-260730-1818/summary.md; 12 Jul deck / workbook commits")

    # 5 — E1 chart
    s = new_slide(prs)
    add_header(s, "What the old tests say", "The audit is done. The answer is still unclear.", 5)
    add_text(s, "Average score on 500 new test questions, using eight seeds per method", 0.68, 1.45, 6.3, 0.28, size=14, color=MUTED)
    means = e1["means"]
    order = ["GRPO", "DAPO", "GSPO", "Dr.GRPO", "AERO"]
    colors = [BLUE, TEAL, LAV, AMBER, RED]
    chart_x, chart_y, chart_w, chart_h = 0.95, 2.08, 6.20, 3.85
    low, high = 0.60, 0.66
    for tick in [0.60, 0.62, 0.64, 0.66]:
        yy = chart_y + chart_h - (tick - low) / (high - low) * chart_h
        add_line(s, chart_x, yy, chart_x + chart_w, yy, GRID, 0.6)
        add_text(s, f"{tick*100:.0f}%", 0.52, yy - 0.10, 0.35, 0.20, size=8, color=MUTED, align=PP_ALIGN.RIGHT)
    bar_w = 0.82
    for i, (arm, color) in enumerate(zip(order, colors)):
        val = means[arm]
        bh = max(0.10, (val - low) / (high - low) * chart_h)
        bx = chart_x + 0.35 + i * 1.18
        by = chart_y + chart_h - bh
        add_shape(s, MSO_SHAPE.RECTANGLE, bx, by, bar_w, bh, color, color, radius=True)
        add_text(s, f"{val*100:.1f}%", bx - 0.16, by - 0.30, bar_w + 0.32, 0.22, size=10, color=INK, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, arm, bx - 0.16, chart_y + chart_h + 0.12, bar_w + 0.32, 0.26, size=9.2, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "The scores are close. The labels on the right are more important.", 0.95, 6.33, 6.2, 0.28, size=12, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    # verdict panel
    add_shape(s, MSO_SHAPE.RECTANGLE, 7.62, 2.08, 5.05, 3.86, PANEL, GRID, radius=True)
    add_text(s, "What the audit says", 7.95, 2.42, 4.3, 0.30, size=16, color=INK, bold=True, font=FONT_DISPLAY)
    verdict_items = [(arm, "INCONCLUSIVE", AMBER) for arm in ("DAPO", "GSPO", "Dr.GRPO", "AERO")]
    for idx, (arm, verdict, color) in enumerate(verdict_items):
        yy = 3.12 + idx * 0.57
        add_text(s, arm, 7.98, yy, 1.3, 0.26, size=12, color=INK, bold=True)
        add_shape(s, MSO_SHAPE.RECTANGLE, 9.58, yy + 0.01, 2.58, 0.28, color, color, radius=True)
        add_text(s, verdict, 9.58, yy + 0.015, 2.58, 0.23, size=8.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
    add_text(s, "I can report the audit. I cannot report a clear improvement.", 7.98, 5.55, 4.2, 0.32, size=11.5, color=MUTED, italic=True)
    add_footer(s, "Source: zvf-program/audit/COLAB_EXECUTION_STATUS.md; values parsed from the eight-seed held-out table")

    # 6 — design upgrade
    s = new_slide(prs)
    add_header(s, "The new checklist", "I added checks before spending more GPU time", 6)
    add_text(s, "I used the RLHF Book and CS2824 notes to write down what must be true before I trust a result.",
             0.66, 1.43, 11.6, 0.30, size=14, color=MUTED)
    # ladder
    ladder = [
        ("S0", "separate", "keep the new work apart", TEAL),
        ("S1", "write it down", "state the assumptions", BLUE),
        ("S2", "check it", "compare with independent answers", LAV),
        ("S3", "make it fair", "make answer lengths vary", AMBER),
        ("S4", "train", "same cost, fresh data", RED),
        ("S5", "recheck", "use untouched questions", GREEN),
    ]
    x0, y0 = 0.77, 2.15
    for i, (code, label, detail, color) in enumerate(ladder):
        x = x0 + i * 1.98
        if i < len(ladder) - 1:
            add_line(s, x + 0.92, y0 + 0.43, x + 1.95, y0 + 0.43, GRID, 1.3)
        add_shape(s, MSO_SHAPE.OVAL, x, y0, 0.86, 0.86, color, color)
        add_text(s, code, x, y0 + 0.18, 0.86, 0.24, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, label, x - 0.15, y0 + 1.06, 1.15, 0.22, size=10.8, color=INK, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, detail, x - 0.43, y0 + 1.38, 1.72, 0.42, size=8.3, color=MUTED, align=PP_ALIGN.CENTER)
    add_metric(s, 0.78, 4.62, 2.55, 1.40, "7", "ideas checked", LAV, "each has a source, an assumption and a way to fail")
    add_metric(s, 3.57, 4.62, 2.55, 1.40, "12", "source files pinned", BLUE, "so I know which notes and code I used")
    add_metric(s, 6.36, 4.62, 2.55, 1.40, "YES", "independent answer check", TEAL, "the checker does not train the model")
    add_metric(s, 9.15, 4.62, 3.05, 1.40, "NOT YET", "small offline check", RED, "no new GPU run until this is done")
    add_text(s, "The plan is ready. The experiment has not started.", 0.80, 6.43, 7.0, 0.28, size=12.5, color=AMBER, bold=True)
    add_footer(s, "Sources: experiments-next/RLHFBOOK_IMPROVEMENT_AUDIT.md; HARVARD_CS2824_IMPROVEMENT_AUDIT.md; theory_transfer_ledger.json; offline_falsification_packet.json")

    # 7 — preflight matrix
    s = new_slide(prs)
    add_header(s, "Remote runs", "I made the setup easier to check", 7)
    add_text(s, "The same checks now run on four remote services, and each run leaves a small receipt.",
             0.66, 1.43, 11.7, 0.30, size=14, color=MUTED)
    # provider row
    providers = [("COLAB", TEAL), ("GCP A100", BLUE), ("HF JOBS", LAV), ("KAGGLE", AMBER)]
    for i, (name, color) in enumerate(providers):
        x = 0.82 + i * 1.95
        add_shape(s, MSO_SHAPE.RECTANGLE, x, 2.06, 1.55, 0.52, PANEL_2, GRID, radius=True)
        add_shape(s, MSO_SHAPE.OVAL, x + 0.15, 2.22, 0.17, 0.17, color, color)
        add_text(s, name, x + 0.39, 2.19, 1.0, 0.18, size=9.4, color=INK, bold=True)
    add_text(s, "Four test cells", 9.00, 2.19, 3.2, 0.22, size=11, color=AMBER, bold=True, align=PP_ALIGN.RIGHT)
    # matrix grid
    gx, gy = 0.82, 3.02
    cols = ["GSM8K /\ncontrast", "GSM8K /\nbase", "MATH /\ncontrast", "MATH /\nbase"]
    widths = [2.15, 2.15, 2.15, 2.15]
    for i, (label, width) in enumerate(zip(cols, widths)):
        xx = gx + sum(widths[:i])
        add_shape(s, MSO_SHAPE.RECTANGLE, xx, gy, width - 0.07, 0.62, PANEL_2, GRID, radius=True)
        add_text(s, label, xx + 0.05, gy + 0.11, width - 0.17, 0.38, size=10, color=INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    values = [("SETUP", GREEN), ("SETUP +\nUPDATE", TEAL), ("SETUP", GREEN), ("SETUP", GREEN)]
    updates = [False, True, False, False]
    for i, ((label, color), has_update) in enumerate(zip(values, updates)):
        xx = gx + sum(widths[:i])
        fill = color if has_update else PANEL
        add_shape(s, MSO_SHAPE.RECTANGLE, xx, gy + 0.82, widths[i] - 0.07, 1.30, fill, GRID, radius=True)
        add_text(s, label, xx + 0.05, gy + 1.12, widths[i] - 0.17, 0.35, size=14, color=NAVY if has_update else INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(s, "real update seen" if has_update else "real update missing", xx + 0.10, gy + 1.60, widths[i] - 0.27, 0.28, size=8.7, color=NAVY if has_update else MUTED, align=PP_ALIGN.CENTER)
    add_metric(s, 0.82, 5.55, 2.15, 1.05, "5", "receipts", BLUE, "independently checked")
    add_metric(s, 3.20, 5.55, 2.15, 1.05, "4/4", "setup checks", GREEN, "all four cells")
    add_metric(s, 5.58, 5.55, 2.15, 1.05, "1/4", "real updates", TEAL, "only GSM8K base has one")
    add_metric(s, 7.96, 5.55, 2.15, 1.05, "3", "still missing", RED, "so I cannot compare yet")
    add_metric(s, 10.34, 5.55, 2.33, 1.05, "WAIT", "new test", RED, "setup is not a result")
    add_footer(s, "Sources: next-submission/README.md; results/preflight/preflight_gate.json; verify_preflight_matrix.py; focused provider tests")

    # 8 — current blockers
    s = new_slide(prs)
    add_header(s, "Current state", "What is done, and what is still missing", 8)
    # left done
    add_shape(s, MSO_SHAPE.RECTANGLE, 0.68, 1.50, 5.75, 4.98, PANEL, GRID, radius=True)
    add_text(s, "FINISHED", 0.98, 1.82, 2.7, 0.26, size=10, color=GREEN, bold=True)
    add_bullet(s, "Defense and demo material", 1.00, 2.28, 4.85, size=13, h=0.38, bullet_color=GREEN)
    add_bullet(s, "Old-run audit: 40/40 units", 1.00, 2.82, 4.85, size=13, h=0.38, bullet_color=GREEN)
    add_bullet(s, "New experiment checklist", 1.00, 3.36, 4.85, size=13, h=0.38, bullet_color=GREEN)
    add_bullet(s, "88 focused tests passed today", 1.00, 3.90, 4.85, size=13, h=0.38, bullet_color=GREEN)
    add_bullet(s, "Remote receipts and cleanup", 1.00, 4.44, 4.85, size=13, h=0.38, bullet_color=GREEN)
    add_text(s, "These show that the setup works.", 1.00, 5.45, 4.9, 0.30, size=11, color=MUTED, italic=True)
    # right blocked
    add_shape(s, MSO_SHAPE.RECTANGLE, 6.77, 1.50, 5.88, 4.98, PANEL, GRID, radius=True)
    add_text(s, "STILL MISSING", 7.07, 1.82, 3.2, 0.26, size=10, color=RED, bold=True)
    add_text(s, "3", 7.08, 2.22, 0.85, 0.70, size=42, color=RED, bold=True, font=FONT_DISPLAY)
    add_text(s, "real updates still missing", 8.02, 2.38, 3.9, 0.30, size=14, color=INK, bold=True)
    missing = gate["missing_scientific_seams"]
    friendly_names = {
        "gsm8k/contrast_early_stop_g2_to_g8:mixed_reward_optimizer_update": "GSM8K contrast",
        "math500/contrast_early_stop_g2_to_g8:mixed_reward_optimizer_update": "MATH contrast",
        "math500/grpo_g8:mixed_reward_optimizer_update": "MATH base",
    }
    for i, item in enumerate(missing):
        label = friendly_names.get(item, item)
        add_bullet(s, label, 7.10, 3.14 + i * 0.49, 4.9, size=11.2, h=0.34, bullet_color=RED)
    add_text(s, "Why this matters", 7.10, 4.82, 2.0, 0.24, size=12, color=AMBER, bold=True)
    add_text(s, "Sometimes all answers in a group get the same score. Then a zero update is expected. That proves the setup ran; it does not prove the new method learns.",
             7.10, 5.15, 4.95, 0.82, size=11, color=MUTED)
    add_footer(s, "Source: preflight_gate.json (confirmatory_execution_gate=blocked; evidence_class=preflight-gate-not-scientific-evidence)")

    # 9 — eighteen-paper review
    s = new_slide(prs)
    add_header(s, "Paper review", "I checked all 18 drafts", 9)
    add_text(s, "Several drafts tell the same story in different formats, so the 18 files collapse into fewer paper ideas.",
             0.68, 1.43, 11.7, 0.32, size=14.5, color=MUTED)
    add_metric(s, 0.72, 2.00, 2.35, 1.25, "18/18", "drafts rebuilt", TEAL, "every current PDF compiles")
    add_metric(s, 3.30, 2.00, 2.35, 1.25, "868", "pages reviewed", BLUE, "including the 239-page internal book")
    add_metric(s, 5.88, 2.00, 2.35, 1.25, "329", "source files", LAV, "all included files entered the review")
    add_metric(s, 8.46, 2.00, 2.35, 1.25, "6 + 1", "real story groups", AMBER, "six research themes + one internal book")
    add_metric(s, 11.04, 2.00, 1.60, 1.25, "0", "ready now", RED, "none unchanged")
    buckets = [
        ("KEEP", "3", "R08 audit\nR02 short note\nR04 artifact", GREEN),
        ("MERGE", "8", "useful pieces,\nbut overlapping", BLUE),
        ("TEST FIRST", "5", "evidence still\nmissing", AMBER),
        ("ARCHIVE", "2", "R01 + U01", RED),
    ]
    for i, (heading, count, body, color) in enumerate(buckets):
        x = 0.76 + i * 3.04
        add_shape(s, MSO_SHAPE.RECTANGLE, x, 3.73, 2.70, 2.33, PANEL, GRID, radius=True)
        add_text(s, heading, x + 0.22, 3.98, 1.55, 0.24, size=9, color=color, bold=True)
        add_text(s, count, x + 1.82, 3.88, 0.54, 0.45, size=27, color=color, bold=True,
                 font=FONT_DISPLAY, align=PP_ALIGN.RIGHT)
        add_text(s, body, x + 0.22, 4.55, 2.20, 0.92, size=13, color=INK, bold=True)
    add_text(s, "The page count got bigger. The publishable story got smaller and clearer.",
             0.78, 6.43, 9.0, 0.30, size=12.5, color=AMBER, bold=True)
    add_footer(s, "Source: autoresearch/deli-neurips-tmlr-260802/audits/18_PAPER_PORTFOLIO_REVIEW.md; paper_portfolio inventory and overlap scan")

    # 10 — publication route
    s = new_slide(prs)
    add_header(s, "Publication plan", "The route I would take", 10)
    add_text(s, "The best journal paper is the separate flagship audit. The current NeurIPS review has to finish first.",
             0.68, 1.43, 11.7, 0.34, size=14.5, color=MUTED)
    route = [
        ("NOW", "Keep the evidence frozen", "Do not submit a second overlapping paper while the NeurIPS review is open.", RED),
        ("NEXT", "Send the flagship to TMLR", "If NeurIPS rejects or the paper is withdrawn, rerun the overlap and anonymity checks.", TEAL),
        ("LATER", "Build the companions", "R08 audit, a short R02 note, and the R04 artifact each need their own missing gate.", BLUE),
    ]
    for i, (when, title, body, color) in enumerate(route):
        y = 2.03 + i * 1.34
        add_shape(s, MSO_SHAPE.OVAL, 0.82, y, 0.72, 0.72, color, color)
        add_text(s, str(i + 1), 0.82, y + 0.19, 0.72, 0.26, size=15, color=NAVY,
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(s, when, 1.85, y - 0.01, 1.10, 0.23, size=8.5, color=color, bold=True)
        add_text(s, title, 1.85, y + 0.27, 4.60, 0.31, size=15.5, color=INK, bold=True,
                 font=FONT_DISPLAY)
        add_text(s, body, 6.08, y + 0.09, 5.85, 0.56, size=11.2, color=MUTED)
        if i < 2:
            add_line(s, 1.18, y + 0.75, 1.18, y + 1.27, GRID, 1.1)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0.82, 6.08, 11.80, 0.58, PANEL_2, GRID, radius=True)
    add_text(s, "The evidence supports a careful methods and failure-analysis paper; the controller result is still missing.",
             1.06, 6.23, 11.30, 0.25, size=12.5, color=AMBER, bold=True,
             align=PP_ALIGN.CENTER)
    add_footer(s, "Sources: PUBLICATION_READINESS.md; 18_PAPER_PORTFOLIO_REVIEW.md; TMLR editorial policies; active NeurIPS overlap record")

    # 11 — next steps
    s = new_slide(prs)
    add_header(s, "Next week", "My next three steps", 11)
    add_text(s, "I will do these in order. I will not skip ahead just because the setup works.", 0.68, 1.42, 10.0, 0.30, size=14, color=MUTED)
    steps = [
        ("01", "Get the three missing updates", "Run the three cells that still have setup receipts but no real update.", RED),
        ("02", "Run the small offline check", "Compare the extra score with an independent answer check before using more GPU time.", LAV),
        ("03", "Only then train more", "If the small check passes, run the matched-cost, multi-seed experiment and check it on held-out questions.", TEAL),
    ]
    for i, (num, title, body, color) in enumerate(steps):
        y = 2.12 + i * 1.33
        add_shape(s, MSO_SHAPE.OVAL, 0.82, y, 0.70, 0.70, color, color)
        add_text(s, num, 0.82, y + 0.19, 0.70, 0.25, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, title, 1.83, y - 0.02, 5.45, 0.30, size=16, color=INK, bold=True, font=FONT_DISPLAY)
        add_text(s, body, 1.83, y + 0.37, 8.9, 0.54, size=11.5, color=MUTED)
        if i < 2:
            add_line(s, 1.17, y + 0.73, 1.17, y + 1.23, GRID, 1.1)
    add_shape(s, MSO_SHAPE.RECTANGLE, 10.95, 2.05, 1.70, 4.06, PANEL, GRID, radius=True)
    add_text(s, "MY RULE", 11.18, 2.36, 1.25, 0.40, size=8.4, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "No\nclaim\nuntil\nall 4\ncells\nupdate.", 11.12, 3.16, 1.37, 2.18, size=18, color=INK, bold=True, font=FONT_DISPLAY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(s, "Status today:\nwaiting on 3 cells", 0.84, 6.47, 5.0, 0.34, size=13, color=RED, bold=True)
    add_footer(s, "Sources: preflight_gate.json; rlhfbook_followup_preregistration.json; offline_falsification_packet.json")

    # 12 — close / backup references
    s = new_slide(prs)
    add_header(s, "Bottom line", "I have a safer experiment. The result is still to be earned.", 12)
    add_text(s, "The setup is cleaner, the old work is documented, and I know exactly what is still missing.",
             0.68, 1.48, 11.7, 0.52, size=21, color=INK, bold=True, font=FONT_DISPLAY)
    cols = [
        ("DONE", ["Defense slides", "Old-run audit", "Evidence workbook"], GREEN),
        ("STRONGER", ["New checklist", "Source tracking", "Independent checks"], BLUE),
        ("OPEN", ["3 real updates", "Small offline check", "No new claim yet"], RED),
    ]
    for i, (heading, items, color) in enumerate(cols):
        x = 0.85 + i * 4.18
        add_shape(s, MSO_SHAPE.RECTANGLE, x, 2.55, 3.60, 2.65, PANEL, GRID, radius=True)
        add_text(s, heading, x + 0.25, 2.88, 3.1, 0.25, size=10, color=color, bold=True)
        for j, item in enumerate(items):
            add_shape(s, MSO_SHAPE.OVAL, x + 0.26, 3.40 + j * 0.49, 0.11, 0.11, color, color)
            add_text(s, item, x + 0.50, 3.31 + j * 0.49, 2.75, 0.25, size=12.2, color=INK, bold=True)
    add_text(s, "Thank you", 0.86, 5.90, 3.0, 0.45, size=25, color=TEAL, bold=True, font=FONT_DISPLAY)
    add_text(s, "Questions / discussion", 0.88, 6.37, 3.5, 0.25, size=11, color=MUTED)
    add_text(s, "Repository: /Users/arvind/Developer/tinker-rl-lab", 7.25, 5.98, 5.0, 0.24, size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    add_text(s, "02 Aug 2026", 7.25, 6.30, 5.0, 0.24, size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    add_footer(s, "Primary references: audit/COLAB_EXECUTION_STATUS.md; experiments-next/*; next-submission/results/preflight/preflight_gate.json")

    return prs


if __name__ == "__main__":
    deck = build()
    OUT.parent.mkdir(parents=True, exist_ok=True)
    deck.save(OUT)
    print(f"saved {OUT}")
    print(f"slides {len(deck.slides)}")
