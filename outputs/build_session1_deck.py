#!/usr/bin/env python3
"""Phase 1 First Review deck — DEFENSE FORMAT (examiner rubric):
title -> base paper & understanding -> architecture -> what I implemented
-> results achieved -> demo. 14 slides, ~20 minutes, timed speaker notes."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
import os

# ---------- matplotlib assets (deterministic, regenerated on every build) ----
os.makedirs('outputs/deck_assets', exist_ok=True)
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np

# U-shape kernel h_G(p) = p^G + (1-p)^G
fig, ax = plt.subplots(figsize=(5.2, 3.1), dpi=160)
p_ = np.linspace(0, 1, 400)
for G, c in [(2, '#1f4e79'), (8, '#c0392b'), (16, '#e67e22')]:
    ax.plot(p_, p_**G + (1-p_)**G, color=c, lw=2, label=f'G={G}')
ax.set_xlabel('per-prompt success rate p', fontsize=9)
ax.set_ylabel(r'$h_G(p)=p^G+(1{-}p)^G$', fontsize=9)
ax.text(0.03, 0.55, 'all-incorrect\nwall', fontsize=7.5, color='#555')
ax.text(0.80, 0.55, 'all-correct\nwall', fontsize=7.5, color='#555')
ax.legend(fontsize=8, loc='lower center'); ax.tick_params(labelsize=8)
ax.set_title('Zero-variance probability is U-shaped in difficulty', fontsize=9.5)
fig.tight_layout(); fig.savefig('outputs/deck_assets/ushape.png'); plt.close(fig)

# Matched-budget trajectory schematic (corrected E-R2b shape)
fig, ax = plt.subplots(figsize=(5.2, 3.1), dpi=160)
x = np.linspace(0, 1, 400)
ax.plot(x, np.minimum(1, 0.35+0.95*x), color='#1f4e79', lw=2, label='G=2: train reward')
ax.plot(x, 0.15+0.8*x**3, color='#1f4e79', lw=2, ls='--', label='G=2: ZVF')
ax.plot(x, 0.33+0.12*x, color='#c0392b', lw=2, label='G=16: train reward')
ax.plot(x, 0.10+0.10*x, color='#c0392b', lw=2, ls='--', label='G=16: ZVF')
ax.axvline(0.72, color='gray', ls=':', lw=1.4)
ax.text(0.735, 0.06, 'G=2 signal gone', fontsize=7.5, color='#555')
ax.set_xlabel('fraction of the 2,560-rollout budget spent', fontsize=9)
ax.set_ylabel('value', fontsize=9); ax.set_ylim(0, 1.05)
ax.legend(fontsize=7.5, loc='upper left'); ax.tick_params(labelsize=8)
ax.set_title('Same budget, opposite endings (schematic of measured runs)', fontsize=9.5)
fig.tight_layout(); fig.savefig('outputs/deck_assets/budget_traj.png'); plt.close(fig)

# poster frame for the opening demo video
import subprocess
VIDEO = 'outputs/project_defense_live_hf_wandb_demo_2026-07-12.mp4'
POSTER = 'outputs/deck_assets/demo_poster.png'
if os.path.exists(VIDEO):
    subprocess.run(['ffmpeg', '-y', '-ss', '75', '-i', VIDEO, '-vframes', '1', POSTER],
                   capture_output=True)
FALLBACK = 'thesis/viva/demo_walkthrough.mp4'
FPOSTER = 'outputs/deck_assets/walkthrough_poster.png'
if os.path.exists(FALLBACK) and not os.path.exists(FPOSTER):
    subprocess.run(['ffmpeg', '-y', '-i', FALLBACK, '-vf', 'select=eq(n\\,30)', '-vframes', '1', FPOSTER],
                   capture_output=True)

BLUE = RGBColor(0x1F, 0x4E, 0x79); INK = RGBColor(0x21, 0x21, 0x21)
MUTED = RGBColor(0x5A, 0x5A, 0x5A); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0x8A, 0x6D, 0x00)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def para(tf, text, size=16, color=INK, bold=False, bullet=False, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    r = p.add_run(); r.text = ('•  ' if bullet else '') + text
    f = r.font; f.size = Pt(size); f.color.rgb = color; f.bold = bold; f.name = 'Calibri'
    p.alignment = align; p.space_after = Pt(7)
    return p

def linkline(tf, label, url, size=13, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    r = p.add_run(); r.text = '\u2197  ' + label
    r.font.size = Pt(size); r.font.color.rgb = RGBColor(0x0B, 0x5C, 0xAB)
    r.font.underline = True; r.font.name = 'Calibri'
    r.hyperlink.address = url
    p.space_after = Pt(6)
    return p

def slide(title=None, page=None, notes=None):
    s = prs.slides.add_slide(BLANK)
    if title:
        tb = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.1), Inches(0.8))
        para(tb.text_frame, title, size=27, color=BLUE, bold=True, first=True)
    if page:
        pb = s.shapes.add_textbox(Inches(12.5), Inches(7.05), Inches(0.6), Inches(0.35))
        para(pb.text_frame, str(page), size=11, color=MUTED, first=True)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s

def bullets(s, items, top=1.25, size=15, width=11.8, left=0.8):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.7))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (t, kw) in enumerate(items):
        para(tf, t, size=size, bullet=True, bold=kw, first=(i == 0))
    return s

def box(s, text, left, top, w, h, fill=BLUE, font=WHITE, size=12, bold=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = s.shapes.add_shape(shape, Inches(left), Inches(top), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.color.rgb = fill
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(45720); tf.margin_right = Emu(45720)
    para(tf, text, size=size, color=font, bold=bold, align=PP_ALIGN.CENTER, first=True)
    return sh

def arrow(s, x1, y1, x2, y2, color=MUTED, w=2.0):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(w)
    return c

def table(s, rows, top=1.35, left=0.9, width=11.5, col_widths=None, size=13):
    n_r, n_c = len(rows), len(rows[0])
    shp = s.shapes.add_table(n_r, n_c, Inches(left), Inches(top), Inches(width), Inches(0.36*n_r))
    t = shp.table
    if col_widths:
        for i, w in enumerate(col_widths): t.columns[i].width = Inches(w)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci); cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(size); r.font.name = 'Calibri'
                    r.font.bold = (ri == 0); r.font.color.rgb = INK if ri else WHITE
    return t

# ------------------------------------------------------- 0 Opening demo video
if os.path.exists(VIDEO):
    sv = slide(notes="(optional 90 s) Recorded live demo: HuggingFace adapter repos + W&B zvf-training "
                     "panel walkthrough, captured 12 Jul. Click to play. Use while people settle, or skip "
                     "straight to the title if he wants to start immediately - the live demo slot is slide 14.")
    tb = sv.shapes.add_textbox(Inches(0.6), Inches(0.22), Inches(12.1), Inches(0.6))
    para(tb.text_frame, 'Opening demo — live HuggingFace + W&B walkthrough (recorded today, 91 s)',
         size=20, color=BLUE, bold=True, align=PP_ALIGN.CENTER, first=True)
    sv.shapes.add_movie(VIDEO, Inches(1.57), Inches(0.95), Inches(10.2), Inches(5.74),
                        poster_frame_image=POSTER if os.path.exists(POSTER) else None,
                        mime_type='video/mp4')
    tb = sv.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.45))
    para(tb.text_frame, 'Click to play — embedded in the deck; no network needed.',
         size=12, color=MUTED, align=PP_ALIGN.CENTER, first=True)

# ---------------------------------------------------------------- 1 Title
s = slide(notes="(1 min) Read the title, then decode it in one breath: GRPO trains on groups; "
                "when every completion in a group gets the same reward the group carries zero gradient — "
                "that is signal starvation; ZVF is the statistic that measures it; 'stack-conditioned' is "
                "the reproducibility finding that the same algorithm label gives different results on "
                "different software stacks. One measurement discipline, two claims, all artifacts audited.")
tb = s.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(11.3), Inches(4.2))
tf = tb.text_frame; tf.word_wrap = True
para(tf, 'M.Tech Project — Phase 1 Review (Defense Format)', size=30, color=BLUE, bold=True, align=PP_ALIGN.CENTER, first=True)
para(tf, 'RL Post-Training of LLMs: Signal Starvation and Stack-Conditioned GRPO', size=21, color=INK, bold=True, align=PP_ALIGN.CENTER)
para(tf, 'Title decoded: GRPO learns from reward contrast inside groups of completions; when a group is all-correct or all-wrong it teaches nothing (signal starvation, measured by the Zero-Variance Fraction), and what a run "shows" depends on the software stack it ran on (stack-conditioned).', size=14, color=MUTED, align=PP_ALIGN.CENTER)
para(tf, 'Arvind C R (Arvind Chitra Rajasekaran)  ·  SRN: PES2PGE24DS140', size=15, color=MUTED, align=PP_ALIGN.CENTER)
para(tf, 'Guide: Ramesh Prakash Guledgudd  ·  Dept. of CSE, PES University  ·  M.Tech Data Science & AI', size=14, color=MUTED, align=PP_ALIGN.CENTER)

# ------------------------------------------------- 2 Base paper & understanding
s2 = slide('Base Paper & My Understanding of It', 2, notes=
  "(2 min) Base paper: GRPO from DeepSeekMath (Shao et al., 2024), the algorithm behind DeepSeek-R1. "
  "Explain the mechanism from first principles: no critic; sample G completions per prompt; each advantage "
  "is its reward minus the group mean. My understanding goes one step further than the paper: that subtraction "
  "has a structural blind spot — identical rewards zero out every advantage. Sell it with the basketball "
  "analogy at the bottom: drills, shots, a coach who only compares shots within a drill; all-makes drills "
  "teach nothing precisely when the scoreboard looks perfect - that is the whole thesis in one image. "
  "Secondary anchor: Dr.GRPO (Liu et al., 2025) which critiques GRPO's normalisation terms — Result 4 tests it. If asked about std-normalisation: base GRPO also divides by the group std; with identical rewards the numerator is zero either way, so the blind spot is identical — and the panels also test the normalisation-free Dr.GRPO form directly.")
_S2_BULLETS = [
    ('Base paper — GRPO, from "DeepSeekMath" (Shao et al., 2024; basis of DeepSeek-R1): replaces PPO\'s learned critic with a group-relative baseline — sample G completions per prompt, advantage = own reward − group mean.', True),
    ('Why it matters: critic-free means cheap and stable at LLM scale with verifiable (binary) rewards; it is now the default RL post-training family. The study runs GRPO against PPO, REINFORCE, DPO, GSPO, Dr.GRPO and audited variant labels (slide 5).', False),
    ('My understanding — the structural blind spot: if all G completions earn the SAME reward (all-correct or all-wrong), every centred advantage is exactly zero: the group consumes compute but contributes zero gradient. The reward curve cannot show this; it can read "success" precisely while learning has stopped.', True),
    ('Secondary anchor — Dr.GRPO (Liu et al., 2025): claims GRPO\'s per-length/std normalisation biases updates toward verbosity. I test this claim under controlled conditions (Result 4).', False),
    ('Thesis position: measure the blind spot (Zero-Variance Fraction, ZVF), calibrate it, budget it, and show what it changes in practice.', True),
]
bullets(s2, _S2_BULLETS, width=7.5, size=13)

# mechanism diagram (right panel): contrast group vs all-same group
GREEN = RGBColor(0x2E, 0x7D, 0x32); RED = RGBColor(0xB3, 0x2D, 0x2D); GRAY = RGBColor(0xB0, 0xB7, 0xC0)
dx = 8.55
tb = s2.shapes.add_textbox(Inches(dx), Inches(1.15), Inches(4.3), Inches(0.4))
para(tb.text_frame, 'one prompt → G=4 completions, reward r ∈ {0,1}', size=11, color=MUTED, align=PP_ALIGN.CENTER, first=True)
tb = s2.shapes.add_textbox(Inches(dx), Inches(1.62), Inches(4.3), Inches(0.35))
para(tb.text_frame, 'group with contrast', size=11, color=INK, bold=True, first=True)
for i, r in enumerate([1, 0, 1, 0]):
    box(s2, str(r), dx + i*0.62, 2.0, 0.5, 0.5, fill=BLUE if r else GRAY, size=13)
arrow(s2, dx+2.6, 2.25, dx+3.15, 2.25)
box(s2, 'A ≠ 0\ngradient flows', dx+3.2, 1.95, 1.15, 0.62, fill=GREEN, size=9)
tb = s2.shapes.add_textbox(Inches(dx), Inches(2.85), Inches(4.3), Inches(0.35))
para(tb.text_frame, 'all-correct group (the blind spot)', size=11, color=INK, bold=True, first=True)
for i, r in enumerate([1, 1, 1, 1]):
    box(s2, str(r), dx + i*0.62, 3.22, 0.5, 0.5, fill=BLUE, size=13)
arrow(s2, dx+2.6, 3.47, dx+3.15, 3.47)
box(s2, 'A = 0,0,0,0\nzero gradient', dx+3.2, 3.17, 1.15, 0.62, fill=RED, size=9)
tb = s2.shapes.add_textbox(Inches(dx), Inches(4.0), Inches(4.35), Inches(1.1))
para(tb.text_frame, 'A_i = r_i − mean(r) (the original also divides by the group std — the zero/nonzero structure is identical). Identical rewards zero every advantage; the reward curve still reads 1.0 while training has silently stopped.', size=10.5, color=MUTED, first=True)

# the basketball analogy (lay explanation of the blind spot)
ab = s2.shapes.add_textbox(Inches(0.8), Inches(5.35), Inches(12.0), Inches(1.6))
tfa = ab.text_frame; tfa.word_wrap = True
para(tfa, 'Basketball analogy:', size=13, color=GOLD, bold=True, first=True)
para(tfa, 'Each prompt is a shooting drill; the G completions are G shots at the same basket. GRPO has no absolute standard (no critic) — the coach grades each shot only against the others in that drill. Mixed makes and misses → he can tell what worked. All makes or all misses → nothing to compare, the drill taught zero. As the player improves, more drills end all-makes: practice silently stops teaching exactly when the scoreboard looks best. ZVF counts the dead drills; more shots per drill (bigger G) keeps the contrast alive longer.', size=12.5, color=INK)

# ---------------------------------------------------------- 3 Problem & RQs
bullets(slide('Problem & Research Questions', 3, notes=
  "(45 s) Compress. Two facts motivate everything: starvation is invisible in reward curves, and "
  "published comparisons are stack-conditioned. Then read the four RQs quickly."), [
    ('Signal starvation is silent: mean reward reads "success" exactly when the all-correct wall starves training. You need a second coordinate.', True),
    ('The reproducibility anchor: the SAME RL training, run by different people in different environments, can produce wildly different results — measured here as a 17× final-reward span from one undisclosed backend+checkpoint swap. Unless the details of what actually ran are reported, comparisons are meaningless.', True),
    ('RQ1 same-stack control · RQ2 ZVF as practical diagnostic · RQ3 group size G as the starvation dial · RQ4 do training gains survive held-out evaluation?', False),
])

# ---------------------------------------------------------- 4 Architecture
s4 = slide('Overall Architecture', 4, notes=
  "(1.25 min) Walk the four layers left to right: training on the managed Tinker API (LoRA, closed loss "
  "kernel — an audit constraint I exploit deliberately); evaluation on three vLLM backends so no single "
  "backend's quirks own the numbers; telemetry: per-step ZVF/GU next to reward, mirrored to W&B; and the "
  "audit layer: run manifests, checkpoint/resume, the runs-audit workbook. Everything downstream cites this.")
LY = 1.35; LH = 1.02; LW = 6.4; LX = 0.9
box(s4, 'TRAINING — Tinker managed API\nLoRA rank-4 GRPO / Dr.GRPO fleets · full-state checkpoint · kill-and-resume (verified live)', LX, LY, LW, LH, size=11)
box(s4, 'TELEMETRY — per-step (reward, ZVF, GU)\nlogged beside every reward point · parser v2 with false-positive audit', LX, LY+1.22, LW, LH, size=11)
box(s4, 'EVALUATION — vLLM pass@k harness\nModal · Lightning AI · Colab — three independent backends · clustered bootstrap CIs', LX, LY+2.44, LW, LH, size=11)
box(s4, 'AUDIT TRAIL — 983 Tinker runs + external Modal / Lightning AI / Colab runs\n(26 pass@k evals · 4 H100 cross-library baselines · Colab experiment packs)\n19 claim-critical runs linked: W&B page + checkpoint + result JSON', LX, LY+3.66, LW, LH+0.12, size=9.5)
for k in range(3):
    arrow(s4, LX+LW/2, LY+LH+1.22*k, LX+LW/2, LY+1.22*(k+1))
RX = 8.1; RW = 4.6
box(s4, 'zvf-triage (Apache-2.0, 82 tests)\nregime classifier · adaptive-G · dead-prompt drop · auto-stop\nveRL / OpenRLHF / NeMo-RL adapters', RX, LY+0.55, RW, 1.25, fill=RGBColor(0x2E,0x7D,0x32), size=10.5)
box(s4, 'MIRRORS & ARTIFACTS\nW&B: 1,034 runs / 17 projects\nHuggingFace: 49 adapter repos', RX, LY+2.35, RW, 1.15, fill=MUTED, size=10.5)
arrow(s4, LX+LW, LY+1.22+LH/2, RX, LY+0.55+0.62)
arrow(s4, LX+LW, LY+2.44+LH/2, RX, LY+2.35+0.57)
tb = s4.shapes.add_textbox(Inches(0.9), Inches(6.35), Inches(11.9), Inches(0.7))
para(tb.text_frame, 'Design rule: no number enters a paper without a path back through this stack — run → telemetry → artifact → audit row.', size=12, color=MUTED, first=True)

# ------------------------------------------------- 5 Algorithms & methods run
s = slide('Not Only GRPO — Algorithms & Methods Run, and Why', 5, notes=
  "(45 s) Fast table walk — lead with the WHY column: each method answers a specific question the thesis "
  "needed closed. PPO answers 'was the critic needed?'; GSPO isolates one knob on the same stack; Dr.GRPO "
  "tests the base paper's strongest critique; DPO is the no-rollouts counterfactual (no groups, so no ZVF "
  "by construction); REINFORCE and SFT are floors; the DAPO audit and the AERO/AREAL/GIFT traces exist to "
  "measure how much a method LABEL underdetermines what actually ran. Every row has runnable artifacts.")
table(s, [
    ['algorithm / method', 'where it ran', 'why it was run (question it answers)'],
    ['GRPO (base paper)', 'Tinker fleets (Qwen3-8B); TRL on Modal H100; open trainers', 'The subject: does its group-relative update starve, when, and how do you see it? (Claims 1-2)'],
    ['Dr.GRPO', 'Tinker, six-arm uncapped panel, 3 seeds/loss', "Test the base paper's strongest critique: does the length-bias fix change length/ZVF at our scale? (No footprint - Result 4)"],
    ['GSPO', 'Tinker, same-stack head-to-head vs GRPO (G=8)', 'Isolate ONE knob (token- vs sequence-level IS ratio) with everything else fixed - the same-stack control discipline (RQ1)'],
    ['PPO', 'SB3 / CleanRL / Tianshou (Modal H100) + bench arm (HF: tinker-rl-bench-ppo_*)', 'Was dropping the critic actually free? Value-based baseline + cross-library reproducibility anchor (RQ1)'],
    ['REINFORCE', 'Colab baselines notebook', 'Simplest policy gradient: a floor showing gains are not just any-PG artifacts'],
    ['DPO / IDPO', 'TRL via unified script generator (platform_local)', 'The counterfactual family: preference optimisation without rollouts has no groups - no ZVF by construction'],
    ['SFT + off-policy distillation', 'capstone tool-call LoRAs; Llama-3.2-1B distillation probe', 'Non-RL control: how much of the gain needs RL at all?'],
    ['"DAPO" (label audit)', 'open trainer w/ dynamic sampling vs closed-stack surrogate', 'How much does a method NAME underdetermine the executed update? Measured: ZVF 0.00 vs 0.55-0.58 under one label'],
    ['AERO / AREAL / GIFT (traces); audit pilot: GRPO / Dr.GRPO / DAPO / adaptive-G', 'open-stack method traces; 4-arm single-stack audit pilot (T4)', 'Can telemetry alone distinguish method variants (P8 detector)? Is collapse universal across them (P2)? Does the survival protocol work end-to-end (pilot)?'],
], col_widths=[2.5, 3.9, 5.1], size=11, top=1.3)

# ------------------------------------------- 6 What I implemented (attribution)
bullets(slide('What I Implemented (Sem-4 Solo, on the Sem-3 Foundation)', 6, notes=
  "(1.5 min) Attribution first: Sem 3 was the group capstone — the multi-framework bench and survey, frozen at "
  "tag capstone-final-2026-04-25. Everything on this slide is Sem-4 solo work. Be specific: these are files "
  "and packages he can open, not concepts. Datasets first if he asks what you trained ON: GSM8K train "
  "split with binary boxed-answer rewards is the finetuning task; synthetic arithmetic for the small panels; "
  "xLAM-60k was the capstone SFT corpus. MATH-500/MBPP/HumanEval are evaluation-only - never trained on. "
  "Sem-3 credit line if asked: 79 logged runs across 7 libraries, submitted to NeurIPS 2026 Main Track as "
  "'A Unified Benchmark for RL Post-Training of Language Models' (group work, 8 authors); Sem 4 is solo."), [
    ('Datasets used for finetuning — RL: GSM8K train split (openai/gsm8k "main"; 512-prompt train pool, binary boxed-answer reward; training samples G completions/prompt per step, and the same pool is offline-characterised at 32 rollouts/prompt for the estimation studies) + synthetic arithmetic (easy/medium) for the 0.5B-1.5B panels. SFT (capstone): Salesforce xLAM-Function-Calling-60k (tool-call LoRAs, Qwen 0.5B-7B).', True),
    ('Evaluation-only — never trained on: GSM8K test (1,319 problems, disjoint from the reward environment), MATH-500, MBPP, HumanEval. Held-out discipline is an explicit reporting-standard item.', True),
    ('Inherited (Sem 3, Group 6): TinkerRL-Bench — 79 logged GRPO-style runs across 7 RL libraries (0.6B–671B), submitted to NeurIPS 2026 Main Track as "A Unified Benchmark for RL Post-Training of Language Models"; plus literature survey and baseline GRPO runs. Frozen at tag capstone-final-2026-04-25.', True),
    ('ZVF measurement stack: per-step ZVF/GU telemetry in the trainer, calibrated confidence intervals (Wilson), waiting-time reliability budget, stratified batch analysis.', True),
    ('Experiment infrastructure: matched-budget runner with --resume (state + optimiser + RNG fast-forward), per-(step,prompt) seeding, W&B resume; loss-form panel runner (GRPO vs Dr.GRPO).', True),
    ('zvf-triage: packaged library (callback, controller, regime classifier, framework adapters, 82-test suite) — publication to PyPI staged.', True),
    ('Standards & tooling: MIN-REPORT-RL 8-item reporting standard, GRPO-Registry (machine-readable stack catalog), stackdiff flip-risk grader, run-audit workbook.', False),
    ('Theory: T1 estimator calibration, T2 reliability budget, T3 optimal-G analysis — plus two corrections found by external adversarial review, adopted and reported openly.', False),
], size=13, top=1.15)

# ------------------------------------------------------- 6 Result 1: Claim 1
s = slide('Result 1 — ZVF Sees What the Reward Curve Cannot (Claim 1)', 7, notes=
  "(2 min) THE core result. Walk the table: late in training the G=2 arms read reward ~1.0 — by the reward "
  "axis, perfect. ZVF says 75-100% of groups are all-correct: zero gradient. Same budget, G=16 arms are "
  "mid-learning with ZVF under 0.25 and signal intact. Read as a pair, (reward, ZVF) separates 'policy is good' "
  "from 'training is still moving'. ZVF alone aliases mastery with incapacity — always read the pair.")
table(s, [
    ['', 'late-run mean reward', 'late-run ZVF', 'gradient signal'],
    ['G=2 × 160 steps', '≈ 0.9–1.0 (pool mastered)', '0.75–1.0 (all-correct wall)', 'effectively zero'],
    ['G=16 × 20 steps', '≈ 0.3–0.5 (mid-learning)', '0.00–0.25', 'sustained'],
], col_widths=[2.6, 3.4, 3.2, 2.3])
bullets(s, [
    ('Same rollout budget (2,560/arm), seeds 123/456: reward alone declares G=2 the winner; the (reward, ZVF) pair shows its lead ended in zero-gradient compute.', True),
    ('ZVF is a diagnostic, not a predictor: ZVF rose AFTER the reward plateau in every measured collapse — a cheap alarm, not a cause.', False),
    ('Population form: the U-shaped kernel — starvation at both walls (too hard / mastered); larger G narrows both.', False),
], top=3.3, width=7.2, size=13)
s.shapes.add_picture('outputs/deck_assets/ushape.png', Inches(8.3), Inches(3.3), width=Inches(4.5))

# ------------------------------------------------------- 7 Result 2: Claim 2
s8 = slide('Result 2 — Group Size Is a Schedule Variable (Claim 2)', 8, notes=
  "(2 min) The decisive experiment design point: hold the ROLLOUT BUDGET fixed, not the step count. "
  "Small G converts the budget into more optimiser steps early — then exhausts its own signal as accuracy "
  "rises (the p->1 wall of the kernel). Large G pays for contrast it doesn't need early and retains signal late. "
  "So group size controls WHICH END of training starves — a schedule question, not a constant. The naive "
  "static sweep is confounded by what the budget is held in — I show both views. ANTICIPATED ATTACK: "
  "'G=16 only avoids the wall because it has not learned the task yet.' Answer: exactly — and that is the "
  "claim. At a FIXED budget, G chooses which end starves: G=2 spends its final steps on dead groups while "
  "reading reward 1.0; G=16 keeps every update signal-bearing. No winner is declared; the pair tells you "
  "where the budget went.")
bullets(s8, [
    ('Design: matched budget of 2,560 rollouts per arm — G=2×160 steps vs G=16×20 steps (batch 8, 512-token completions, LoRA rank 4, seeds 123/456).', True),
    ('Finding: G=2 races to reward ≈1.0 on the sampled pool, then terminates inside the all-correct zero-variance wall; G=16 ends mid-learning with ZVF ≤ 0.25 and signal intact.', True),
    ('Interpretation: small G buys early steps and starves the endgame; large G holds signal throughout. Group size selects WHICH END of training starves — a schedule variable.', True),
    ('Honesty check: static fixed-step sweeps are confounded by what the budget is held in; controller efficacy is NOT claimed (pre-registered test is future work).', False),
], width=7.2, size=13)
s8.shapes.add_picture('outputs/deck_assets/budget_traj.png', Inches(8.25), Inches(1.6), width=Inches(4.6))

# ------------------------------------------------- 8 Result 3: theory calibrated
bullets(slide('Result 3 — The Estimator Is Calibrated, the Budget Is Exact', 9, notes=
  "(1.5 min) Three theory results, each validated on real 512-prompt pools. T1: Wilson interval covers 0.95-0.98 "
  "in every tested setting — report which ZVF the interval covers under curriculum ordering. T2: geometric "
  "waiting-time budget N = G ln(delta)/ln(ZVF) matched observed quantiles at ratio 1.00 in all six difficulty "
  "strata — hardest stratum needs 160 rollouts for a 90%-guaranteed informative group. T3 is the honest one: "
  "our signal-per-rollout objective turns out to have a UNIVERSAL argmax G* in {2,3} for every prior — an "
  "algebraic identity, found by external review of our own theory. We report it as a negative result. ANTICIPATED ATTACK: Wilson assumes iid — answer: under curriculum ordering the interval stays calibrated for the LOCAL stage-level ZVF (0.944), and stratified batch composition restores global validity (0.996); the requirement is labelling WHICH estimand the interval covers."), [
    ('T1 (calibration): ZVF_t is an unbiased binomial-proportion estimator; Wilson CI covers 0.95–0.98 in every tested setting (Wald marginal at m=32). Curriculum ordering is an estimand-labelling requirement, not an invalidation.', True),
    ('T2 (reliability budget): rollouts-to-next-informative-group is geometric; N(ZVF)=G⌈ln δ/ln ZVF⌉ matched observed quantiles at ratio 1.00 across all six difficulty strata (hardest: ZVF=0.886 ⇒ 160 rollouts) — both quantiles are integer counts of G-rollout groups, so exact agreement is attainable, not a suspicious continuous fit. A budget, NOT an impossibility bound.', True),
    ('T3 (honest negative): the signal-per-rollout objective satisfies J(2)=J(3) for EVERY difficulty prior — its argmax is universally {2,3}, so it cannot yield a data-adaptive G — per-rollout accounting ALWAYS favours tiny groups, which is exactly why Result 2 treats G as a schedule, not a static optimum. Found by adversarial review; reported as a result, not buried.', True),
    ('Two earlier statement errors (a quantifier confusion in T2; a GU sign slip) were caught by external review, corrected, re-validated — and are documented in the thesis as part of the method.', False),
])

# --------------------------------------- 9 Result 4: loss panel + the incident
bullets(slide('Result 4 — GRPO vs Dr.GRPO: No Footprint at This Scale (+ the Incident)', 10, notes=
  "(2 min) Tests the base-paper critique directly. Six uncapped arms (1,024 tokens), 3 seeds per loss: "
  "completion lengths SHRINK 3.8-12.2% in all six arms — no verbosity trap at this scale — and no late-ZVF "
  "separation. Then own the incident proudly: the first panel was invalid because a documented --loss flag "
  "was never wired to the loss; no output-level trace revealed it — only reading the runner did. We "
  "invalidated loudly, preserved artifacts under .invalid names, reran same-day, and one conclusion REVERSED. "
  "That incident is now a case study and the seed of the reporting standard.", ), [
    ('Six-arm uncapped panel (Qwen3-8B, 1,024-token cap, 3 seeds/loss): GRPO lengths 1004→905, 981→944, 996→900; Dr.GRPO 999→931, 972→902, 1000→878 — lengths SHRINK 3.8–12.2% in every arm; no length inflation, no late-ZVF separation between losses.', True),
    ('Reading: at this scale the loss-form choice has no observable footprint on length or ZVF — evidence about reporting, not superiority. Comparisons between these losses measure stack noise unless controlled far more tightly than the labels suggest.', False),
    ('The incident: the first panel ran with a documented --loss drgrpo flag that was never wired in — both "arms" silently trained identical GRPO. Caught only by reading OUR client-side runner script (the managed loss kernel stays closed — the flag died in our code before reaching it); artifacts preserved under .invalid_actually_grpo names; corrected rerun REVERSED one conclusion.', True),
    ('Response became protocol: invalidate loudly → preserve → rerun → record. This failure is a first-class result feeding the reproducibility standard (Result 5).', False),
])

# ------------------------------------------- 10 Result 5: reproducibility results
bullets(slide('Result 5 — Setting the Standard: What Must Be Reported When Training', 11, notes=
  "(1.5 min) Three measured instances, each a lever that flipped a result: 17x reward span from an undisclosed "
  "backend swap that also bundled a checkpoint change; the same 'DAPO' label yielding ZVF 0.00 on an open "
  "trainer vs 0.55-0.58 on a closed stack; and reward micro-jitter below the verifier's resolution collapsing "
  "batch ZVF 0.158 to 0. Every item of the 8-item MIN-REPORT-RL standard exists because one of these levers "
  "moved a result in OUR OWN data — not from taste. Land the framing: this project SETS THE STANDARD for what to report when training - walk the 8 items fast, then say every one earned its place by flipping a result here."), [
    ('Anchor: same algorithm, same task, different environment or different hands → wildly different results. The fix is not more compute — it is reporting the details of what actually ran. Three measured instances from our own data:', True),
    ('Backend swap (undisclosed, bundled a base-checkpoint change): final training reward moved across a 17× span — 85.6% vs 5.0% — under the same label.', True),
    ('Same "DAPO" label: mean ZVF 0.00 on an open trainer with true dynamic sampling vs 0.55–0.58 on a closed stack running an asymmetric-clip surrogate.', True),
    ('Reward-parser sensitivity: micro-jitter ε~U(0,1e-4) below verifier resolution collapses batch ZVF 0.158 → 0.000 — reported ZVF must name its verifier.', False),
    ('The contribution: MIN-REPORT-RL sets the standard for what details every RL training report must include — plus GRPO-Registry (machine-readable stack catalog, 20 seed entries) and stackdiff (pairwise flip-risk verdicts R0–R5) to enforce it mechanically.', True),
    ('The 8 items: ① loss form (ratio/clip/mask/normalisation) ② reference policy & KL handling ③ sampler/backend/precision incl. base-checkpoint identity ④ per-step ZVF/GU trajectory ⑤ group-size schedule ⑥ held-out split disjoint from the reward environment ⑦ decontamination + parser-robustness probe ⑧ held-out pass@k curves, not just pass@1. Every item earned its place by flipping a result in our own corpus.', False),
], size=13)

# ------------------------------------------- 11 Result 6: held-out evaluation
bullets(slide('Result 6 — Held-Out Evaluation: Gains, Transfer, and an Honest Boundary', 12, notes=
  "(1.25 min) RQ4. Base Qwen3-8B on GSM8K: pass@1 30.4% but pass@32 91% — the base model already solves "
  "almost everything at k=32, so GSM8K alone cannot demonstrate capability expansion; that scope discipline "
  "is itself a finding. Post-RL adapters: zero forgetting on MBPP and a +1.5-2.5 point pass@32 frontier "
  "improvement, within noise at single-seed — exactly why the standard mandates pass@k curves with CIs. "
  "MATH-500 partial: GSM8K-trained gains do NOT replicate on hard math — distribution sharpening, "
  "not capability expansion."), [
    ('Baseline capability (deterministic first-200 test slice, n=32, clustered bootstrap): pass@1 30.4% [27.5, 33.1] but pass@32 91.0% — corroborated on the 512-prompt train pool (30.3%). GSM8K is nearly saturated at k=32; ~9 points of headroom bounds what training can claim here.', True),
    ('Transfer: post-RL adapters show zero forgetting and mild positive transfer on MBPP (pass@32 within noise of or above base for all G); pass@1-only reporting would misread the G=2 arm as a regression (all pass@1 deltas −0.5 to +1.9pp, within single-seed noise).', False),
    ('Hard-task boundary (MATH-500, partial): GSM8K-trained frontier gains do not carry — consistent with distribution sharpening rather than capability expansion. Stated as a non-claim in the thesis.', True),
    ('Cross-scale observations (the 79-run cross-library corpus, 7 libraries, 0.6B–671B — frontier entries are single-seed LoRA case studies via the managed API): scale does not uniformly reduce ZVF; starvation is (difficulty × G × phase) geometry. Observations, never claims (scope slide).', False),
])

# ---------------------------------------------------- 12 Scale & evidence trail
s13 = slide('Implementation Scale & Evidence Trail (audited 12 Jul)', 13, notes=
  "(45 s) Fast slide. 983 runs on the Tinker account — audited and classified this morning; the thesis "
  "claims rest on 19 identified claim-critical runs, each linked to W&B and its artifact. If asked about "
  "any number: the workbook key_runs sheet has the run id, checkpoint, W&B link, and result JSON.")
_S13_BULLETS = [
    ('983 Tinker training runs enumerated via REST API — 26 base models, 0.6B → 1T; all 65 corrupted runs predate June 8; every thesis-supporting run is clean and checkpointed.', True),
    ('External backends: Modal / Lightning AI / Colab pass@k panels (26 runs) + 4 cross-library H100 baselines; 1,034 W&B runs across 17 projects; 49 HuggingFace artifact repos.', False),
    ('zvf-triage: 82/82 tests green; wheel + sdist built, twine-checked; PyPI publication staged.', False),
    ('Traceability: 19 claim-critical runs highlighted in the audit workbook, each with W&B link + checkpoint + result artifact — embedded below (double-click to open).', True),
]
bullets(s13, _S13_BULLETS, size=14)

# embed the audit workbook as a real OLE object (opens in Excel on double-click)
from pptx.enum.shapes import PROG_ID
s13.shapes.add_ole_object(
    'outputs/tinker_runs_audit_2026-07-12.xlsx', PROG_ID.XLSX,
    left=Inches(0.95), top=Inches(5.15), icon_width=Inches(0.85), icon_height=Inches(0.95))
tb = s13.shapes.add_textbox(Inches(1.95), Inches(5.3), Inches(10.5), Inches(0.8))
para(tb.text_frame, 'tinker_runs_audit_2026-07-12.xlsx — embedded copy (983 runs · key_runs · external_runs · wandb_runs · hf_artifacts · insights). Canonical file: outputs/ in the repo.', size=12, color=MUTED, first=True)

# ------------------------------------------------------------------ 13 Demo
s14 = slide('Demo (Live)', 14, notes=
  "(1 min + live demo) Run the one-command offline demo FIRST — it cannot fail on the network. "
  "./submission/demo/demo.sh: mechanism fixture (4 groups, ZVF=0.5, GU=0.5), recorded artifact (80 rewards, "
  "mean 0.6875, ZVF 0.30), SHA-256 integrity check, HTML dashboard. Then if time and connectivity allow: "
  "the W&B zvf-training panel with live E-R2b curves, and the audit workbook key_runs sheet. "
  "Fallback: the 86-second recorded walkthrough. Verified PASS this morning.")
_S14_BULLETS = [
    ('One command, fully offline: ./submission/demo/demo.sh — mechanism fixture (4 groups → ZVF=0.500, GU=0.500), recorded artifact check (80 rewards, mean 0.6875, ZVF 0.3000), SHA-256 integrity, JSON + HTML dashboard. Status: PASS (re-verified today).', True),
    ('Live artifact tour: W&B zvf-training — the (reward, ZVF) pair diverging on the real E-R2b arms; audit workbook key_runs sheet — every claim-critical run traceable in two clicks.', False),
    ('zvf-triage quickstart: pip-installable package, examples/quickstart.py — the diagnostic as a reusable library, not a one-off script.', False),
    ('Fallback if connectivity fails: the 86-second recorded walkthrough — embedded on this slide (right), plays offline.', False),
]
bullets(s14, _S14_BULLETS, size=14, width=7.4)
if os.path.exists(FALLBACK):
    # 1728x1122 -> keep aspect (h = w * 1122/1728)
    s14.shapes.add_movie(FALLBACK, Inches(8.45), Inches(1.2), Inches(4.3), Inches(2.79),
                         poster_frame_image=FPOSTER if os.path.exists(FPOSTER) else None,
                         mime_type='video/mp4')
    tbv = s14.shapes.add_textbox(Inches(8.45), Inches(4.02), Inches(4.3), Inches(0.4))
    para(tbv.text_frame, 'embedded fallback: 86 s recorded walkthrough (click to play)', size=10.5, color=MUTED, align=PP_ALIGN.CENTER, first=True)

# click-to-open evidence links (live during Q&A)
tb = s14.shapes.add_textbox(Inches(0.85), Inches(4.55), Inches(11.8), Inches(2.3))
tf = tb.text_frame; tf.word_wrap = True
para(tf, 'Click-to-open evidence (live):', size=14, bold=True, first=True)
linkline(tf, 'W&B: E-R2b G=2 arm — reward hits ~1.0 while ZVF climbs into the wall (er2b_g2_s123)',
         'https://wandb.ai/arvindcr4-pes-university/zvf-training/runs/pob7nd05')
linkline(tf, 'W&B: E-R2b G=16 arm — mid-learning, ZVF low, signal intact (er2b_g16_s123)',
         'https://wandb.ai/arvindcr4-pes-university/zvf-training/runs/tiicy3km')
linkline(tf, 'W&B: the invalidated P4 arm, preserved (p4uncap_drgrpo_s42.invalid_actually_grpo)',
         'https://wandb.ai/arvindcr4-pes-university/zvf-training/runs/kmqjbhwn')
linkline(tf, 'W&B: full zvf-training project (45 runs)',
         'https://wandb.ai/arvindcr4-pes-university/zvf-training')
linkline(tf, 'HuggingFace: published bench adapters (e.g. tinker-rl-bench-ppo_gsm8k_Qwen3-8B_s42)',
         'https://huggingface.co/arvindcr4/tinker-rl-bench-ppo_gsm8k_Qwen3-8B_s42')
linkline(tf, 'GitHub: arvindcr4/tinker-rl-lab (code, artifacts, audit workbook)',
         'https://github.com/arvindcr4/tinker-rl-lab')

# ------------------------------------------------- 14 Limitations & close
bullets(slide('Scope, Limitations & Roadmap', 15, notes=
  "(45 s) Close with the honesty that survives cross-examination: one model, one task family, one managed "
  "API, 1-3 seeds — the claims are stated at the stack level and nowhere above it. Roadmap is gated: "
  "diagnostic paper publishable now; controller paper gated on a pre-registered compute-matched win; "
  "survival audit gated on an open stack. Then stop talking and invite questions."), [
    ('Declared scope: Qwen3-8B, GSM8K-family binary rewards, Tinker managed API (LoRA rank 4), 1–3 seeds per result — claims are stated at the stack level and nowhere above it.', True),
    ('No causal or predictive power is claimed for ZVF beyond diagnosis; controller efficacy is explicitly gated on a pre-registered, compute-matched comparison (≥3 seeds, held-out metrics).', False),
    ('Roadmap (gated): bounded diagnostic paper (Claims 1–2) → controller paper (needs the win) → survival audit (needs an open stack). Thesis consolidates all 17 working documents.', False),
    ('The transferable lesson: at this scale the binding constraint is not compute or novelty — it is certainty about what actually ran.', True),
])

s = slide(notes="Thank the panel. Repo and email on screen. Offer the audit workbook or any W&B run on request.")
tb = s.shapes.add_textbox(Inches(1.0), Inches(2.7), Inches(11.3), Inches(2.0))
tf = tb.text_frame
para(tf, 'Thank you — Questions', size=30, color=BLUE, bold=True, align=PP_ALIGN.CENTER, first=True)
pcl = para(tf, '', size=15, color=MUTED, align=PP_ALIGN.CENTER)
r = pcl.add_run(); r.text = 'github.com/arvindcr4/tinker-rl-lab'
r.font.size = Pt(15); r.font.color.rgb = RGBColor(0x0B,0x5C,0xAB); r.font.underline = True; r.font.name='Calibri'
r.hyperlink.address = 'https://github.com/arvindcr4/tinker-rl-lab'
r = pcl.add_run(); r.text = '   ·   arvindcr4@gmail.com'
r.font.size = Pt(15); r.font.color.rgb = MUTED; r.font.name='Calibri'

prs.save('outputs/PESU_MTech_Phase1_Session1_Review_ArvindCR.pptx')
print('slides:', len(prs.slides._sldIdLst))
