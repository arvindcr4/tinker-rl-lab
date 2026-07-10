#!/usr/bin/env python3
"""
Build the M.Tech final thesis review deck for the Tinker RL Lab project.

Every number/claim on the slides is sourced from repository files; see the
inline SRC comments and the accompanying README.md for the file -> claim map.

Usage:
    python build_review_pptx.py          # writes Arvind_MTech_Thesis_Review.pptx
    python build_review_pptx.py --check   # also round-trip loads the result

Design: clean academic. White background, PES-blue (#1F4E79) titles, Calibri,
slide numbers, embedded figure PNGs with captions. If a figure PNG is missing,
a text panel is drawn instead (never a broken image).

Repository root is resolved relative to this file so the script is portable.
"""
from __future__ import annotations

import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

try:
    from PIL import Image
    _HAS_PIL = True
except Exception:  # pragma: no cover
    _HAS_PIL = False

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
REPO_ROOT = os.path.abspath(os.path.join(HERE, "..", "..", ".."))
OUT_PATH = os.path.join(HERE, "Arvind_MTech_Thesis_Review.pptx")


def fig(*parts: str) -> str:
    return os.path.join(REPO_ROOT, *parts)


# ---------------------------------------------------------------------------
# Palette / typography
# ---------------------------------------------------------------------------
PES_BLUE = RGBColor(0x1F, 0x4E, 0x79)     # dark navy title
ACCENT = RGBColor(0x2E, 0x74, 0xB5)       # lighter blue rule / accents
INK = RGBColor(0x22, 0x22, 0x22)          # body text
MUTED = RGBColor(0x5A, 0x5A, 0x5A)        # captions / footer
LIGHT_BG = RGBColor(0xF2, 0xF5, 0xF9)     # panel fill
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xD0, 0xD9, 0xE4)
FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_W = SLIDE_W - 2 * MARGIN

EMBEDDED_FIGURES: list[str] = []   # populated as figures are placed
MISSING_FIGURES: list[str] = []


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------
def _set_bg(slide, color=WHITE):
    fillel = slide.background.fill
    fillel.solid()
    fillel.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    return tb, tf


def _style_run(run, size, color=INK, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def _para(tf, text, size=16, color=INK, bold=False, italic=False,
          align=PP_ALIGN.LEFT, space_after=6, space_before=0, level=0,
          bullet=False, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.level = level
    if space_after is not None:
        p.space_after = Pt(space_after)
    if space_before is not None:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    _style_run(run, size, color=color, bold=bold, italic=italic)
    if bullet:
        _apply_bullet(p)
    else:
        _no_bullet(p)
    return p


def _no_bullet(p):
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    buNone = pPr.makeelement(qn("a:buNone"), {})
    pPr.append(buNone)


def _apply_bullet(p, char="•", color=ACCENT):
    pPr = p._p.get_or_add_pPr()
    pPr.set("indent", "-182880")
    pPr.set("marL", "182880")
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    buClr = pPr.makeelement(qn("a:buClr"), {})
    srgb = pPr.makeelement(qn("a:srgbClr"), {"val": "2E74B5"})
    buClr.append(srgb)
    pPr.append(buClr)
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    pPr.append(buFont)
    buChar = pPr.makeelement(qn("a:buChar"), {"char": char})
    pPr.append(buChar)


def _title_band(slide, title, kicker=None):
    """Title text + accent rule at the top of a content slide."""
    tb, tf = _add_textbox(slide, MARGIN, Inches(0.34), CONTENT_W, Inches(0.9))
    if kicker:
        _para(tf, kicker.upper(), size=12, color=ACCENT, bold=True,
              space_after=2, first=True)
        _para(tf, title, size=27, color=PES_BLUE, bold=True, space_after=0)
    else:
        _para(tf, title, size=28, color=PES_BLUE, bold=True, first=True)
    # accent rule
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.30), Inches(2.2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    line.shadow.inherit = False


def _footer(slide, page_no):
    tb, tf = _add_textbox(slide, MARGIN, Inches(7.06), CONTENT_W, Inches(0.32),
                          anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Tinker RL Lab  —  M.Tech Final Thesis Review  •  Arvind C R"
    _style_run(r, 9, color=MUTED)
    # page number on the right
    tb2, tf2 = _add_textbox(slide, SLIDE_W - Inches(1.4), Inches(7.06),
                            Inches(0.8), Inches(0.32), anchor=MSO_ANCHOR.MIDDLE)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = str(page_no)
    _style_run(r2, 10, color=MUTED, bold=True)


def new_slide(prs, title=None, kicker=None, page_no=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, WHITE)
    if title:
        _title_band(slide, title, kicker=kicker)
    if page_no is not None:
        _footer(slide, page_no)
    return slide


def add_panel(slide, left, top, width, height, fill=LIGHT_BG, line=RULE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    # soften the corner radius
    try:
        shp.adjustments[0] = 0.04
    except Exception:
        pass
    return shp


def add_bullets(slide, left, top, width, height, items, size=16,
                anchor=MSO_ANCHOR.TOP):
    """items: list of (text, level, bold) or plain strings."""
    tb, tf = _add_textbox(slide, left, top, width, height, anchor=anchor)
    first = True
    for it in items:
        if isinstance(it, tuple):
            text, level, bold = (it + (0, False))[:3]
        else:
            text, level, bold = it, 0, False
        sz = size if level == 0 else size - 2
        color = INK if not bold else PES_BLUE
        _para(tf, text, size=sz, color=color, bold=bold, level=level,
              bullet=True, space_after=7 if level == 0 else 4, first=first)
        first = False
    return tb


def add_image_fitted(slide, path, left, top, max_w, max_h, caption=None):
    """Place image scaled to fit within (max_w, max_h), centered; add caption."""
    rel = os.path.relpath(path, REPO_ROOT)
    if not os.path.exists(path):
        MISSING_FIGURES.append(rel)
        return _text_fallback(slide, left, top, max_w, max_h,
                              f"[figure unavailable: {os.path.basename(path)}]",
                              caption)
    # figure aspect
    aspect = None
    if _HAS_PIL:
        try:
            with Image.open(path) as im:
                w, h = im.size
                aspect = w / h
        except Exception:
            aspect = None
    cap_h = Inches(0.32) if caption else Inches(0)
    avail_h = max_h - cap_h
    if aspect:
        draw_w = max_w
        draw_h = Emu(int(draw_w / aspect))
        if draw_h > avail_h:
            draw_h = avail_h
            draw_w = Emu(int(draw_h * aspect))
    else:
        draw_w, draw_h = max_w, avail_h
    img_left = left + (max_w - draw_w) // 2
    slide.shapes.add_picture(path, img_left, top, width=draw_w, height=draw_h)
    EMBEDDED_FIGURES.append(rel)
    if caption:
        cb, cf = _add_textbox(slide, left, top + draw_h + Inches(0.04),
                              max_w, cap_h)
        p = cf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = caption
        _style_run(r, 10.5, color=MUTED, italic=True)
    return True


def _text_fallback(slide, left, top, w, h, msg, caption=None):
    add_panel(slide, left, top, w, h, fill=RGBColor(0xF7, 0xF7, 0xF7))
    tb, tf = _add_textbox(slide, left + Inches(0.2), top, w - Inches(0.4), h,
                          anchor=MSO_ANCHOR.MIDDLE)
    _para(tf, msg, size=13, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
          first=True)
    if caption:
        _para(tf, caption, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    return False


def add_table(slide, left, top, width, rows, col_widths=None, header=True,
              font_size=13, row_height=Inches(0.34)):
    n_rows = len(rows)
    n_cols = len(rows[0])
    height = row_height * n_rows
    gtbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = gtbl.table
    # disable banded styling for a clean academic look
    tbl.first_row = header
    tbl.horz_banding = False
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Emu(int(width * cw / total))
    for r in range(n_rows):
        tbl.rows[r].height = row_height
        for c in range(n_cols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.01)
            cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(rows[r][c])
            is_head = header and r == 0
            _style_run(run, font_size, color=WHITE if is_head else INK,
                       bold=is_head)
            if is_head:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PES_BLUE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT_BG
    return gtbl


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------
def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    # top color band
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.28))
    band.fill.solid(); band.fill.fore_color.rgb = PES_BLUE
    band.line.fill.background(); band.shadow.inherit = False
    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.22), SLIDE_W, Inches(0.28))
    bottom.fill.solid(); bottom.fill.fore_color.rgb = PES_BLUE
    bottom.line.fill.background(); bottom.shadow.inherit = False

    tb, tf = _add_textbox(slide, Inches(0.9), Inches(1.15), Inches(11.5), Inches(0.5))
    _para(tf, "M.TECH DATA SCIENCE & AI  —  FINAL THESIS REVIEW",
          size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, first=True)

    tb, tf = _add_textbox(slide, Inches(0.7), Inches(1.75), Inches(11.9), Inches(1.9))
    _para(tf, "Tinker RL Lab", size=46, color=PES_BLUE, bold=True,
          align=PP_ALIGN.CENTER, space_after=6, first=True)
    _para(tf, "A Multi-Framework Benchmark and Study of GRPO-Style "
              "Reinforcement-Learning Post-Training of Large Language Models",
          size=21, color=INK, align=PP_ALIGN.CENTER, space_after=0)

    # rule
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.17), Inches(3.95),
                                  Inches(3.0), Pt(2.5))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background(); line.shadow.inherit = False

    tb, tf = _add_textbox(slide, Inches(1.2), Inches(4.25), Inches(10.9), Inches(2.4),
                          anchor=MSO_ANCHOR.TOP)
    _para(tf, "Arvind C R  (Arvind Chitra Rajasekaran)", size=20, color=INK,
          bold=True, align=PP_ALIGN.CENTER, space_after=2, first=True)
    _para(tf, "SRN: PES2PGE24DS140", size=15, color=MUTED, align=PP_ALIGN.CENTER,
          space_after=10)
    _para(tf, "Project Guide:  Ramesh Prakash Guledgudd", size=16, color=INK,
          align=PP_ALIGN.CENTER, space_after=2)
    _para(tf, "Department of Computer Science & Engineering", size=15,
          color=INK, align=PP_ALIGN.CENTER, space_after=2)
    _para(tf, "PES University, Bengaluru", size=15, color=INK,
          align=PP_ALIGN.CENTER, space_after=10)
    _para(tf, "July 2026", size=14, color=MUTED, align=PP_ALIGN.CENTER)


def slide_02_agenda(prs, n):
    slide = new_slide(prs, "Agenda", page_no=n)
    left_items = [
        "1.  Problem statement & motivation",
        "2.  Objectives",
        "3.  Literature context (RLHF → PPO → DPO → GRPO)",
        "4.  System architecture (6 RL frameworks)",
        "5.  Methodology: the four de-confound pillars",
        "6.  Key results (5 slides)",
    ]
    right_items = [
        "7.  P1–P8 contribution map",
        "8.  Applied study: LLM vs XGBoost fraud",
        "9.  Reproducibility & audit apparatus",
        "10. Publications & submissions",
        "11. Demo, ownership, limitations",
        "12. Conclusions & Q&A",
    ]
    add_panel(slide, MARGIN, Inches(1.65), Inches(5.95), Inches(4.9))
    add_panel(slide, Inches(6.85), Inches(1.65), Inches(5.9), Inches(4.9))
    add_bullets(slide, MARGIN + Inches(0.25), Inches(1.95), Inches(5.5),
                Inches(4.4), left_items, size=17)
    add_bullets(slide, Inches(7.1), Inches(1.95), Inches(5.5),
                Inches(4.4), right_items, size=17)


def slide_03_problem(prs, n):
    slide = new_slide(prs, "Problem Statement & Motivation", page_no=n)
    items = [
        ("GRPO (critic-free, group-relative RL) is now the default for LLM "
         "post-training — but its reported gains are easily confounded.", 0, True),
        ("A headline number hides the stack: same nominal “GRPO” gives "
         "TRL 73.4% vs Tinker 99.9% on the same task (p=0.0014).", 0, False),
        ("Most published RL-for-LLM results are single-seed, single-library, and "
         "train-set-only — not auditable or comparable.", 0, False),
        ("Does GRPO actually add capability, or surface pre-existing capability? "
         "On held-out GSM8K, GRPO adds only +1.3 pts over the base model "
         "(83.3% vs 82.0%, p=0.26).", 0, False),
        ("Need: an honest, multi-framework harness that runs identical GRPO/PPO/DPO "
         "workloads and isolates what actually drives the effect.", 0, True),
    ]
    add_bullets(slide, MARGIN, Inches(1.7), CONTENT_W, Inches(4.8), items, size=18)
    # source strip
    tb, tf = _add_textbox(slide, MARGIN, Inches(6.5), CONTENT_W, Inches(0.4))
    _para(tf, "Sources: LIMITATIONS_AND_IMPACT.md §8; reports/final/ held-out "
              "eval; experiments/experiment_summary.md", size=10.5, color=MUTED,
          italic=True, first=True)


def slide_04_objectives(prs, n):
    slide = new_slide(prs, "Objectives", page_no=n)
    items = [
        ("O1  Build a multi-framework benchmark that runs identical GRPO/PPO/DPO "
         "workloads across Tinker, SkyRL, verl, OpenRLHF, TRL and Atropos.", 0, True),
        ("O2  De-confound the GRPO effect via four controlled pillars "
         "(same-stack algorithm, ZVF, group size, held-out generalization).", 0, True),
        ("O3  Introduce diagnostics for signal starvation — the Zero-Variance "
         "Fraction (ZVF) and gradient utilization = 1−ZVF.", 0, True),
        ("O4  Characterize GRPO scaling behavior from 0.6B to ~671B parameters "
         "across five model families.", 0, True),
        ("O5  Enforce reproducibility: multi-seed protocol, 13-audit integrity "
         "suite, Docker, and copy-pasteable REPRODUCE.md.", 0, True),
        ("O6  Probe transfer of the measurement discipline to an applied domain "
         "(LLM vs XGBoost credit-card fraud).", 0, True),
    ]
    add_bullets(slide, MARGIN, Inches(1.7), CONTENT_W, Inches(5.0), items, size=18)


def slide_05_literature(prs, n):
    slide = new_slide(prs, "Literature Context", kicker="Algorithmic lineage",
                      page_no=n)
    rows = [
        ["Method", "Key idea", "Reference"],
        ["RLHF / InstructGPT", "Align LMs to human preference via a learned reward + RL",
         "Ouyang et al. 2022"],
        ["PPO", "Clipped policy-gradient with a value critic (the RLHF workhorse)",
         "Schulman et al. 2017"],
        ["DPO", "Skip the reward model: preference as a closed-form classification loss",
         "Rafailov et al. 2023"],
        ["GRPO", "Critic-free; group-relative baseline over G sampled completions",
         "Shao et al. 2024 (DeepSeekMath)"],
        ["R1-style RL", "GRPO at scale elicits emergent chain-of-thought reasoning",
         "DeepSeek-R1, 2025"],
    ]
    add_table(slide, MARGIN, Inches(1.75), CONTENT_W, rows,
              col_widths=[0.22, 0.5, 0.28], font_size=14.5,
              row_height=Inches(0.72))
    tb, tf = _add_textbox(slide, MARGIN, Inches(6.35), CONTENT_W, Inches(0.6))
    _para(tf, "This work studies the GRPO node of this lineage: critic-free RL "
              "removes the value network but introduces group-relative variance as "
              "the new failure surface (Pillars 1–2).", size=13.5, color=INK,
          italic=True, first=True)
    tb2, tf2 = _add_textbox(slide, MARGIN, Inches(6.92), CONTENT_W, Inches(0.3))
    _para(tf2, "Source: capstone-literature-survey/chapter2_foundation.tex, "
               "references.bib", size=10.5, color=MUTED, italic=True, first=True)


def slide_06_architecture(prs, n):
    slide = new_slide(prs, "System Architecture", kicker="Multi-framework benchmark",
                      page_no=n)
    # left: framework cards; right: pipeline description
    add_panel(slide, MARGIN, Inches(1.7), Inches(6.35), Inches(4.75))
    header_tb, header_tf = _add_textbox(slide, MARGIN + Inches(0.2), Inches(1.8),
                                        Inches(6.0), Inches(0.4))
    _para(header_tf, "Six RL backends behind one Tinker-style API", size=15,
          color=PES_BLUE, bold=True, first=True)
    fw = [
        ("Tinker", "managed GRPO/PPO/REINFORCE (hosted API)"),
        ("SkyRL tx", "local Tinker API on own / vast.ai / Colab GPUs"),
        ("verl", "Volcano Engine / HybridFlow, Ray + vLLM"),
        ("OpenRLHF", "Ray + vLLM; PPO / DAPO / REINFORCE++"),
        ("TRL", "HuggingFace reference same-stack runner"),
        ("Atropos", "NousResearch RL environments (GSM8K/MATH/tool-use)"),
    ]
    tb, tf = _add_textbox(slide, MARGIN + Inches(0.2), Inches(2.25),
                          Inches(6.0), Inches(4.1))
    first = True
    for name, desc in fw:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.space_after = Pt(8)
        r1 = p.add_run(); r1.text = f"{name}  "
        _style_run(r1, 15.5, color=INK, bold=True)
        r2 = p.add_run(); r2.text = f"—  {desc}"
        _style_run(r2, 13.5, color=MUTED)
        _apply_bullet(p)
        first = False

    # right pipeline
    add_panel(slide, Inches(7.15), Inches(1.7), Inches(5.6), Inches(4.75),
              fill=WHITE)
    tb, tf = _add_textbox(slide, Inches(7.35), Inches(1.85), Inches(5.2), Inches(4.5))
    _para(tf, "Unified pipeline", size=15, color=PES_BLUE, bold=True,
          space_after=8, first=True)
    steps = [
        "unified.launcher → dispatch to any backend",
        "Envs: GSM8K, MATH-500, HumanEval, xLAM / synthetic tool-use",
        "Models: Qwen 0.6B–235B, Llama-3.1-8B, DeepSeek-V3.1, Nemotron-120B",
        "Telemetry: per-step reward, ZVF, gradient norm, length",
        "Aggregate → rliable + bootstrap → LaTeX tables + figures",
        "Held-out evaluator + 9 variance-mitigation baselines",
    ]
    for s in steps:
        _para(tf, s, size=13.5, color=INK, bullet=True, space_after=8)
    tb2, tf2 = _add_textbox(slide, Inches(7.35), Inches(6.02), Inches(5.2), Inches(0.4))
    _para(tf2, "Backends span Local GPU, vast.ai, Modal H100 and Colab.",
          size=12, color=MUTED, italic=True, first=True)


def slide_07_pillars(prs, n):
    slide = new_slide(prs, "Methodology: The Four De-Confound Pillars",
                      page_no=n)
    rows = [
        ["Pillar", "What it isolates", "Design"],
        ["1 · PPO vs GRPO (same-stack)",
         "Only the advantage estimator differs",
         "Identical model / data / stack"],
        ["2 · Zero-Variance Fraction",
         "Measured signal starvation + confounders",
         "Per-step ZVF across 9 libraries"],
        ["3 · Trainability / group size",
         "Preference density vs variance knob",
         "G ∈ {2,4,8,16} × seeds, held-out"],
        ["4 · Held-out generalization",
         "Does GRPO add real capability?",
         "Dr.GRPO vs GRPO, pre→post, McNemar"],
    ]
    add_table(slide, MARGIN, Inches(1.75), CONTENT_W, rows,
              col_widths=[0.30, 0.40, 0.30], font_size=14.5,
              row_height=Inches(0.78))
    tb, tf = _add_textbox(slide, MARGIN, Inches(6.5), CONTENT_W, Inches(0.5))
    _para(tf, "Each pillar holds the whole stack fixed and varies exactly one "
              "factor — so a difference cannot be blamed on the framework. "
              "Scripts: experiments/modal/.", size=13, color=INK, italic=True,
          first=True)


def slide_08_ppo_grpo(prs, n):
    slide = new_slide(prs, "Result 1 — Same-Stack PPO vs GRPO",
                      kicker="Pillar 1", page_no=n)
    items = [
        ("Same stack, only the estimator differs: PPO and GRPO are "
         "statistically indistinguishable.", 0, True),
        ("Welch t-test p = 0.7605; paired Δ = −0.002, p = 0.374.", 0, False),
        ("The advantage of “GRPO” in the wild is model-dependent, not "
         "universal:", 0, True),
        ("Qwen3-8B: GRPO 34.4% vs PPO 22.5% last-10  (GRPO +11.9 pp)", 1, False),
        ("Llama-3.1-8B-Instruct: GRPO 84.4% vs PPO 97.5%  (PPO +13.1 pp)", 1, False),
        ("Takeaway: report the stack, not the label — there is no universally "
         "superior algorithm at these scales.", 0, True),
    ]
    add_bullets(slide, MARGIN, Inches(1.7), Inches(6.5), Inches(5.0), items, size=16)
    add_image_fitted(slide, fig("paper", "figures", "v2", "ppo_vs_grpo.png"),
                     Inches(7.2), Inches(1.75), Inches(5.55), Inches(4.7),
                     caption="PPO vs GRPO across models (paper/figures/v2/)")
    tb, tf = _add_textbox(slide, MARGIN, Inches(6.62), Inches(6.5), Inches(0.4))
    _para(tf, "Source: experiment_summary.md; _shared_methods.tex row 20; "
              "frontier_synthesis_scaling.tex", size=10, color=MUTED, italic=True,
          first=True)


def slide_09_zvf(prs, n):
    slide = new_slide(prs, "Result 2 — Zero-Variance Fraction Diagnostic",
                      kicker="Pillar 2", page_no=n)
    items = [
        ("ZVF = fraction of rollout groups where all G completions get identical "
         "reward → zero gradient. Gradient utilization = 1 − ZVF.", 0, True),
        ("Vanilla GRPO wastes ~half its rollouts: mean ZVF 0.481.", 0, False),
        ("AERO halves it to 0.220 and is the only mitigation that keeps vanilla's "
         "last-10 accuracy (0.399 vs 0.379).", 0, False),
        ("Aggregated over N=80 (experiment, condition, seed) rows on one managed "
         "runtime — identical reward parser, β=0 KL, template.", 0, False),
        ("ZVF is a descriptive early-warning signal for signal starvation, "
         "portable across libraries.", 0, True),
    ]
    add_bullets(slide, MARGIN, Inches(1.7), Inches(6.4), Inches(5.0), items, size=16)
    add_image_fitted(slide, fig("paper", "figures", "zvf_by_library.png"),
                     Inches(7.15), Inches(1.75), Inches(5.6), Inches(4.7),
                     caption="Mean ZVF per library (paper/figures/)")
    tb, tf = _add_textbox(slide, MARGIN, Inches(6.62), Inches(6.4), Inches(0.4))
    _para(tf, "Source: paper/sections/zvf.tex, experiments/results/zvf_by_library.tsv",
          size=10, color=MUTED, italic=True, first=True)


def slide_10_group_size(prs, n):
    slide = new_slide(prs, "Result 3 — Group Size & Trainability",
                      kicker="Pillar 3", page_no=n)
    rows = [
        ["G", "Peak", "Last-10"],
        ["2", "50.0%", "37.5%"],
        ["4", "75.0%", "52.1%"],
        ["8", "100%", "84.4%"],
        ["16", "71.9%", "38.0%"],
    ]
    add_table(slide, MARGIN, Inches(1.85), Inches(3.7), rows,
              col_widths=[0.3, 0.35, 0.35], font_size=14, row_height=Inches(0.5))
    items = [
        ("Trainability is non-monotone in G — intermediate sizes beat the "
         "smallest and largest tested; no universal optimum.", 0, True),
        ("Effective signal-to-noise rises only sublinearly: ~52% of the √G "
         "ideal.", 0, False),
        ("On the near-ceiling task, held-out accuracy is flat G=2→G=16 "
         "(retention ratio ≈ 1.00) — a saturated-regime null, not "
         "equivalence on hard tasks.", 0, False),
        ("Framing: G is a preference-density dial (an implicit all-pairs / DPO-like "
         "contrast), not merely a variance knob.", 0, True),
    ]
    add_bullets(slide, Inches(4.55), Inches(1.75), Inches(8.2), Inches(3.2),
                items, size=15)
    add_image_fitted(slide, fig("paper", "figures", "v2", "group_size_ablation.png"),
                     Inches(4.6), Inches(4.55), Inches(8.15), Inches(1.95),
                     caption="Group-size ablation (paper/figures/v2/)")
    tb, tf = _add_textbox(slide, MARGIN, Inches(4.7), Inches(3.7), Inches(1.6))
    _para(tf, "Table: REPRODUCE.md §4.1 (paper Table 2, group-size block).",
          size=11, color=MUTED, italic=True, first=True)
    tb2, tf2 = _add_textbox(slide, MARGIN, Inches(6.62), CONTENT_W, Inches(0.4))
    _para(tf2, "Source: REPRODUCE.md §4.1; paper/sections/p3_abstract.tex",
          size=10, color=MUTED, italic=True, first=True)


def slide_11_length_heldout(prs, n):
    slide = new_slide(prs, "Result 4 — Length Bias & Held-Out Generalization",
                      kicker="Pillar 4", page_no=n)
    items = [
        ("Length bias (verbosity trap): 4 of 11 GRPO runs peak before 65% of "
         "training then decay — optimizing length, not correctness.", 0, True),
        ("Dr.GRPO is the only algorithm whose length keeps a negative slope "
         "(diff −0.0028, 95% CI [−0.0045, −0.0014]).", 0, False),
        ("Held-out GSM8K (Qwen3-8B, 5 seeds × 200 test items, greedy):", 0, True),
        ("GRPO 83.3% (SD 2.2%) vs base 82.0% — only +1.3 pp", 1, False),
        ("t = 1.32, p = 0.26 → not significant; capability is pre-existing", 1, False),
        ("Honest negative: GRPO surfaces capability more than it creates it.", 0, True),
    ]
    add_bullets(slide, MARGIN, Inches(1.7), Inches(6.5), Inches(5.0), items, size=15.5)
    add_image_fitted(slide, fig("paper", "figures", "length_vs_reward.png"),
                     Inches(7.2), Inches(1.75), Inches(5.55), Inches(4.7),
                     caption="Length vs reward coupling (paper/figures/)")
    tb, tf = _add_textbox(slide, MARGIN, Inches(6.62), Inches(6.5), Inches(0.4))
    _para(tf, "Source: reports/final/ held-out eval (grpo_agentic_llm_paper.tex); "
              "length_bias_iter48.tex", size=10, color=MUTED, italic=True, first=True)


def slide_12_scaling(prs, n):
    slide = new_slide(prs, "Result 5 — GRPO Scaling Behavior",
                      kicker="Pillar / Paper P1", page_no=n)
    items = [
        ("Studied across ~2.4 orders of magnitude (0.6B → ~671B) over 70+ runs, "
         "7 libraries, 5 model families.", 0, True),
        ("Central result is a conservative negative: no reliable positive cross-scale "
         "slope in mean GRPO reward — flat within (wide) uncertainty.", 0, False),
        ("A clean three-phase saturation law is falsified: only 2 of 12 anchors "
         "match; no single saturation exponent is identifiable.", 0, False),
        ("Nemotron-120B is a distinct collapse phase (zero-reward step fraction "
         "0.55 vs ≤ 0.067 elsewhere), not a point on a trend.", 0, False),
        ("Categorical capability structure carries more signal than log₁₀N "
         "— a local, stack-conditioned taxonomy, not a Chinchilla power law.", 0, True),
    ]
    add_bullets(slide, MARGIN, Inches(1.7), Inches(6.5), Inches(5.0), items, size=15)
    add_image_fitted(slide, fig("paper", "figures", "scaling_law_fit.png"),
                     Inches(7.2), Inches(1.75), Inches(5.55), Inches(4.7),
                     caption="Cross-scale GRPO reward fit (paper/figures/)")
    tb, tf = _add_textbox(slide, MARGIN, Inches(6.62), Inches(6.5), Inches(0.4))
    _para(tf, "Source: paper/sections/p1_abstract.tex; experiment_summary.md size ladder",
          size=10, color=MUTED, italic=True, first=True)


def slide_13_contrib_map(prs, n):
    slide = new_slide(prs, "P1–P8 Contribution Map",
                      kicker="Semester 4 paper series", page_no=n)
    rows = [
        ["#", "Title (abbreviated)"],
        ["P1", "Scaling Laws for GRPO Post-Training: A Cross-Library, Cross-Scale Study"],
        ["P2", "The Zero-Variance Fraction: A Descriptive Diagnostic for Signal Starvation"],
        ["P3", "Group Size in GRPO: Contrast Density and the Bridge to DPO"],
        ["P4", "Length Bias and Held-Out Generalization in GRPO and Dr. GRPO"],
        ["P5", "Report the Stack, Not the Label: RL-for-LLM Results Are Stack-Conditioned"],
        ["P6", "GRPO-Registry: A Machine-Readable Catalog of Group-Relative RL Stacks"],
        ["P7", "From Diagnostic to Controller: Adaptive Group-Size Intervention"],
        ["P8", "LLM vs. XGBoost in Credit-Card Fraud: Sensor and Scribe, Not Scorer"],
    ]
    add_table(slide, MARGIN, Inches(1.7), CONTENT_W, rows,
              col_widths=[0.07, 0.93], font_size=14, row_height=Inches(0.535))
    tb, tf = _add_textbox(slide, MARGIN, Inches(6.55), CONTENT_W, Inches(0.5))
    _para(tf, "Eight standalone papers, freshly compiled from canonical LaTeX "
              "sources. Map: sem 4 work/EXPERIMENTS.md (paper → source → evidence).",
          size=12.5, color=MUTED, italic=True, first=True)


def slide_14_fraud(prs, n):
    slide = new_slide(prs, "Applied Study — LLM vs XGBoost Fraud",
                      kicker="Paper P8", page_no=n)
    items = [
        ("Head-to-head on a synthetic card-fraud set (50,000 transactions, "
         "~1.4% fraud).", 0, True),
        ("XGBoost: test AUC 0.7955, F1 0.356, precision 0.723, recall 0.236 "
         "— trains in 0.5 s.", 0, False),
        ("Qwen3.5-4B SFT scorer: accuracy 0.792 but AUC 0.483 — chance-level "
         "ranking. LLM loses the real-time scorer seat.", 0, False),
        ("But “LLM vs XGBoost” is the wrong question. The LLM is a "
         "sensor (vision/doc evidence), scribe (SAR / adverse-action narration), "
         "and cold-start triage — ~85× cheaper than a human analyst.", 0, True),
        ("Hybrid: LLM = sensor + scribe, XGBoost = scorer.", 0, True),
    ]
    add_bullets(slide, MARGIN, Inches(1.7), Inches(6.5), Inches(5.0), items, size=15.5)
    add_image_fitted(slide, fig("paper", "figures", "p8_cost_per_caught.png"),
                     Inches(7.2), Inches(1.75), Inches(5.55), Inches(4.7),
                     caption="Cost per fraud caught (paper/figures/)")
    tb, tf = _add_textbox(slide, MARGIN, Inches(6.62), Inches(6.5), Inches(0.4))
    _para(tf, "Source: paper/sections/p8_abstract.tex; xgboost_results.json",
          size=10, color=MUTED, italic=True, first=True)


def slide_15_reproducibility(prs, n):
    slide = new_slide(prs, "Reproducibility & Audit Apparatus", page_no=n)
    add_panel(slide, MARGIN, Inches(1.7), Inches(5.95), Inches(4.75))
    tb, tf = _add_textbox(slide, MARGIN + Inches(0.22), Inches(1.82),
                          Inches(5.55), Inches(4.5))
    _para(tf, "Integrity suite", size=15, color=PES_BLUE, bold=True,
          space_after=8, first=True)
    for s in [
        "run_all_audits.py → 13-audit integrity suite",
        "17 standalone *_audit.py checks; each prints METRIC …=N",
        "Covers paper / claim / abstract-scope / held-out / anonymization / "
        "submission / blind-review",
        "Multi-seed protocol: canonical 10-seed set",
    ]:
        _para(tf, s, size=13.5, color=INK, bullet=True, space_after=8)

    add_panel(slide, Inches(6.85), Inches(1.7), Inches(5.9), Inches(4.75), fill=WHITE)
    tb2, tf2 = _add_textbox(slide, Inches(7.05), Inches(1.82), Inches(5.5), Inches(4.5))
    _para(tf2, "One-command reproduction", size=15, color=PES_BLUE, bold=True,
          space_after=8, first=True)
    for s in [
        "Pinned Dockerfile (~6 min cold build)",
        "scripts/smoke_test.sh — <10 min, CPU-only",
        "Headline: GRPO Qwen3-8B GSM8K, peak 62.5%, last-10 34.4% "
        "(±5 pts tolerance)",
        "Data pinned by HF revision (GSM8K 7,473 train / 1,319 test)",
        "Full paper reproduction budget: ~446 GPU-h; headline ~5 GPU-h",
    ]:
        _para(tf2, s, size=13.5, color=INK, bullet=True, space_after=8)
    tb3, tf3 = _add_textbox(slide, MARGIN, Inches(6.55), CONTENT_W, Inches(0.4))
    _para(tf3, "Sources: run_all_audits.py; REPRODUCE.md; Dockerfile",
          size=10.5, color=MUTED, italic=True, first=True)


def slide_16_publications(prs, n):
    slide = new_slide(prs, "Publications & Submissions", page_no=n)
    add_panel(slide, MARGIN, Inches(1.7), Inches(5.95), Inches(4.75), fill=LIGHT_BG)
    tb, tf = _add_textbox(slide, MARGIN + Inches(0.22), Inches(1.82),
                          Inches(5.55), Inches(4.5))
    _para(tf, "Semester 3 — Group 6", size=15, color=PES_BLUE, bold=True,
          space_after=6, first=True)
    for s in [
        "NeurIPS 2026 main-track submission (anonymized blind-review package)",
        "Blind-review manifest dated April 19, 2026",
        "Frozen at tag capstone-final-2026-04-25 (commit 21a99ef7)",
        "Capstone final report + group benchmark paper",
    ]:
        _para(tf, s, size=13.5, color=INK, bullet=True, space_after=7)

    add_panel(slide, Inches(6.85), Inches(1.7), Inches(5.9), Inches(4.75), fill=WHITE)
    tb2, tf2 = _add_textbox(slide, Inches(7.05), Inches(1.82), Inches(5.5), Inches(4.5))
    _para(tf2, "Semester 4 — solo continuation", size=15, color=PES_BLUE,
          bold=True, space_after=6, first=True)
    for s in [
        "NeurIPS 2026 workshop submission (16-page anonymous artifact)",
        "“RL-Finetuning Bench: An Exploratory Workshop Artifact for "
        "GRPO-Style Post-Training Across Libraries, Models, and Backends”",
        "First committed June 21, 2026 (commit b0ac85bf)",
        "Eight standalone papers P1–P8 + evidence map",
    ]:
        _para(tf2, s, size=13.5, color=INK, bullet=True, space_after=7)
    tb3, tf3 = _add_textbox(slide, MARGIN, Inches(6.55), CONTENT_W, Inches(0.4))
    _para(tf3, "Sources: PROJECT_HISTORY.md; sem 4 work/submissions/neurips-workshop/README.md",
          size=10.5, color=MUTED, italic=True, first=True)


def slide_17_demo(prs, n):
    slide = new_slide(prs, "Demo Pointer", page_no=n)
    items = [
        ("Executed demo notebook: demo_recording/thesis_demo.ipynb — 11 cells, "
         "all real outputs (results tables, ZVF, group-size, P8 metrics); "
         "HTML export + recorded walkthrough thesis_demo.mp4.", 0, True),
        ("Live HuggingFace Space (pre-computed, CPU-only): "
         "arvindcr4-tinkerrl-bench-demo.hf.space", 0, True),
        ("Four tabs: tool-use comparison · ZVF diagnostic · team artifacts "
         "· headline numbers.", 0, False),
        ("Phase-1 defense demo runbook: submission/demo/demo.sh + "
         "reports/esa_phase1/CODE_WALKTHROUGH.md (DEMO STATUS: PASS).", 0, True),
        ("Headline shown live: tool-call JSON validity 0% → 92% after "
         "SFT+GRPO.", 0, False),
    ]
    add_bullets(slide, MARGIN, Inches(1.7), CONTENT_W, Inches(4.5), items, size=17)
    tb, tf = _add_textbox(slide, MARGIN, Inches(6.4), CONTENT_W, Inches(0.5))
    _para(tf, "Sources: demo_recording/thesis_demo.ipynb; reports/esa_phase1/CODE_WALKTHROUGH.md",
          size=10.5, color=MUTED, italic=True, first=True)


def slide_18_ownership(prs, n):
    slide = new_slide(prs, "Semester 3 vs Semester 4 Ownership", page_no=n)
    rows = [
        ["", "Semester 3 — Group 6", "Semester 4 — Solo"],
        ["Team", "Six students; guides Paduri & Darapaneni",
         "Arvind C R; guide Ramesh Prakash Guledgudd"],
        ["Deliverables", "Lit foundation, multi-framework env, first experiments, "
         "group paper, NeurIPS main-track",
         "P1–P8 papers, NeurIPS workshop, expanded diagnostics & controllers, "
         "fraud study"],
        ["Boundary", "Frozen: tag capstone-final-2026-04-25 (21a99ef7)",
         "Starts after boundary; workshop from June 21, 2026"],
        ["Evidence", "Immutable historical deliverables",
         "206 recorded analysis iterations; new audits & prototypes"],
    ]
    add_table(slide, MARGIN, Inches(1.7), CONTENT_W, rows,
              col_widths=[0.16, 0.42, 0.42], font_size=13, row_height=Inches(0.95))
    tb, tf = _add_textbox(slide, MARGIN, Inches(6.75), CONTENT_W, Inches(0.35))
    _para(tf, "Shared code stays at repo root; the two semester folders are curated "
              "academic views. Source: PROJECT_HISTORY.md; sem 4 work/README.md",
          size=11, color=MUTED, italic=True, first=True)


def slide_19_limitations(prs, n):
    slide = new_slide(prs, "Limitations & Future Work", page_no=n)
    add_panel(slide, MARGIN, Inches(1.7), Inches(5.95), Inches(4.75), fill=LIGHT_BG)
    tb, tf = _add_textbox(slide, MARGIN + Inches(0.22), Inches(1.82),
                          Inches(5.55), Inches(4.5))
    _para(tf, "Limitations", size=15, color=PES_BLUE, bold=True, space_after=6,
          first=True)
    for s in [
        "Single-seed Tinker runs; no variance/significance on Tinker data",
        "Short horizons (30–50 steps) — no asymptotic claims",
        "Closed Tinker internals: results measure the platform, not the algorithm",
        "Saturated tasks compress variance (G-equivalence is a null, not proof)",
        "Exploratory, not pre-registered; narrow benchmark coverage",
    ]:
        _para(tf, s, size=13, color=INK, bullet=True, space_after=7)

    add_panel(slide, Inches(6.85), Inches(1.7), Inches(5.9), Inches(4.75), fill=WHITE)
    tb2, tf2 = _add_textbox(slide, Inches(7.05), Inches(1.82), Inches(5.5), Inches(4.5))
    _para(tf2, "Future work", size=15, color=PES_BLUE, bold=True, space_after=6,
          first=True)
    for s in [
        "Repeat ZVF / group-size on harder, non-saturating tasks",
        "Direct gradient-vector validation of the DPO-like interpretation",
        "Full-FT vs LoRA re-runs across every baseline library",
        "Exhaustive LR × rank × batch × G sweeps",
        "P7 adaptive ZVF controller: from diagnostic to intervention",
    ]:
        _para(tf2, s, size=13, color=INK, bullet=True, space_after=7)
    tb3, tf3 = _add_textbox(slide, MARGIN, Inches(6.55), CONTENT_W, Inches(0.4))
    _para(tf3, "Source: LIMITATIONS_AND_IMPACT.md §7", size=10.5, color=MUTED,
          italic=True, first=True)


def slide_20_conclusions(prs, n):
    slide = new_slide(prs, "Conclusions", page_no=n)
    items = [
        ("A multi-framework GRPO benchmark (6 backends, 44+ curated runs, 70+ across "
         "the roster) that runs identical workloads and reports the stack.", 0, True),
        ("Four de-confound pillars show “GRPO” gains are stack- and "
         "model-conditioned: same-stack PPO≡GRPO (p=0.76).", 0, True),
        ("ZVF (and gradient utilization = 1−ZVF) is a portable diagnostic for "
         "signal starvation; group size is a preference-density dial.", 0, True),
        ("Honest held-out negative: GRPO adds only +1.3 pts on GSM8K (p=0.26) — "
         "it surfaces capability more than it creates it.", 0, True),
        ("Reproducibility-first: 13-audit suite, Docker, REPRODUCE.md, and eight "
         "papers P1–P8 + a NeurIPS workshop artifact.", 0, True),
    ]
    add_bullets(slide, MARGIN, Inches(1.7), CONTENT_W, Inches(5.0), items, size=17.5)


def slide_21_thanks(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, PES_BLUE)
    tb, tf = _add_textbox(slide, Inches(1.0), Inches(2.6), Inches(11.3), Inches(2.4),
                          anchor=MSO_ANCHOR.MIDDLE)
    _para(tf, "Thank You", size=52, color=WHITE, bold=True,
          align=PP_ALIGN.CENTER, space_after=10, first=True)
    _para(tf, "Questions & Discussion", size=24, color=RGBColor(0xCF, 0xDE, 0xF0),
          align=PP_ALIGN.CENTER, space_after=24)
    _para(tf, "Arvind C R  ·  SRN <SRN>  ·  M.Tech Data Science & AI, PES University",
          size=15, color=WHITE, align=PP_ALIGN.CENTER, space_after=4)
    _para(tf, "Guide: Ramesh Prakash Guledgudd  ·  github.com/arvindcr4/tinker-rl-lab",
          size=14, color=RGBColor(0xCF, 0xDE, 0xF0), align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------
def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_agenda(prs, 2)
    slide_03_problem(prs, 3)
    slide_04_objectives(prs, 4)
    slide_05_literature(prs, 5)
    slide_06_architecture(prs, 6)
    slide_07_pillars(prs, 7)
    slide_08_ppo_grpo(prs, 8)
    slide_09_zvf(prs, 9)
    slide_10_group_size(prs, 10)
    slide_11_length_heldout(prs, 11)
    slide_12_scaling(prs, 12)
    slide_13_contrib_map(prs, 13)
    slide_14_fraud(prs, 14)
    slide_15_reproducibility(prs, 15)
    slide_16_publications(prs, 16)
    slide_17_demo(prs, 17)
    slide_18_ownership(prs, 18)
    slide_19_limitations(prs, 19)
    slide_20_conclusions(prs, 20)
    slide_21_thanks(prs, 21)

    return prs


def main():
    prs = build()
    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH}")
    print(f"Slides: {len(prs.slides._sldIdLst)}")
    print(f"Figures embedded ({len(EMBEDDED_FIGURES)}):")
    for f in EMBEDDED_FIGURES:
        print(f"  + {f}")
    if MISSING_FIGURES:
        print(f"Figures MISSING ({len(MISSING_FIGURES)}) — text fallback used:")
        for f in MISSING_FIGURES:
            print(f"  ! {f}")

    if "--check" in sys.argv:
        reloaded = Presentation(OUT_PATH)
        print(f"Round-trip load OK: {len(reloaded.slides._sldIdLst)} slides")


if __name__ == "__main__":
    main()
