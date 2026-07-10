#!/usr/bin/env python3
"""Build the ZVF-aligned ESA Phase-1 deck from the existing draft (same official
ESA template layout), rewriting content to the ZVF framing + corrected findings."""
import copy
from pptx import Presentation
from pptx.util import Pt

SRC = "reports/esa_phase1/ESA_Phase1_ArvindCR_DRAFT.pptx"
OUT = "reports/submissions/PESU_Phase1_ESA_ZVF_ArvindCR.pptx"

HEADINGS = {
    "problem statement", "abstract and scope", "literature survey",
    "suggestions from review - 3", "suggestions from review – 3", "design approach",
    "design constraints, assumptions & dependencies", "design details",
    "proposed methodology / approach", "architecture (if applicable)",
    "design description (if applicable)", "technologies used", "project progress",
    "references", "any other information", "agenda",
}

def set_content(shape, lines):
    tf = shape.text_frame
    paras = tf.paragraphs
    # capture a template paragraph element (with its bullet/indent formatting) and font size
    tmpl_p = copy.deepcopy(paras[0]._p)
    base_sz = None
    for p in paras:
        for r in p.runs:
            if r.font.size:
                base_sz = r.font.size; break
        if base_sz: break
    # remove all existing paragraphs
    for p in list(paras):
        p._p.getparent().remove(p._p)
    body = tf._txBody
    for line in lines:
        newp = copy.deepcopy(tmpl_p)
        # strip existing runs from the cloned paragraph
        for r in newp.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}r'):
            newp.remove(r)
        body.append(newp)
    # now set text via python-pptx paragraph API
    for para, line in zip(tf.paragraphs, lines):
        run = para.add_run()
        run.text = line
        if base_sz:
            run.font.size = base_sz

def content_shape(slide):
    """Return the non-heading, non-empty text shape (the content box)."""
    cands = [sh for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
    for sh in cands:
        t = sh.text_frame.text.strip().lower()
        if t not in HEADINGS and not t.startswith("ue20cs971"):
            return sh
    return None

prs = Presentation(SRC)
S = prs.slides

# ---- Slide 1: Title ----
for sh in S[0].shapes:
    if sh.has_text_frame and "Project Title" in sh.text_frame.text:
        set_content(sh, [
            "Project Title : Zero-Variance Fraction (ZVF):",
            "                A Diagnostic Benchmark for GRPO Post-Training of Language Models",
            "Project Guide : Ramesh Prakash Guledgudd",
            "Student : Arvind C R   (SRN: PES2PGE24DS140)",
            "Programme : M.Tech, Data Science and Machine Learning  —  Individual Project",
            "Course : UE20CS971 — Project Phase-1",
        ])

# ---- content per heading ----
CONTENT = {
 "problem statement": [
  "RLVR / GRPO is now the de facto default recipe for post-training reasoning LLMs, yet reported gains are hard to reproduce and the internal training signal that would explain them is invisible.",
  "Reproducibility crisis: results are spread across TRL / veRL / OpenRLHF / Tinker with different, often undocumented reward parsers, temperatures, KL and clip ranges — a headline number cannot be attributed to the ALGORITHM vs the STACK.",
  "Gradient waste is invisible: standard logs report reward and loss, not what fraction of the batch produced any learning signal at all.",
  "No provenance standard: there is no machine-readable record of what was trained, on which verifier version, with which rollouts — so contamination and reward drift go undetected.",
 ],
 "abstract and scope": [
  "We build a unified, provenance-tracked benchmark + per-step diagnostic harness for GRPO/RLVR, organised around ONE measurable quantity: the Zero-Variance Fraction (ZVF) — the fraction of prompt groups whose within-group reward variance is zero (hence no policy-gradient signal).",
  "The same task, reward and decoding are held fixed across back-ends so differences are attributable to the stack; every step emits ZVF, gradient utilisation, entropy, length and KL, and raw reward tensors are saved for recomputation.",
  "On top of the harness we run a portfolio of eight studies (P1–P8). Framed as MEASUREMENT / diagnostics + reproducibility tooling — NOT a 'GRPO wins' leaderboard.",
  "Scope: small–mid open models (Qwen2.5 / Qwen3.5, 1.5B–4B), GSM8K + a HumanEval subset; held-out sets kept small and every gain reported with its n; honest nulls included.",
 ],
 "literature survey": [
  "DeepSeekMath (Shao et al., 2024) — introduces GRPO, our baseline; DeepSeek-R1 (2025) — motivates scale + reproducibility questions.",
  "RLOO (Ahmadian et al., 2024) — leave-one-out within-group baseline; directly relevant to the P2 re-baselining analysis.",
  "Dr.GRPO (Liu 2025) — length/difficulty bias fix; DAPO (Yu 2025) — DYNAMIC SAMPLING to avoid zero-variance groups (closest prior art for ZVF).",
  "Collapse variants: GVPO, AERO, NGRPO, lambda-GRPO, MC-GRPO, AVSPO — attack advantage collapse / length bias from different angles (crowded).",
  "RLVR (TULU-3, 2024); provenance: Datasheets (Gebru 2021), Model Cards (Mitchell 2019). Positioning: narrow crowded pillars (P2, P4) to measurement; open ground = systems/provenance (P5/P6/P8).",
 ],
 "suggestions from review - 3": [
  "Panel remarks (Review-3, 4/5 July): tighten statistical claims (single-seed frontier runs are not confirmatory); separate validated results from prototypes; strengthen novelty positioning vs recent GRPO variants; report held-out / generalisation honestly.  [confirm against your panel notes]",
  "Actions taken since Review-3:",
  "• Ran a powered multi-seed campaign that CONVERTED single-seed 'wins' (curriculum, G=4 sweet spot, P1 layer-freeze) into honest nulls.",
  "• Ran an adversarial review with FOUR frontier models (Kimi K2.7, MiniMax-M3, agy, GPT/codex) and applied the fixes.",
  "• Retitled the project to ZVF and produced a 50-page Phase-1 report in the official PES sample format (solo authorship).",
 ],
 "design approach": [
  "Benchmark-harness approach: a unified launcher runs the SAME task / reward / decoding across multiple RL frameworks, so any difference is attributable to the stack, not the experiment.",
  "Per-step telemetry (reward, ZVF / gradient-utilisation, entropy, completion length, KL) + held-out evaluation + a multi-seed statistical protocol (seed as the unit; matched baselines).",
  "Why: it is the only way to separate 'stack effect' from 'algorithm effect', which single-framework studies conflate.",
  "Drawbacks (stated openly): managed back-ends (Tinker) limit internal inspection; compute limits force short horizons and small held-out sets — so gains are reported as noise-limited.",
 ],
 "design constraints, assumptions & dependencies": [
  "Constraints: limited compute → short training horizons (6–10 steps) and small held-out sets (n = 8–20); binary verifiable rewards; capped generation length.",
  "Assumptions: GSM8K / HumanEval approximate reasoning trainability; matched task/reward/decoding make cross-framework runs comparable.",
  "Dependency note: the managed back-end does not expose per-layer gradients → P1 uses a separate white-box (Colab) path.",
  "Dependencies: Tinker training API, Weights & Biases, HuggingFace checkpoints, Google Colab / Modal compute.",
  "Impact: single-seed cells carry no valid CIs and are reported descriptively; headline claims are re-tested multi-seed.",
 ],
 "design details": [
  "Unified launcher — one config drives TRL / veRL / OpenRLHF / Tinker adapters (reward grader + decoding are shared code).",
  "ZVF formalisation: A_i = (r_i − mu_g)/(sigma_g+eps); when sigma_g = 0 every advantage is 0, so the group gives no policy-gradient signal. ZVF = batch-mean of 1[sigma_g = 0].",
  "Telemetry + registry — per-step JSONL (reward, ZVF, gradient utilisation, all-correct/all-wrong collapse, entropy, length, KL) to W&B; raw reward tensors saved for RECOMPUTATION.",
  "Evaluation harness — held-out GSM8K slice, checkpoint selection, and recompute-from-tensors diagnostics (p2_collapse_analysis.py, p8_detector.py).",
 ],
 "proposed methodology / approach": [
  "Approach: controlled sweeps (group size G, seeds, baseline vs curriculum) with identical reward/decoding; measure ZVF + held-out accuracy. Posture: verify → multi-seed → report the honest null.",
  "P2 (measured): mean ZVF = 0.72–0.77; collapse is EASY-driven (65–71% all-correct); re-baselining recovers no real gradient.",
  "P3 / curriculum (NULL): no robust group-size sweet spot; curriculum drives ZVF 0.50→0.00 but at ~5–6× sampling cost and +0.028 vs −0.028 held-out (not significant).",
  "P1 (NULL): step-1 layer-adaptation predictability = 1.0 at 1.5B collapses to 0.11 (≈chance) on a scaled 3B multi-seed re-test.",
  "P8 (POSITIVE): a telemetry-based integrity auditor reaches AUROC 0.84 vs 0.43 for a reward-only baseline.",
 ],
 "architecture (if applicable)": [
  "High-level flow: Unified Launcher → Framework Adapter (TRL / veRL / OpenRLHF / Tinker) → Environment + Reward → Rollout + per-step Telemetry (ZVF / GU / entropy / KL) → Registry / W&B → Evaluation + Figures.",
  "Logical groups: (a) orchestration/config, (b) framework adapters, (c) environments & reward functions, (d) telemetry & registry, (e) evaluation & reporting.",
  "See report Fig. 4.1 (unified harness) and Fig. 4.4 (eight-study portfolio) — insert diagram here.",
 ],
 "design description (if applicable)": [
  "P5 MIN-REPORT schema: a minimal per-run provenance record (model, init, framework, reward, temperature, seed, group size, verifier version, rollout hashes).",
  "P6 GRPO-Registry: machine-readable, versioned entries per run, with a validation schema and a claim-vs-evidence ledger.",
  "Data flow: raw per-step JSONL → aggregated run record → registry entry → tables/figures; every headline number traces back to a stored artefact.",
  "Recompute-from-tensors: published ZVF and P8 AUROC are regenerated by scripts, not trusted from a logged summary.",
 ],
 "technologies used": [
  "RL frameworks: Tinker (managed back-end); TRL, veRL / HybridFlow, OpenRLHF (unified via adapters).",
  "Models: Qwen2.5-1.5B/3B, Qwen3.5-4B (verify exact model IDs before final submission).",
  "Tasks / data: GSM8K (math, exact-match), HumanEval subset (code, unit-test).",
  "Infra: PyTorch, HuggingFace Transformers + PEFT (LoRA), scikit-learn (P8 detector), Weights & Biases (telemetry), Google Colab / Modal (compute).",
  "Analysis / reporting: Python (recompute + figures), LaTeX / TikZ (50-page report + this deck).",
 ],
 "project progress": [
  "Phase-1 core (measurement + diagnostic stack) COMPLETE ≈ 80% code; cross-framework four-way run and the designed levers (P4 length, P7 controller) are Phase-2.",
  "Done: harness + per-step telemetry; P2 ZVF measured; P3 / curriculum / campaign nulls (multi-seed); P1 white-box (1.5B + scaled 3B); P8 detector (AUROC 0.84); P5 provenance prototype.",
  "Deliverables: 50-page Phase-1 report (PES sample format, solo); this ESA deck; code walkthrough; 4-model adversarial review applied.",
  "In progress / Phase-2: adequately-powered eval (full GSM8K, ≥5 seeds, CIs); four-way cross-framework divergence table; harden P8; the P5 standard; publication.",
 ],
 "references": [
  "[1] Z. Shao et al., 'DeepSeekMath: Pushing the Limits of Mathematical Reasoning...', arXiv:2402.03300, 2024. (GRPO)",
  "[2] DeepSeek-AI, 'DeepSeek-R1: Incentivizing Reasoning Capability in LLMs via RL', arXiv:2501.12948, 2025.",
  "[3] A. Ahmadian et al., 'Back to Basics: Revisiting REINFORCE-Style Optimization (RLOO)', arXiv:2402.14740, 2024.",
  "[4] Q. Yu et al., 'DAPO: An Open-Source LLM RL System at Scale', arXiv:2503.14476, 2025. (dynamic sampling)",
  "[5] N. Lambert et al., 'TULU 3: ... Open Language Model Post-Training (RLVR)', arXiv:2411.15124, 2024.",
  "[6] T. Gebru et al., 'Datasheets for Datasets', Commun. ACM, 2021;  M. Mitchell et al., 'Model Cards', ACM FAT*, 2019.",
  "(Full list in the Phase-1 report; IEEE format; arXiv IDs verified 2026-07-06.)",
 ],
 "any other information": [
  "Reproducibility: recompute-from-tensors scripts, W&B logs, and a run registry mapping each result to its producing run; report + code are version-controlled.",
  "Honesty statement: single-seed 'wins' (curriculum, G=4, P1 layer-freeze) have been demoted to nulls; the project is framed as a diagnostic / measurement study, and negative results are reported.",
  "Adversarial review: the report was reviewed by four frontier models (Kimi, MiniMax, agy, GPT); fixes applied — KL-gradient wording, statistics framing, standalone 'Contribution of the Candidate' section, plagiarism + code statements, and page count raised to 50pp.",
  "This is an INDIVIDUAL project (SRN PES2PGE24DS140) — not a group submission.",
 ],
}

for sl in list(S)[1:]:  # skip title; agenda kept as-is unless matched
    # find heading text on the slide
    heading = None
    for sh in sl.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip().lower()
            if t in HEADINGS:
                heading = t; break
    if heading in CONTENT:
        cs = content_shape(sl)
        if cs is not None:
            set_content(cs, CONTENT[heading])

prs.save(OUT)
print("saved", OUT, "slides", len(prs.slides))
