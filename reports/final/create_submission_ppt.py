#!/usr/bin/env python3
"""Create the final capstone submission PPTX."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BG = RGBColor(22, 28, 42)
CARD = RGBColor(31, 42, 61)
TEAL = RGBColor(46, 196, 182)
AMBER = RGBColor(255, 190, 106)
RED = RGBColor(239, 99, 81)
WHITE = RGBColor(246, 248, 250)
MUTED = RGBColor(179, 190, 205)
BLUE = RGBColor(92, 147, 255)
GREEN = RGBColor(97, 210, 133)


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def text(slide, x, y, w, h, value, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    p = frame.paragraphs[0]
    p.text = value
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return box


def card(slide, x, y, w, h, fill=CARD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(54, 68, 91)
    shape.line.width = Pt(1)
    return shape


def title(slide, number, heading, subtitle=None):
    text(slide, 0.5, 0.25, 0.5, 0.3, number, 11, TEAL, True)
    text(slide, 1.0, 0.22, 11.8, 0.45, heading, 28, WHITE, True)
    if subtitle:
        text(slide, 1.0, 0.72, 11.2, 0.35, subtitle, 13, MUTED)


def bullets(slide, x, y, items, color=WHITE, size=15, gap=0.45):
    for i, item in enumerate(items):
        text(slide, x, y + i * gap, 0.22, 0.25, "•", size, TEAL, True)
        text(slide, x + 0.28, y + i * gap, 5.7, 0.35, item, size, color)


def stat(slide, x, y, value, label, color=TEAL):
    card(slide, x, y, 2.45, 1.1)
    text(slide, x, y + 0.12, 2.45, 0.45, value, 28, color, True, PP_ALIGN.CENTER)
    text(slide, x + 0.12, y + 0.66, 2.2, 0.3, label, 10, MUTED, False, PP_ALIGN.CENTER)


def add_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text(slide, 0.9, 0.95, 11.5, 0.55, "Reinforcement Learning for Agentic LLM Fine-Tuning", 28, TEAL, True, PP_ALIGN.CENTER)
    text(slide, 0.9, 1.7, 11.5, 0.9, "GRPO Across Tool Use, Code Generation, and Mathematical Reasoning", 34, WHITE, True, PP_ALIGN.CENTER)
    text(slide, 1.6, 3.0, 10.1, 0.35, "Final Capstone Presentation | Group 6 | PES University", 17, MUTED, False, PP_ALIGN.CENTER)
    stat(slide, 1.3, 4.2, "32+", "structured runs")
    stat(slide, 4.1, 4.2, "5", "task families", BLUE)
    stat(slide, 6.9, 4.2, "4", "training backends", AMBER)
    stat(slide, 9.7, 4.2, "76", "report pages", GREEN)
    text(slide, 1.0, 6.55, 11.3, 0.3, "Guides: Prof. Narayana Darapaneni and Mr. Anwesh Reddy Paduri", 12, MUTED, False, PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "01", "Problem and Research Questions", "When does critic-free RL improve LLM agents, and when does it fail?")
    card(slide, 0.7, 1.35, 5.7, 4.9)
    text(slide, 1.0, 1.65, 5.1, 0.4, "Motivation", 20, TEAL, True)
    bullets(slide, 1.0, 2.25, [
        "Agents need verifiable behavior, not only imitation.",
        "Tool calls must be schema-valid and semantically correct.",
        "Math and code tasks expose whether RL improves reasoning.",
        "GRPO removes the PPO value model, but depends on reward variance.",
    ])
    card(slide, 6.9, 1.35, 5.7, 4.9)
    text(slide, 7.2, 1.65, 5.1, 0.4, "Core Questions", 20, AMBER, True)
    bullets(slide, 7.2, 2.25, [
        "Can GRPO improve tool-use beyond supervised tuning?",
        "How do model size and instruction tuning affect GRPO?",
        "How do GRPO and PPO compare under similar constraints?",
        "Can ZVF and GU diagnose bad RL runs early?",
    ], size=14)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "02", "Methodology", "A shared rollout-score-update loop across tool use, math, and code")
    labels = [
        ("Prompt Dataset", 0.8, 2.1, BLUE),
        ("Policy Samples", 3.25, 2.1, TEAL),
        ("Reward Function", 5.7, 2.1, AMBER),
        ("GRPO / PPO Update", 8.15, 2.1, GREEN),
        ("Evaluation Logs", 10.6, 2.1, RED),
    ]
    for label, x, y, color in labels:
        card(slide, x, y, 1.85, 1.15, color)
        text(slide, x + 0.12, y + 0.36, 1.6, 0.35, label, 13, BG, True, PP_ALIGN.CENTER)
    for x in [2.66, 5.12, 7.56, 10.02]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(2.48), Inches(0.42), Inches(0.28))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MUTED
        arrow.line.fill.background()
    card(slide, 1.2, 4.05, 10.9, 1.45)
    text(slide, 1.45, 4.25, 10.4, 0.35, "Diagnostics", 18, TEAL, True)
    text(slide, 1.45, 4.75, 10.4, 0.35, "Zero-Variance Fraction (ZVF) measures dead groups; Gradient Utilization (GU) measures how much sampled rollout data can produce a useful policy update.", 15, WHITE)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "03", "Experiment Inventory", "Final report consolidates the first chapter, interim scope, literature survey, and all completed experiments")
    cols = [
        ("Tool Use", ["SFT warm-up", "GRPO on schema rewards", "0.5B and 1.5B Qwen variants"], TEAL),
        ("Math", ["Arithmetic", "GSM8K", "MATH and held-out GSM8K"], BLUE),
        ("Code", ["HumanEval-style rewards", "Semantic code", "Framework sensitivity"], AMBER),
        ("Baselines", ["PPO: SB3, CleanRL", "TRL GRPO", "DPO & distillation"], GREEN),
    ]
    for i, (head, lines, color) in enumerate(cols):
        x = 0.55 + i * 3.15
        card(slide, x, 1.5, 2.75, 4.6)
        text(slide, x + 0.15, 1.75, 2.45, 0.35, head, 18, color, True)
        bullets(slide, x + 0.18, 2.35, lines, size=12, gap=0.55)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "04", "Key Result: Tool Use", "GRPO is most useful when the model already has a supervised behavioral prior")
    stat(slide, 0.8, 1.5, "SFT", "establishes schema")
    stat(slide, 3.65, 1.5, "GRPO", "sharpens validity", GREEN)
    stat(slide, 6.5, 1.5, "ZVF", "predicts collapse", AMBER)
    stat(slide, 9.35, 1.5, "GU", "tracks usable signal", BLUE)
    card(slide, 1.0, 3.55, 11.3, 1.85)
    text(slide, 1.35, 3.85, 10.7, 0.35, "Interpretation", 19, TEAL, True)
    text(slide, 1.35, 4.35, 10.7, 0.55, "Tool-call rewards are structured and immediate. Once SFT teaches the output format, GRPO can exploit reward variation to improve valid action emission.", 16, WHITE)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "05", "Key Result: Mathematical Reasoning", "Scale and initialization matter more than the optimizer label")
    rows = [
        ("Small/base models", "Weak or unstable gains"),
        ("Instruct-tuned models", "Higher starting point and easier saturation"),
        ("Larger Qwen / MoE runs", "Better trajectories, but held-out caution remains"),
        ("Held-out GSM8K", "Effect is weaker than training reward suggests"),
    ]
    for i, (left, right) in enumerate(rows):
        y = 1.55 + i * 1.0
        card(slide, 1.0, y, 4.3, 0.7, CARD)
        text(slide, 1.25, y + 0.18, 3.8, 0.25, left, 14, TEAL, True)
        card(slide, 5.6, y, 6.7, 0.7, CARD)
        text(slide, 5.85, y + 0.18, 6.2, 0.25, right, 14, WHITE)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "06", "Key Result: Code and PPO Baselines", "Negative results clarified the structural ceiling")
    card(slide, 0.75, 1.45, 5.65, 4.8)
    text(slide, 1.05, 1.75, 5.05, 0.35, "Code Generation", 20, AMBER, True)
    bullets(slide, 1.05, 2.35, [
        "Reward sparsity made on-policy learning brittle.",
        "Semantic correctness required stronger evaluation than format checks.",
        "Short runs were not enough to establish robust code gains.",
    ], size=14)
    card(slide, 6.9, 1.45, 5.65, 4.8)
    text(slide, 7.2, 1.75, 5.05, 0.35, "PPO and Classical RL", 20, GREEN, True)
    bullets(slide, 7.2, 2.35, [
        "Small CartPole-style PPO baselines did not transfer to LLM behavior.",
        "Framework choice affected stability and throughput.",
        "The bottleneck was reward signal quality, not only algorithm choice.",
    ], size=14)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "07", "Final Conclusions", "GRPO amplifies useful priors; it does not reliably create missing reasoning ability")
    conclusions = [
        ("1", "Best domain", "Structured tool-use after supervised warm-up"),
        ("2", "Main failure mode", "Zero reward variance leads to dead policy updates"),
        ("3", "Scaling lesson", "Bigger and instruction-tuned models give GRPO more usable signal"),
        ("4", "Evaluation lesson", "Training reward saturation is not held-out generalization"),
    ]
    for i, (num, head, body) in enumerate(conclusions):
        y = 1.35 + i * 1.15
        card(slide, 0.9, y, 1.0, 0.72, TEAL)
        text(slide, 0.9, y + 0.16, 1.0, 0.3, num, 18, BG, True, PP_ALIGN.CENTER)
        text(slide, 2.15, y + 0.05, 3.1, 0.3, head, 16, WHITE, True)
        text(slide, 5.1, y + 0.05, 7.1, 0.42, body, 15, MUTED)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "08", "Submission Artifacts", "Everything is packaged for review and reproduction")
    artifacts = [
        ("Overleaf source", "Self-contained TeX, TikZ diagrams, literature survey inputs, bibliography"),
        ("Data and code", "Experiments, notebooks, scripts, logs, tests, and result JSON/CSV files"),
        ("Final report", "76-page compiled PDF"),
        ("PPT", "This final presentation deck"),
    ]
    for i, (head, body) in enumerate(artifacts):
        x = 0.9 + (i % 2) * 6.15
        y = 1.55 + (i // 2) * 2.15
        card(slide, x, y, 5.45, 1.55)
        text(slide, x + 0.25, y + 0.22, 4.95, 0.32, head, 18, TEAL if i != 2 else AMBER, True)
        text(slide, x + 0.25, y + 0.72, 4.95, 0.45, body, 13, MUTED)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text(slide, 1.0, 2.25, 11.3, 0.8, "Thank You", 44, WHITE, True, PP_ALIGN.CENTER)
    text(slide, 1.0, 3.25, 11.3, 0.5, "Questions and Discussion", 24, TEAL, True, PP_ALIGN.CENTER)
    text(slide, 1.0, 4.25, 11.3, 0.35, "Group 6 | Reinforcement Learning for Agentic LLM Fine-Tuning", 15, MUTED, False, PP_ALIGN.CENTER)

    return prs


if __name__ == "__main__":
    deck = add_deck()
    output = "final_capstone_presentation.pptx"
    deck.save(output)
    print(output)
