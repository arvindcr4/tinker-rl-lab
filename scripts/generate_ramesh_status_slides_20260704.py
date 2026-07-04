#!/usr/bin/env python3
"""Generate an 8-slide, ZVF-focused, plain-language status deck for Ramesh (4 July 2026).

Live experiment statuses on slide 5 are read from run logs at generation
time, so re-running this script right before the meeting refreshes them.
"""

import re
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent

BG = RGBColor(0x12, 0x12, 0x20)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
GREEN = RGBColor(0x00, 0xFF, 0x88)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)
RED = RGBColor(0xFF, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)
DARK_CARD = RGBColor(0x1E, 0x1E, 0x32)
LIGHT_CARD = RGBColor(0x26, 0x26, 0x40)


# ---------------------------------------------------------------- live status

def _read(path):
    try:
        return Path(path).read_text(errors="replace")
    except OSError:
        return ""


def n2_status():
    txt = _read(ROOT / "experiments/tinker-runs/logs/n2_reward_tensor_resume_20260704.out")
    prog = {}
    for arm, step, total in re.findall(r"\[n2:(\w+)\]\s+(\d+)/(\d+)", txt):
        prog[arm] = "done" if int(step) >= int(total) else f"{step}/{total}"
    if not prog:
        return "relaunched · logging"
    return " · ".join(f"{a.upper()} {prog[a]}" for a in ("grpo", "aero", "gift", "areal") if a in prog)


def n10_status():
    txt = _read(ROOT / "experiments/tinker-runs/logs/n10_gsm8k_cot_seed_expansion_20260704.out")
    runs = {}
    for run, step, total in re.findall(r"\[(n10_\w+?_s\d+)\]\s+(\d+)/(\d+)", txt):
        runs[run] = (int(step), int(total))
    if not runs:
        return "relaunched · logging"
    done = sum(1 for s, t in runs.values() if s >= t)
    live = [f"{r.split('_', 1)[1]} at {s}/{t}" for r, (s, t) in runs.items() if s < t]
    cur = f" · {live[-1]}" if live else ""
    return f"{done}/16 runs done{cur}"


def mega_status():
    txt = _read(ROOT / "experiments/results/mega_20260704/cells_done.jsonl")
    n = sum(1 for line in txt.splitlines() if line.strip())
    return f"506 cells planned · {n} done since relaunch"


STAMP = datetime.now().strftime("%H:%M, %d %b").lstrip("0")

# ---------------------------------------------------------------- deck helpers

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tb


def add_card(slide, left, top, width, height, fill=DARK_CARD):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def slide_header(slide, num, title):
    add_text(slide, 0.5, 0.35, 0.6, 0.45, f"{num:02d}", 18, CYAN, True)
    add_text(slide, 1.1, 0.3, 11.5, 0.6, title, 34, WHITE, True)
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.02))


def footer(slide, num):
    add_text(slide, 0.5, 7.05, 12.3, 0.3, f"ZVF research status · for Ramesh · 4 July 2026 · {num}/8", 11, GRAY)


# --- Slide 1: Title ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 1, 1.2, 11, 1, "Zero-Variance Fraction (ZVF)", 52, CYAN, True, PP_ALIGN.CENTER)
add_text(slide, 1, 2.35, 11, 0.9, "A warning light for RL training of language models", 26, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, 1, 3.35, 11, 0.6, "Status & Next Steps", 22, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, 1, 4.25, 11, 0.5, "Arvind C R  ·  4 July 2026", 18, WHITE, False, PP_ALIGN.CENTER)
add_card(slide, 4.5, 5.4, 4.3, 0.9, LIGHT_CARD)
add_text(slide, 4.5, 5.55, 4.3, 0.6, "Prepared for Ramesh", 16, CYAN, True, PP_ALIGN.CENTER)
footer(slide, 1)

# --- Slide 2: Executive Summary ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_header(slide, 2, "Executive Summary")

add_card(slide, 0.5, 1.2, 12.3, 1.05, LIGHT_CARD)
add_text(slide, 0.7, 1.33, 11.9, 0.8,
         "GRPO teaches a model by comparing several answers to the same question. ZVF measures how often all those answers get the same score — when that happens, there is nothing to compare, and learning silently stops.",
         15, WHITE, True, PP_ALIGN.CENTER)

wins = [
    ("500+", "training runs completed", "844 run IDs in our Tinker inventory; the papers draw on 70+ carefully curated runs.", CYAN),
    ("0.93", "collapse early-warning score", "Our ZVF-based warning score ranks a failing run above a healthy one 93% of the time (AUROC 0.93).", GREEN),
    ("8", "papers, all building to PDF", "One 4-paper foundation pillar, three method pillars, and one auxiliary study.", YELLOW),
    ("3", "experiments running now", "N2, N10, and a 506-cell sweep are adding evidence as we speak.", CYAN),
]
for i, (val, label, desc, color) in enumerate(wins):
    x = 0.5 + (i % 2) * 6.3
    y = 2.45 + (i // 2) * 2.25
    add_card(slide, x, y, 6.0, 2.0)
    add_text(slide, x + 0.2, y + 0.1, 2.5, 0.8, val, 32, color, True)
    add_text(slide, x + 2.7, y + 0.15, 3.3, 0.75, label, 14, WHITE, True)
    add_text(slide, x + 0.2, y + 1.0, 5.6, 0.9, desc, 12, GRAY)

add_text(slide, 0.5, 6.9, 12.3, 0.4, "ZVF has grown from a warning light into a research program — measure it, predict trouble with it, act on it — aimed at ICLR 2027.", 14, GREEN, True, PP_ALIGN.CENTER)
footer(slide, 2)

# --- Slide 3: The ZVF Story ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_header(slide, 3, "The ZVF Story: Measure → Predict → Act")

acts = [
    ("1 · Measure", CYAN,
     "Each training step we ask: in how many question-groups did every answer get the same score?\n\nThat fraction is ZVF. Same score everywhere = nothing to compare = zero learning signal.\n\nIt is cheap — computed from scores we already collect. It even explained why our tool-use task scored 0%: ZVF was pinned at 100%, so the model never saw a learnable signal."),
    ("2 · Predict", GREEN,
     "A warning score built from ZVF flags a collapsing run before its reward curve visibly breaks.\n\nIt ranks a failing run above a healthy one 93% of the time (AUROC 0.93).\n\nAdding extra signals, like answer length, did not improve it — the simple version already does the job."),
    ("3 · Act", YELLOW,
     "Paper 7 turns the light into a controller.\n\nWhen ZVF climbs, it automatically changes how many answers we sample per question, restoring the contrast GRPO needs.\n\nThat closes the loop: from spotting failure to preventing it."),
]
for i, (title, color, body) in enumerate(acts):
    x = 0.5 + i * 4.2
    add_card(slide, x, 1.2, 4.0, 5.1)
    add_text(slide, x + 0.2, 1.35, 3.6, 0.45, title, 17, color, True)
    add_text(slide, x + 0.2, 2.0, 3.6, 4.1, body, 12, GRAY)

add_text(slide, 0.5, 6.5, 12.3, 0.5, "One signal, one arc: from a per-step statistic to an automatic intervention.", 14, GREEN, True, PP_ALIGN.CENTER)
footer(slide, 3)

# --- Slide 4: Paper Portfolio ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_header(slide, 4, "Paper Portfolio: 4 Pillars + 1 Auxiliary")

add_card(slide, 0.5, 1.2, 12.3, 2.0)
add_text(slide, 0.7, 1.3, 11.9, 0.4, "Pillar I · Foundations (papers 1–4)", 16, CYAN, True)
found = [
    "P1 · Scaling — bigger models don't automatically learn more from RL (0.6B–671B); what matters is whether the model could nearly do the task already",
    "P2 · ZVF — the warning light itself: measures when learning silently stops",
    "P3 · Group size — how many answers per question? No single best number; also tests the claim that GRPO secretly equals DPO",
    "P4 · Length bias — RL training did not teach models to ramble on short math (8/8 checks agree, p = 0.0039)",
]
for j, line in enumerate(found):
    add_text(slide, 0.7 + (j % 2) * 6.1, 1.75 + (j // 2) * 0.62, 5.9, 0.6, line, 10.5, GRAY)

pillars = [
    ("Pillar II · P5", "Report the Stack, Not the Label", "The same algorithm name behaves differently in different libraries — so papers should report the whole software stack, not just the label.", GREEN),
    ("Pillar III · P6", "GRPO-Registry", "A machine-readable catalog of what each RL library actually implements, and how the variants differ from each other.", YELLOW),
    ("Pillar IV · P7", "From Diagnostic to Controller", "A theory of why GRPO runs out of learning signal, plus an automatic group-size adjustment driven by ZVF.", RED),
]
for i, (tag, title, body, color) in enumerate(pillars):
    x = 0.5 + i * 4.2
    add_card(slide, x, 3.45, 4.0, 2.5)
    add_text(slide, x + 0.2, 3.58, 3.6, 0.35, tag, 13, color, True)
    add_text(slide, x + 0.2, 3.98, 3.6, 0.6, title, 14, WHITE, True)
    add_text(slide, x + 0.2, 4.62, 3.6, 1.25, body, 11, GRAY)

add_card(slide, 0.5, 6.15, 12.3, 0.75, LIGHT_CARD)
add_text(slide, 0.7, 6.28, 11.9, 0.5,
         "Auxiliary · P8 — LLM vs. XGBoost in credit-card fraud: the language model is best at sensing and explaining; XGBoost keeps the scoring job.",
         12, GRAY)

add_text(slide, 0.5, 6.92, 12.3, 0.35, "All 8 papers compile to PDF today; ZVF threads through Pillars I and IV.", 13, WHITE, False, PP_ALIGN.CENTER)
footer(slide, 4)

# --- Slide 5: Live Tinker Campaigns ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_header(slide, 5, "Running Right Now — Three Tinker Campaigns")

exps = [
    ("N2 · The full gradebook", "Records every score for every answer — not just averages — across four GRPO-family methods, so ZVF is computed exactly.", n2_status(), CYAN),
    ("N10 · Eight seeds instead of three", "Re-runs the length-bias experiment from 8 random starting points, tracking ZVF each step. Same answer 8 times is hard to call luck.", n10_status(), GREEN),
    ("Mega · The coverage map", "506 cells of model × task × settings, mapping where a learning signal exists at all. Streams to the zvf-audit-v2 dashboard.", mega_status(), YELLOW),
]
for i, (name, desc, status, color) in enumerate(exps):
    y = 1.35 + i * 1.7
    add_card(slide, 0.5, y, 12.3, 1.45)
    add_text(slide, 0.7, y + 0.15, 5.0, 0.35, name, 15, color, True)
    add_text(slide, 0.7, y + 0.62, 7.2, 0.7, desc, 12, GRAY)
    add_text(slide, 8.1, y + 0.62, 4.4, 0.6, status, 12, color, False, PP_ALIGN.RIGHT)

add_text(slide, 0.5, 6.75, 12.3, 0.4,
         f"Every run saves partial results continuously, so nothing is lost if something stops. Status auto-read from run logs at {STAMP}.",
         12, GRAY)
footer(slide, 5)

# --- Slide 6: From Runs to Papers ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_header(slide, 6, "From Runs to Papers")

maps = [
    ("N2 → P2 & P7", CYAN,
     "The full gradebook gives the warning-light paper exact ZVF numbers instead of estimates.\n\nIt also gives the controller paper the data it needs to tune when and how to act."),
    ("N10 → P4 & P2", GREEN,
     "Eight seeds turn “we saw it once” into “we see it every time, with error bars.”\n\nThat was the single biggest weakness reviewers would have flagged in the length-bias paper."),
    ("Mega → P1, P3 & P5", YELLOW,
     "The coverage map backs the scaling and group-size stories with breadth.\n\nEvery cell also writes a standard report card — live worked examples for the report-the-stack paper."),
]
for i, (title, color, body) in enumerate(maps):
    x = 0.5 + i * 4.2
    add_card(slide, x, 1.2, 4.0, 4.9)
    add_text(slide, x + 0.2, 1.35, 3.6, 0.45, title, 17, color, True)
    add_text(slide, x + 0.2, 2.0, 3.6, 3.9, body, 12, GRAY)

add_text(slide, 0.5, 6.4, 12.3, 0.5, "Each campaign lands directly in the sections reviewers are most likely to probe.", 14, GREEN, True, PP_ALIGN.CENTER)
footer(slide, 6)

# --- Slide 7: Rigor, Reproducibility & Honest Limits ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_header(slide, 7, "Keeping Ourselves Honest")

add_card(slide, 0.5, 1.2, 6.0, 5.2)
add_text(slide, 0.7, 1.35, 5.6, 0.4, "Rigor in place", 16, GREEN, True)
add_text(slide, 0.7, 1.9, 5.6, 4.0,
         "• Every headline claim is double-checked against the data files on disk; mismatches get corrected\n\n"
         "• A citation check flagged 2 references for correction; fixes tracked before submission\n\n"
         "• All 8 papers build to PDF with zero errors; three still show reference warnings that are being fixed\n\n"
         "• Anyone can re-run the work: Docker image, pinned seeds, W&B logs, checksums, step-by-step REPRODUCE.md\n\n"
         "• Submission package list and reviewer README are ready; final PDFs get generated before submission",
         11, GRAY)

add_card(slide, 6.8, 1.2, 6.0, 5.2)
add_text(slide, 7.0, 1.35, 5.6, 0.4, "Known limits — said out loud, with fixes", 16, YELLOW, True)
add_text(slide, 7.0, 1.9, 5.6, 4.0,
         "• Most headline results used one seed → N10 is expanding to eight right now\n\n• Training runs are short by design → presented as early-training snapshots, not final outcomes\n\n• One health metric (KL) was missing mid-campaign → disclosed; stand-ins used where possible\n\n• Tool-use scored 0% → explained by ZVF: no learnable signal, and we say so\n\n• Headline scores use training questions → held-out checks reported where we have them",
         11, GRAY)

add_text(slide, 0.5, 6.6, 12.3, 0.5, "We lead with limitations rather than hiding them — reviewers consistently reward this.", 14, GREEN, True, PP_ALIGN.CENTER)
footer(slide, 7)

# --- Slide 8: Next Steps & Decisions ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_header(slide, 8, "Next Steps & Decisions")

add_card(slide, 0.5, 1.2, 6.0, 2.6)
add_text(slide, 0.7, 1.35, 5.6, 0.4, "This week", 16, CYAN, True)
add_text(slide, 0.7, 1.85, 5.6, 1.7,
         "• Let the three campaigns finish\n• Fold their ZVF results into the papers as they land\n• Double-check numbers in anything we touch\n• Lock the ICLR 2027 plan — deadlines expected ~mid-Sept 2026 (last cycle: abstract 19 Sept, paper 24 Sept)",
         12, GRAY)

add_card(slide, 6.8, 1.2, 6.0, 2.6)
add_text(slide, 7.0, 1.35, 5.6, 0.4, "Next 2–4 weeks", 16, GREEN, True)
add_text(slide, 7.0, 1.85, 5.6, 1.7,
         "• Error bars across the whole portfolio\n• Connect to the wider post-training literature (Tulu-3, DPO, PPO)\n• Strengthen the tool-use evaluation\n• Finish reference cleanup so every paper is submission-clean",
         12, GRAY)

add_card(slide, 0.5, 4.0, 12.3, 2.6, LIGHT_CARD)
add_text(slide, 0.7, 4.15, 11.9, 0.4, "Questions for you", 16, YELLOW, True)
add_text(slide, 0.7, 4.65, 11.9, 1.7,
         "1. For ICLR 2027: one unified ZVF-led paper, or ZVF (P2 + P7) leading with companion submissions?\n2. Any concerns to address before we lock the ICLR 2027 timeline?",
         13, GRAY)

add_text(slide, 0.5, 6.8, 12.3, 0.5, "Target: ICLR 2027 (deadlines expected ~mid-Sept 2026; CFP not yet posted). Fallbacks: ICML 2027 D&B (~Jan 2027), NeurIPS 2027 E&D (~May 2027).", 11, GRAY, False, PP_ALIGN.RIGHT)
footer(slide, 8)

out = ROOT / "reports/ZVF_status_Ramesh_2026-07-04.pptx"
prs.save(str(out))
print(f"Saved {out} ({len(prs.slides)} slides)")
