import urllib.request
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

DECK = Path("/Users/arvind/Developer/agentic_repos/tinker-rl-lab/outputs/PESU_MTech_Phase1_ZVF_Defense_ArvindCR.pptx")
ASSETS = Path("/Users/arvind/Developer/agentic_repos/tinker-rl-lab/outputs/deck_assets")

TEAL = RGBColor(0x2D, 0xD4, 0xBF)

# Free-to-use Pexels photos (Pexels License: free, no attribution required)
SOURCES = {
    "title_cookoff.jpg": "https://images.pexels.com/photos/37138037/pexels-photo-37138037/free-photo-of-culinary-students-in-a-professional-kitchen.jpeg?auto=compress&cs=tinysrgb&w=1400",
    "title_jumpshot.jpg": "https://images.pexels.com/photos/29479238/pexels-photo-29479238/free-photo-of-intense-indoor-basketball-game-action-shot.jpeg?auto=compress&cs=tinysrgb&w=1400",
}

# Tile geometry (inches) — lower-left band under the subtitle, clear of the author panel (L>=8.96)
TILE_ASPECT = 4.04 / 1.42  # width / height
LAYOUT = [
    # key, left, top, width, height, vertical_focus (0=top .. 1=bottom)
    ("title_cookoff.jpg",  0.52, 5.52, 4.04, 1.42, 0.50),
    ("title_jumpshot.jpg", 4.80, 5.52, 4.04, 1.42, 0.42),
]

def fetch(name, url):
    dst = ASSETS / name
    req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
    with urllib.request.urlopen(req, timeout=60) as r:
        dst.write_bytes(r.read())
    return dst

def crop_to_aspect(path, aspect, vfocus):
    im = Image.open(path).convert("RGB")
    w, h = im.size
    if w / h > aspect:                      # too wide -> trim width, keep full height
        new_w = int(round(h * aspect))
        x0 = (w - new_w) // 2
        box = (x0, 0, x0 + new_w, h)
    else:                                   # too tall -> trim height around vfocus
        new_h = int(round(w / aspect))
        y0 = int(round((h - new_h) * vfocus))
        box = (0, y0, w, y0 + new_h)
    im.crop(box).save(path, quality=88)

if __name__ == "__main__":
    # download + crop
    for name, l, t, w, h, vf in LAYOUT:
        p = fetch(name, SOURCES[name])
        crop_to_aspect(p, TILE_ASPECT, vf)
        print("processed", name, Image.open(p).size)

    prs = Presentation(str(DECK))
    slide = prs.slides[0]
    for name, l, t, w, h, vf in LAYOUT:
        pic = slide.shapes.add_picture(str(ASSETS / name), Inches(l), Inches(t), Inches(w), Inches(h))
        pic.line.color.rgb = TEAL
        pic.line.width = Pt(1.25)
        pic.shadow.inherit = False
    prs.save(str(DECK))
    print("saved", DECK)
