"""ZVF Program — 6-slide LIGHTNING progress + next-steps deck (2026-06-14).
Author: Arvind. Same dark theme as the 12-slide deck, condensed to the highest-signal cut.
Nothing fabricated.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG     = RGBColor(0x0E, 0x11, 0x16); PANEL = RGBColor(0x16, 0x1B, 0x22)
PANEL2 = RGBColor(0x1C, 0x23, 0x2C); INK   = RGBColor(0xE6, 0xED, 0xF3)
MUTE   = RGBColor(0x8B, 0x94, 0x9E); TEAL  = RGBColor(0x2D, 0xD4, 0xBF)
AMBER  = RGBColor(0xF5, 0xA6, 0x23); GREEN = RGBColor(0x3F, 0xB9, 0x50)
YELLOW = RGBColor(0xD2, 0x99, 0x22); RED   = RGBColor(0xF8, 0x51, 0x49)
LINECLR= RGBColor(0x30, 0x36, 0x3D); FONT  = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def spaced(s): return "  ".join(list(s))


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    return s


def box(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tf


def para(tf, text, size, color=INK, bold=False, font=FONT, align=PP_ALIGN.LEFT,
         first=False, space_before=0, space_after=4, italic=False, spacing=1.0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_before = Pt(space_before); p.space_after = Pt(space_after)
    p.line_spacing = spacing
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.name = font
    f.color.rgb = color
    return p


def rect(s, l, t, w, h, fill=PANEL, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def pill(s, l, t, text, color):
    w = 0.16 + 0.082 * len(text)
    sp = rect(s, l, t, w, 0.32, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try: sp.adjustments[0] = 0.5
    except Exception: pass
    tf = sp.text_frame; tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(10); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = BG
    return w


def chrome(s, eyebrow, n):
    rect(s, 0.55, 0.5, 0.12, 0.12, fill=TEAL, shape=MSO_SHAPE.OVAL)
    para(box(s, 0.78, 0.46, 9.5, 0.4), spaced(eyebrow), 10.5, MUTE, bold=True, first=True)
    para(box(s, 11.2, 0.46, 1.6, 0.4), f"{n:02d} / 06", 10.5, MUTE, bold=True,
         align=PP_ALIGN.RIGHT, first=True)
    rect(s, 0.55, 6.92, 12.23, 0.018, fill=LINECLR)
    para(box(s, 0.55, 7.0, 12.23, 0.35),
         "ZVF Program  ·  Progress & Next Steps  ·  2026-06-14        agentic_grpo_bench · Arvind",
         9, MUTE, first=True)


def header(s, eyebrow, title, n, subtitle=None):
    chrome(s, eyebrow, n)
    tf = box(s, 0.55, 0.95, 12.2, 1.1)
    para(tf, title, 29, INK, bold=True, first=True)
    if subtitle: para(tf, subtitle, 13.5, MUTE, space_before=2)
    rect(s, 0.57, 1.92, 1.5, 0.05, fill=TEAL)


# ============================================================ 1 · TITLE
s = slide()
rect(s, 0, 0, 13.333, 0.09, fill=TEAL)
para(box(s, 0.9, 0.62, 11, 0.4), spaced("GRPO AUDIT · NEXT MOVES"), 11, MUTE, bold=True, first=True)
tf = box(s, 0.9, 2.3, 11.5, 2.2)
para(tf, "The ZVF Program", 52, INK, bold=True, first=True)
para(tf, "Progress & Next Steps — Lightning Review", 22, TEAL, bold=True, space_before=6)
para(tf, "From an audit diagnostic to the standard telemetry, theory, and tooling for RL post-training.",
     15, MUTE, space_before=12)
rect(s, 0.92, 5.3, 2.0, 0.05, fill=AMBER)
tf = box(s, 0.9, 5.5, 11.5, 1.2)
para(tf, "One method · one theorem · one library · one community position", 14, INK, bold=True, first=True)
para(tf, "June 14, 2026     ·     Arvind     ·     agentic_grpo_bench     ·     NeurIPS 2026 track",
     12, MUTE, space_before=8)

# ============================================================ 2 · SCORECARD
s = slide()
header(s, "STATUS", "Where the four pillars stand — 3 of 4 now ship artifacts", 2)
rows = [
    ("PILLAR", "ARTIFACT ON DISK", "STATUS"),
    ("3 · OSS library", "zvf-triage package — 71/71 tests, builds wheel", "SHIPPABLE"),
    ("1 · Method / M1", "sweep/ harness — 403-cell dry-run, cell_runner wired", "LAUNCH-READY"),
    ("2 · Theory", "zvf_theory.tex compiles + THEORY_NOTES.md", "DRAFT · SKETCH"),
    ("4 · Position", "min_report_rl.tex compiles + CHECKLIST.md", "DRAFT · TODO"),
]
clr = {"SHIPPABLE": GREEN, "LAUNCH-READY": TEAL, "DRAFT · SKETCH": YELLOW, "DRAFT · TODO": AMBER}
top = 2.35; rh = 0.86; left = 0.55; widths = [3.0, 6.3, 2.93]
for ri, row in enumerate(rows):
    x = left; head = ri == 0
    for ci, cell in enumerate(row):
        fill = PANEL2 if head else (PANEL if ri % 2 else BG)
        rect(s, x, top + ri * rh, widths[ci], rh, fill=fill, line=LINECLR, line_w=0.75)
        if ci == 2 and not head:
            pill(s, x + 0.2, top + ri * rh + (rh - 0.32) / 2, cell, clr[cell])
        else:
            tf = box(s, x + 0.2, top + ri * rh, widths[ci] - 0.36, rh, anchor=MSO_ANCHOR.MIDDLE)
            para(tf, cell, 12 if head else 12.5, MUTE if head else INK, bold=head or ci == 0, first=True)
        x += widths[ci]
para(box(s, 0.55, top + len(rows) * rh + 0.22, 12.2, 0.5),
     "Predecessor v1 audit (EAI Endorsed Transactions, ~95 runs, 97/100) is submission-ready — the corpus the follow-ups build on.",
     11.5, MUTE, italic=True, first=True)

# ============================================================ 3 · SHIPPED
s = slide()
header(s, "WHAT SHIPPED TODAY", "The library + the harness — the two artifacts the rest feed on", 3)
# left: zvf-triage
rect(s, 0.55, 2.3, 6.0, 3.95, fill=PANEL); rect(s, 0.55, 2.3, 6.0, 0.09, fill=GREEN)
tf = box(s, 0.8, 2.55, 5.55, 3.6)
para(tf, "zvf-triage  ·  Pillar 3 (OSS)", 15, GREEN, bold=True, first=True)
for t in ["Installable package, 71/71 tests pass, builds a wheel",
          "ZVFController state machine + ZVFCallback (lazy TRL adapter)",
          "Pure-numpy ZVF/GU core — reused from your formalization",
          "Apache-2.0 · CI py3.9–3.12 · example · adapter stubs"]:
    para(tf, "•  " + t, 12, INK, space_before=8, spacing=1.03)
para(tf, "Next: TestPyPI → upstream PR to TRL.", 11.5, AMBER, bold=True, space_before=12)
# right: sweep
rect(s, 6.78, 2.3, 6.0, 3.95, fill=PANEL); rect(s, 6.78, 2.3, 6.0, 0.09, fill=TEAL)
tf = box(s, 7.03, 2.55, 5.55, 3.6)
para(tf, "sweep harness  ·  Pillar 1 / M1", 15, TEAL, bold=True, first=True)
for t in ["403-cell grid (385 audit + 18 matched-compute)",
          "cell_runner wired into live_zvf_probe — no training faked",
          "GRPO vs DAPO vs GSPO at matched rollouts/tokens",
          "Dry-run verified: enumerates 403, launches nothing"]:
    para(tf, "•  " + t, 12, INK, space_before=8, spacing=1.03)
para(tf, "Needs you: TINKER_API_KEY → run_sweep.py --execute.", 11.5, AMBER, bold=True, space_before=12)

# ============================================================ 4 · DRAFTED
s = slide()
header(s, "WHAT'S DRAFTED", "Theory + position — real drafts, hard parts honestly fenced", 4)
# left: theory
rect(s, 0.55, 2.3, 6.0, 3.95, fill=PANEL); rect(s, 0.55, 2.3, 6.0, 0.09, fill=AMBER)
tf = box(s, 0.8, 2.55, 5.55, 3.6)
para(tf, "Theory  ·  Pillar 2", 15, AMBER, bold=True, first=True)
para(tf, "Compiles. PROOF-SKETCH ONLY — gaps flagged:", 11.5, MUTE, space_before=4)
for t in ["T1 estimator/CI — degenerates at ZVF→1",
          "T2 lower bound — proves nonzero-grad, not improvement",
          "T3 optimal G* — signal function unjustified"]:
    para(tf, "•  " + t, 12, INK, space_before=7, spacing=1.03)
para(tf, "16-item proof ledger in THEORY_NOTES.md. You verify the math.", 11, RED, bold=True, space_before=10, spacing=1.03)
# right: position
rect(s, 6.78, 2.3, 6.0, 3.95, fill=PANEL); rect(s, 6.78, 2.3, 6.0, 0.09, fill=AMBER)
tf = box(s, 7.03, 2.55, 5.55, 3.6)
para(tf, "Position  ·  Pillar 4", 15, AMBER, bold=True, first=True)
para(tf, "MIN-REPORT-RL — every GRPO paper reports:", 11.5, MUTE, space_before=4)
for t in ["Loss form · ref policy + KL · sampler/backend",
          "Per-step ZVF/GU · group-size schedule",
          "Held-out split · decontamination probe"]:
    para(tf, "✓  " + t, 12, INK, space_before=7, spacing=1.03)
para(tf, "+ reproducibility audit of DAPO/GSPO/Dr.GRPO/MAD-GRPO. Results + cites = TODO.",
     11, RED, bold=True, space_before=10, spacing=1.03)

# ============================================================ 5 · BOUNDARY + DECISIONS
s = slide()
header(s, "INTEGRITY · YOUR CALL", "What's yours to finish — and 3 decisions", 5)
# needs-you (left)
rect(s, 0.55, 2.3, 5.6, 3.95, fill=PANEL); rect(s, 0.55, 2.3, 5.6, 0.09, fill=RED)
tf = box(s, 0.8, 2.55, 5.15, 3.6)
para(tf, "Needs you — can't be faked", 14, RED, bold=True, first=True)
for t in ["The 403 GPU runs — your Tinker/Modal compute",
          "The theorem proofs — verify T1/T2/T3",
          "The audit table — fills from sweep output",
          "Real GitHub URL + PyPI account to publish"]:
    para(tf, "•  " + t, 12, INK, space_before=9, spacing=1.03)
para(tf, "Nothing was fabricated.", 11.5, TEAL, bold=True, space_before=12)
# decisions (right)
rect(s, 6.38, 2.3, 6.4, 3.95, fill=PANEL); rect(s, 6.38, 2.3, 6.4, 0.09, fill=AMBER)
tf = box(s, 6.63, 2.55, 5.95, 3.6)
para(tf, "3 decisions to confirm", 14, AMBER, bold=True, first=True)
decs = [
    ("1 · Sweep runner", "live_zvf_probe (only one logging ZVF/GU) — OK?", TEAL),
    ("2 · DAPO/GSPO fidelity", "arms are surrogates, not canonical — tighten?", RED),
    ("3 · License", "Apache-2.0 (matches TRL/verl) — keep or MIT?", TEAL),
]
for name, body, acc in decs:
    para(tf, name, 12.5, acc, bold=True, space_before=10)
    para(tf, body, 11.5, MUTE, space_before=1, spacing=1.02)

# ============================================================ 6 · NEXT STEPS
s = slide()
header(s, "ROADMAP · NEXT STEPS", "From shipped infrastructure to four papers", 6)
para(box(s, 0.55, 1.95, 12.2, 0.4),
     "This week:  5-cell Tinker smoke-run via cell_runner · publish zvf-triage to TestPyPI · lock DAPO/GSPO configs.",
     12.5, TEAL, bold=True, first=True)
steps = [
    ("M1 · wk 1–2", "Strengthen the audit", "Full 403-run sweep + held-out suite", "v2 paper + zvf-triage v0.1"),
    ("M2 · wk 2–3", "Ship the method", "Adaptive controller; matched-compute wins; PRs", "ICLR '27 method paper"),
    ("M3 · wk 2–3", "Land the theory", "Prove T1 CI; fix T2; derive G*", "AISTATS / COLT paper"),
    ("M4 · wk 3–4", "Set the standard", "Fill audit table; MIN-REPORT-RL position", "NeurIPS Position + MLRC"),
]
y = 2.5
for when, title, body, out in steps:
    rect(s, 0.55, y, 12.23, 0.82, fill=PANEL); rect(s, 0.55, y, 0.07, 0.82, fill=TEAL)
    tf = box(s, 0.8, y, 2.4, 0.82, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, when, 12.5, TEAL, bold=True, first=True)
    tf = box(s, 3.2, y, 2.7, 0.82, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, title, 12.5, INK, bold=True, first=True)
    tf = box(s, 5.9, y, 4.0, 0.82, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, body, 11, MUTE, first=True, spacing=1.0)
    tf = box(s, 9.95, y, 2.75, 0.82, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "→ " + out, 11, AMBER, bold=True, first=True, spacing=1.0)
    y += 0.9
rect(s, 0.55, 6.18, 12.23, 0.58, fill=PANEL2)
para(box(s, 0.8, 6.18, 11.8, 0.58, anchor=MSO_ANCHOR.MIDDLE),
     "Defensible · Compounding · Ownable — today the library + harness shipped; the rest feeds on them.",
     13.5, TEAL, bold=True, align=PP_ALIGN.CENTER, first=True)

OUT = "/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14_lightning.pptx"
prs.save(OUT)
print("saved:", OUT)
print("slides:", len(prs.slides._sldIdLst))
