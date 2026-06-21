import os
import sys
from pptx import Presentation
from groq import Groq

def enrich_text(client, original_text):
    # Skip very short text or typical footer/header text like "10 / 12"
    if len(original_text.strip()) < 15 or ' / ' in original_text or 'ZVF Program' in original_text:
        return original_text
    
    prompt = f"Enhance and professionally enrich the following presentation slide text. Keep it concise, punchy, and suitable for a presentation. Return ONLY the enriched text, no quotes or additional commentary:\n\n{original_text}"
    
    try:
        chat_completion = client.chat.completions.create(
            messages=[
                {"role": "system", "content": "You are a helpful assistant expert in presentation copywriting."},
                {"role": "user", "content": prompt}
            ],
            model="kimi-k2-0905-preview",
        )
        return chat_completion.choices[0].message.content.strip()
    except Exception as e:
        print(f"Error calling Groq API: {e}")
        return original_text

def replace_paragraph_text(paragraph, new_text):
    if not paragraph.runs:
        paragraph.text = new_text
        return
    # Preserve formatting of the first run and put all text there
    first_run = paragraph.runs[0]
    first_run.text = new_text
    # Remove subsequent runs to avoid leftover text
    for i in range(len(paragraph.runs) - 1, 0, -1):
        p = paragraph._p
        p.remove(paragraph.runs[i]._r)

def enrich_slide(pptx_path, slide_index):
    if "GROQ_API_KEY" not in os.environ:
        print("Error: GROQ_API_KEY environment variable is not set.")
        sys.exit(1)
        
    client = Groq(api_key=os.environ["GROQ_API_KEY"])
    prs = Presentation(pptx_path)
    
    if slide_index >= len(prs.slides):
        print(f"Error: Slide {slide_index + 1} does not exist in {pptx_path}. It only has {len(prs.slides)} slides.")
        return

    slide = prs.slides[slide_index]
    
    print(f"Enriching slide {slide_index + 1} in {pptx_path}...")
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        for paragraph in shape.text_frame.paragraphs:
            original_text = paragraph.text.strip()
            if original_text:
                enriched = enrich_text(client, original_text)
                if enriched and enriched != original_text:
                    print(f"\nOriginal: {original_text}")
                    print(f"Enriched: {enriched}")
                    replace_paragraph_text(paragraph, enriched)

    out_path = pptx_path.replace(".pptx", "_enriched.pptx")
    prs.save(out_path)
    print(f"\nSaved enriched presentation to {out_path}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 enrich_slide.py <pptx_file> [slide_number_1_indexed]")
        sys.exit(1)
        
    pptx_path = sys.argv[1]
    # Default to slide 10 as requested
    slide_num = int(sys.argv[2]) if len(sys.argv) > 2 else 10
    
    # slide_num is 1-indexed, so we subtract 1 for 0-indexed.
    enrich_slide(pptx_path, slide_num - 1)
