import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def enrich_slide_15():
    file_path = '/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14_lightning.pptx'
    
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}")
        return

    prs = Presentation(file_path)
    print(f"Original slide count: {len(prs.slides)}")

    # Add blank slides if there are fewer than 15 slides
    # layout 6 is typically a blank slide
    blank_layout = prs.slide_layouts[6] 
    while len(prs.slides) < 15:
        s = prs.slides.add_slide(blank_layout)
        # Apply the dark theme background to match your other slides
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = RGBColor(0x0E, 0x11, 0x16)

    # Access Slide 15 (0-indexed, so 14)
    slide_15 = prs.slides[14]

    # Enrich Slide 15: Add a Text Box for a Title
    title_box = slide_15.shapes.add_textbox(Inches(0.55), Inches(0.95), Inches(12.2), Inches(1.1))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "ENRICHED SLIDE 15"
    title_p.font.bold = True
    title_p.font.size = Pt(30)
    title_p.font.name = "Arial"
    title_p.font.color.rgb = RGBColor(0xE6, 0xED, 0xF3) # INK color

    # Add an accent line under the title (TEAL)
    line = slide_15.shapes.add_shape(
        1, # MSO_SHAPE.RECTANGLE
        Inches(0.57), Inches(1.95), Inches(1.5), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x2D, 0xD4, 0xBF) # TEAL
    line.line.fill.background()
    line.shadow.inherit = False

    # Add a Content Text Box
    content_box = slide_15.shapes.add_textbox(Inches(0.55), Inches(2.5), Inches(10), Inches(4))
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    p1 = content_tf.paragraphs[0]
    p1.text = "• This slide was programmatically generated/enriched via python-pptx."
    p1.font.size = Pt(16)
    p1.font.color.rgb = RGBColor(0x8B, 0x94, 0x9E) # MUTE color
    p1.font.name = "Arial"

    p2 = content_tf.add_paragraph()
    p2.text = "• If you meant a different slide, simply change the slide index in this script."
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0x8B, 0x94, 0x9E)
    p2.font.name = "Arial"
    p2.space_before = Pt(10)

    # Save to a new file so we don't overwrite the original accidentally
    output_path = '/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14_lightning_enriched.pptx'
    prs.save(output_path)
    print(f"Presentation successfully saved to: {output_path}")
    print(f"New slide count: {len(prs.slides)}")

if __name__ == '__main__':
    enrich_slide_15()
