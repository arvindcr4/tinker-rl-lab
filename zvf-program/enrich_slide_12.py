import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def enrich_slide_12(pptx_path):
    print(f"Loading {pptx_path}...")
    prs = Presentation(pptx_path)
    
    if len(prs.slides) < 12:
        print(f"Error: Presentation only has {len(prs.slides)} slides. Slide 12 does not exist.")
        print("Note: The lightning deck only contains 6 slides. Please run this on the main deck instead.")
        sys.exit(1)
        
    slide = prs.slides[11]  # Slide 12 is at index 11
    
    # Add a Call-to-Action box at the bottom right
    print("Enriching Slide 12 with Call-to-Action box...")
    left = Inches(8.5)
    top = Inches(0.5)
    width = Inches(4.2)
    height = Inches(0.8)
    
    # Add shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    
    # Style the shape (Orange / Amber background)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xA6, 0x23)  # AMBER from original deck
    shape.line.color.rgb = RGBColor(0xE6, 0xED, 0xF3)       # INK from original deck
    shape.line.width = Pt(1.5)
    
    # Add text
    text_frame = shape.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "IMMEDIATE ACTION REQUIRED"
    run.font.bold = True
    run.font.size = Pt(12)
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x0E, 0x11, 0x16) # BG color
    
    p2 = text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Confirm 3 decisions to launch Tinker sweep."
    run2.font.bold = False
    run2.font.size = Pt(10)
    run2.font.name = "Arial"
    run2.font.color.rgb = RGBColor(0x0E, 0x11, 0x16) # BG color
    
    # Add a visual 'star' or 'badge' to the top left of the title
    star = slide.shapes.add_shape(
        MSO_SHAPE.STAR_5_POINT, Inches(0.55), Inches(0.95), Inches(0.5), Inches(0.5)
    )
    star.fill.solid()
    star.fill.fore_color.rgb = RGBColor(0x2D, 0xD4, 0xBF)  # TEAL
    star.line.fill.background()
    
    # Save the enriched presentation
    output_path = pptx_path.replace(".pptx", "_enriched.pptx")
    prs.save(output_path)
    print(f"Success! Enriched presentation saved to {output_path}")

if __name__ == "__main__":
    # Default to the main deck as the lightning deck has only 6 slides
    default_path = "/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14.pptx"
    
    path = sys.argv[1] if len(sys.argv) > 1 else default_path
    enrich_slide_12(path)
