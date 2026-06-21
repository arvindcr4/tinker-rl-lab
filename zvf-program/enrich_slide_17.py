import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def enrich_slide_17(filepath, output_filepath):
    print(f"Loading {filepath}...")
    prs = Presentation(filepath)
    
    if len(prs.slides) < 17:
        print(f"Error: The presentation only has {len(prs.slides)} slides.")
        print("Cannot enrich slide 17 because it does not exist.")
        return
        
    print("Found slide 17. Enriching...")
    # 0-indexed, so slide 17 is index 16
    slide = prs.slides[16]
    
    # Add a title if needed or an enrichment box
    left = Inches(0.55)
    top = Inches(6.0)
    width = Inches(12.2)
    height = Inches(1.0)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "Enriched Content for Slide 17"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0x2D, 0xD4, 0xBF) # TEAL
    
    p2 = text_frame.add_paragraph()
    p2.text = "This content was added programmatically via python-pptx."
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(0x8B, 0x94, 0x9E) # MUTE

    prs.save(output_filepath)
    print(f"Saved enriched presentation to {output_filepath}")

if __name__ == "__main__":
    input_file = "/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14_lightning.pptx"
    output_file = "/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14_lightning_slide17_enriched.pptx"
    enrich_slide_17(input_file, output_file)
