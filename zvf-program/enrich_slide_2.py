import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def enrich_slide():
    file_path = '/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14_lightning.pptx'
    out_path = '/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14_lightning_enriched.pptx'
    
    prs = Presentation(file_path)
    slide = prs.slides[1] # Slide 2
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
            
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                # Color code statuses
                if 'SHIPPABLE' in run.text:
                    run.font.color.rgb = RGBColor(0, 200, 0)
                    run.font.bold = True
                elif 'LAUNCH-READY' in run.text:
                    run.font.color.rgb = RGBColor(0, 150, 255)
                    run.font.bold = True
                elif 'DRAFT · SKETCH' in run.text:
                    run.font.color.rgb = RGBColor(255, 140, 0)
                    run.font.bold = True
                elif 'DRAFT · TODO' in run.text:
                    run.font.color.rgb = RGBColor(255, 0, 0)
                    run.font.bold = True
                    
    prs.save(out_path)
    print(f"Slide 2 enriched successfully! Saved to: {out_path}")

if __name__ == '__main__':
    enrich_slide()
