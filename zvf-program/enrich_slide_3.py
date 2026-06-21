from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

def add_pill(slide, left, top, width, height, text, bg_color, text_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    
    try:
        shape.adjustments[0] = 0.5
    except Exception:
        pass

    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.name = "Arial"
    r.font.color.rgb = text_color
    return shape

def enrich_slide_3():
    base_dir = "/Users/arvind/Developer/tinker-rl-lab/zvf-program"
    ppt_path = os.path.join(base_dir, "ZVF_Program_Progress_2026-06-14_lightning.pptx")
    prs = Presentation(ppt_path)
    
    # 0-indexed, so slides[2] is Slide 3
    slide = prs.slides[2]

    # Colors based on the existing theme
    GREEN = RGBColor(0x3F, 0xB9, 0x50)
    TEAL = RGBColor(0x2D, 0xD4, 0xBF)
    BG = RGBColor(0x0E, 0x11, 0x16)

    # Left panel bounds: l=0.55, t=2.3, w=6.0
    # Add 'SHIPPED' badge to the top right of the left panel
    add_pill(slide, left=5.35, top=2.5, width=1.0, height=0.3, text="SHIPPED", bg_color=GREEN, text_color=BG)

    # Right panel bounds: l=6.78, t=2.3, w=6.0
    # Add 'READY' badge to the top right of the right panel
    add_pill(slide, left=11.58, top=2.5, width=1.0, height=0.3, text="READY", bg_color=TEAL, text_color=BG)

    # Save to the original file to "enrich" it
    out_path = os.path.join(base_dir, "ZVF_Program_Progress_2026-06-14_lightning.pptx")
    prs.save(out_path)
    print(f"Slide 3 successfully enriched and saved to {out_path}")

if __name__ == "__main__":
    enrich_slide_3()
