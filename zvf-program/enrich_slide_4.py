import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def enrich_slide():
    file_path = '/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14_lightning.pptx'
    output_path = '/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14_lightning_enriched.pptx'
    
    print(f"Loading {file_path}...")
    prs = Presentation(file_path)
    
    # Slide 4 is at index 3
    slide = prs.slides[3]
    
    # 1. Highlight specific text like "TODO" or "PROOF-SKETCH ONLY"
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if 'TODO' in run.text:
                    run.font.color.rgb = RGBColor(255, 0, 0) # Red
                    run.font.bold = True
                if 'PROOF-SKETCH' in run.text:
                    run.font.color.rgb = RGBColor(255, 140, 0) # Dark Orange
                    run.font.bold = True

    # 2. Add an "Action Item" callout shape to the bottom right
    left = Inches(7.5)
    top = Inches(5.5)
    width = Inches(2.2)
    height = Inches(1.0)
    
    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(255, 243, 205) # Light yellow
    callout.line.color.rgb = RGBColor(133, 100, 4) # Dark brownish yellow
    
    tf = callout.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Critical Path:\nPrioritize Theory proofs & Audit table before NeurIPS!"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(133, 100, 4)
    p.font.bold = True
    
    prs.save(output_path)
    print(f"Enriched presentation saved to:\n{output_path}")

if __name__ == "__main__":
    enrich_slide()
