from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

def enrich_slide():
    prs = Presentation('/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14_lightning.pptx')
    slide = prs.slides[4] # Slide 5
    
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
            
        tf = shape.text_frame
        
        # Enrich the "Needs you" section
        if "Needs you" in tf.text:
            for paragraph in tf.paragraphs:
                # Highlight the header
                if "Needs you" in paragraph.text:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(220, 80, 0) # Deep Orange
                
                # Emphasize the final statement
                if "Nothing was fabricated." in paragraph.text:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.italic = True
                        run.font.color.rgb = RGBColor(0, 128, 0) # Green
                        
        # Enrich the "3 decisions" section
        if "3 decisions to confirm" in tf.text:
            for paragraph in tf.paragraphs:
                # Highlight the header
                if "3 decisions to confirm" in paragraph.text:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0, 102, 204) # Deep Blue
                
                # Highlight the decision titles
                elif "1 ·" in paragraph.text or "2 ·" in paragraph.text or "3 ·" in paragraph.text:
                    for run in paragraph.runs:
                        run.font.bold = True

    output_path = '/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14_lightning_enriched.pptx'
    prs.save(output_path)
    print(f"Slide 5 enriched successfully! Saved to:\n{output_path}")

if __name__ == '__main__':
    enrich_slide()
