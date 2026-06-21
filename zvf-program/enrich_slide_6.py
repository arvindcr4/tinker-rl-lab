import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

def enrich_slide():
    file_path = '/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14_lightning.pptx'
    out_path = '/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14_lightning_enriched.pptx'
    
    prs = Presentation(file_path)
    slide = prs.slides[5] # Slide 6
    
    # Define a subtle color palette for the 4 milestones
    # Colors: Blue, Green, Orange, Purple
    palette = [
        {"bg": RGBColor(240, 248, 255), "accent": RGBColor(70, 130, 180), "text": RGBColor(25, 25, 112)}, # M1 (Blue)
        {"bg": RGBColor(245, 255, 250), "accent": RGBColor(60, 179, 113), "text": RGBColor(0, 100, 0)},     # M2 (Green)
        {"bg": RGBColor(255, 250, 240), "accent": RGBColor(255, 140, 0), "text": RGBColor(139, 69, 0)},     # M3 (Orange)
        {"bg": RGBColor(248, 248, 255), "accent": RGBColor(147, 112, 219), "text": RGBColor(75, 0, 130)}    # M4 (Purple)
    ]
    
    # Shape indices for each row
    rows = [
        (8, 9, 10, 11, 12, 13),   # M1
        (14, 15, 16, 17, 18, 19), # M2
        (20, 21, 22, 23, 24, 25), # M3
        (26, 27, 28, 29, 30, 31)  # M4
    ]
    
    for i, row in enumerate(rows):
        colors = palette[i]
        
        bg_rect = slide.shapes[row[0]]
        accent_rect = slide.shapes[row[1]]
        
        # Add background colors
        bg_rect.fill.solid()
        bg_rect.fill.fore_color.rgb = colors["bg"]
        # Remove borders for a cleaner look
        bg_rect.line.fill.background()
        
        # Color the accent marker on the left
        accent_rect.fill.solid()
        accent_rect.fill.fore_color.rgb = colors["accent"]
        accent_rect.line.fill.background()
        
        # Bold milestone titles and apply text color
        m_title = slide.shapes[row[2]]
        if m_title.has_text_frame:
            for paragraph in m_title.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.color.rgb = colors["text"]
        
        # Bold the outcome descriptions
        outcome = slide.shapes[row[5]]
        if outcome.has_text_frame:
            for paragraph in outcome.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.color.rgb = colors["text"]
                    
    # Emphasize the "This week:" callout
    this_week = slide.shapes[7]
    if this_week.has_text_frame:
        text = this_week.text
        if text.startswith("This week:"):
            # Clear text frame to rebuild it with styled runs
            this_week.text_frame.clear()
            p = this_week.text_frame.paragraphs[0]
            
            # Bold and color "This week:"
            r1 = p.add_run()
            r1.text = "This week:"
            r1.font.bold = True
            r1.font.color.rgb = RGBColor(220, 20, 60) # Crimson
            
            # Add the rest of the text normally
            r2 = p.add_run()
            r2.text = text[len("This week:"):]
            
    # Add an accent color to the overall slide title bar (Shape 4: ZVF Program...)
    slide_title = slide.shapes[4]
    if slide_title.has_text_frame:
        for paragraph in slide_title.text_frame.paragraphs:
            for run in paragraph.runs:
                # Assuming light text on dark background or just make it standard dark if not set
                run.font.bold = True
            
    prs.save(out_path)
    print(f"Successfully enriched presentation saved to {out_path}")

if __name__ == '__main__':
    enrich_slide()
