import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def enrich_slide_8(filepath):
    # Load the presentation
    prs = Presentation(filepath)
    
    # Ensure there are at least 8 slides
    # If the presentation has fewer slides, append blank slides to reach 8
    while len(prs.slides) < 8:
        # Layout 5 is typically a blank slide in default templates, 
        # but fallback to layout 6 or 0 if out of bounds
        try:
            layout = prs.slide_layouts[6] 
        except IndexError:
            layout = prs.slide_layouts[0]
        prs.slides.add_slide(layout)
        
    # Access slide 8 (index 7)
    slide = prs.slides[7]
    
    # 1. Add an Enriched Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_tf = title_box.text_frame
    p = title_tf.paragraphs[0]
    p.text = "A P P E N D I X  ·  D E E P  D I V E"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # 2. Add a subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
    subtitle_tf = subtitle_box.text_frame
    p2 = subtitle_tf.paragraphs[0]
    p2.text = "Additional context on ZVF/GU Telemetry & Audit Harness"
    p2.font.size = Pt(20)
    p2.font.italic = True
    p2.font.color.rgb = RGBColor(100, 100, 100)
    
    # 3. Add a content box with bullet points summarizing the deep dive
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(8.5), Inches(4))
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    bullets = [
        "zvf-triage Package Details: Pure-numpy ZVF/GU core, validated with 71/71 tests.",
        "Audit Harness Scaling: Capable of expanding from 403-cell grid to thousands of runs using Tinker/Modal.",
        "Theory Validation: Needs rigorous peer review for T1 (CI) and T2 (lower bound) estimators.",
        "Reporting Standards: MIN-REPORT-RL checklist will standardize GRPO telemetry across publications."
    ]
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = content_tf.paragraphs[0]
        else:
            p = content_tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.level = 0
        p.space_after = Pt(14)
        
    # 4. Add a decorative action item shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(6.5), Inches(5.5), Inches(3), Inches(1.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 102, 204)
    text_frame = shape.text_frame
    text_frame.text = "Action Item:\nFinalize configs & launch!"
    for p in text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(16)
        p.font.bold = True
        
    # Save the enriched presentation
    output_path = filepath.replace(".pptx", "_enriched.pptx")
    prs.save(output_path)
    print(f"Successfully appended slides (if necessary) and enriched slide 8!")
    print(f"Saved to: {output_path}")

if __name__ == '__main__':
    filepath = '/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14_lightning.pptx'
    if os.path.exists(filepath):
        enrich_slide_8(filepath)
    else:
        print(f"File not found: {filepath}")
