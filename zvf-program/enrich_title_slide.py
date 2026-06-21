from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import sys

def main():
    pptx_path = "/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14_lightning.pptx"
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Failed to load pptx: {e}")
        return

    slide = prs.slides[0]

    # Colors based on the dark theme
    BG = RGBColor(0x0E, 0x11, 0x16)
    PANEL = RGBColor(0x16, 0x1B, 0x22)
    TEAL = RGBColor(0x2D, 0xD4, 0xBF)
    RED = RGBColor(0xF8, 0x51, 0x49)
    INK = RGBColor(0xE6, 0xED, 0xF3)
    LINECLR = RGBColor(0x30, 0x36, 0x3D)
    FONT = "Arial"

    # Helper function to add a rectangle
    def rect(s, l, t, w, h, fill=PANEL, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
        sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line
            sp.line.width = Pt(line_w)
        sp.shadow.inherit = False
        return sp

    # 1. Add a vertical accent line to the left of the main title block
    # Main title block starts around left=0.9, top=2.3
    rect(slide, 0.6, 2.3, 0.05, 1.8, fill=TEAL)

    # 2. Add a CONFIDENTIAL pill at the top right
    pill_text = "CONFIDENTIAL · DRAFT"
    w = 0.16 + 0.082 * len(pill_text)
    pill_sp = rect(slide, 12.5 - w, 0.55, w, 0.32, fill=RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        pill_sp.adjustments[0] = 0.5
    except Exception:
        pass
    
    tf = pill_sp.text_frame
    tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = pill_text
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = BG

    # 3. Add a background panel for the bottom text area to frame it nicely
    # Bottom text is at left=0.9, top=5.5
    # Let's send the panel to the back? python-pptx appends to the end of shapes (so it will cover text unless reordered)
    # Reordering shapes in python-pptx is not natively supported without manipulating xml.
    # So instead of a solid panel over the text, let's add a decorative outline box or some top/bottom borders.
    rect(slide, 0.9, 5.15, 11.5, 0.02, fill=LINECLR)
    rect(slide, 0.9, 6.25, 11.5, 0.02, fill=LINECLR)

    # 4. Add a small logo-like circle in the top left or a dot matrix pattern
    # Let's add three small dots to the right of the top eyebrow text
    for i in range(3):
        rect(slide, 4.0 + i*0.15, 0.72, 0.08, 0.08, fill=PANEL if i < 2 else TEAL, shape=MSO_SHAPE.OVAL)

    # Save the modified presentation
    out_path = pptx_path.replace(".pptx", "_enriched.pptx")
    prs.save(out_path)
    print(f"Slide 1 enriched and saved to: {out_path}")

if __name__ == '__main__':
    main()
