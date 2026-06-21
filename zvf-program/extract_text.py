import sys
from pptx import Presentation

prs = Presentation('/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14_lightning_enriched.pptx')
for i, slide in enumerate(prs.slides):
    print(f"\n--- SLIDE {i+1} ---")
    for shape in slide.shapes:
        if shape.has_text_frame:
            print(shape.text)
