from pptx import Presentation
prs = Presentation('/Users/arvind/Developer/tinker-rl-lab/zvf-program/ZVF_Program_Progress_2026-06-14_lightning.pptx')
slide = prs.slides[1]
print("--- SLIDE 2 CONTENT ---")
for shape in slide.shapes:
    if shape.has_text_frame:
        print(shape.text)
