from pptx import Presentation

# Open the presentation
prs = Presentation('ZVF_Program_Progress_2026-06-14.pptx')

print(f"Number of slides in full deck: {len(prs.slides)}")

if len(prs.slides) >= 13:
    slide = prs.slides[12]
    print("\n--- SLIDE 13 ---")
    for j, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            print(f"Text {j}: ", shape.text.replace('\n', ' '))
